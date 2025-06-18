import os
import shutil
import tempfile
import hashlib
import mimetypes
import logging
import time
from datetime import datetime
from pathlib import Path
import urllib.parse
import requests
import asyncio
import stat
import platform
import zipfile
import tarfile
import re
import chardet
import traceback
from typing import Dict, Any, List, Optional, Union
import pdfplumber
import pypandoc
from pptx import Presentation
import html2text
from openpyxl import load_workbook
import subprocess

from agents.tool.tool_base import ToolBase
from agents.utils.logger import logger

class FileParserError(Exception):
    """æ–‡ä»¶è§£æå¼‚å¸¸"""
    pass

class FileValidator:
    """æ–‡ä»¶éªŒè¯å™¨"""
    
    # æ”¯æŒçš„æ–‡ä»¶ç±»å‹å’Œå¯¹åº”çš„MIMEç±»å‹
    SUPPORTED_FORMATS = {
        '.pdf': 'application/pdf',
        '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        '.doc': 'application/msword',
        '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        '.ppt': 'application/vnd.ms-powerpoint',
        '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        '.xls': 'application/vnd.ms-excel',
        '.txt': 'text/plain',
        '.csv': 'text/csv',
        '.json': 'application/json',
        '.xml': 'application/xml',
        '.html': 'text/html',
        '.htm': 'text/html',
        '.md': 'text/markdown',
        '.markdown': 'text/markdown',
        '.rtf': 'application/rtf',
        '.odt': 'application/vnd.oasis.opendocument.text',
        '.epub': 'application/epub+zip',
        '.latex': 'application/x-latex',
        '.tex': 'application/x-tex'
    }
    
    # æ–‡ä»¶å¤§å°é™åˆ¶ (MB)
    MAX_FILE_SIZE = {
        '.pdf': 50,
        '.docx': 25,
        '.doc': 25,
        '.pptx': 100,
        '.ppt': 100,
        '.xlsx': 25,
        '.xls': 25,
        '.txt': 10,
        '.csv': 50,
        '.json': 10,
        '.xml': 10,
        '.html': 5,
        '.htm': 5,
        '.md': 5,
        '.markdown': 5,
        '.rtf': 10,
        '.odt': 25,
        '.epub': 50,
        '.latex': 10,
        '.tex': 10
    }
    
    @staticmethod
    def validate_file(file_path: str) -> Dict[str, Any]:
        """éªŒè¯æ–‡ä»¶çš„æœ‰æ•ˆæ€§"""
        try:
            path = Path(file_path)
            
            # æ£€æŸ¥æ–‡ä»¶æ˜¯å¦å­˜åœ¨
            if not path.exists():
                return {"valid": False, "error": f"æ–‡ä»¶ä¸å­˜åœ¨: {file_path}"}
            
            # æ£€æŸ¥æ˜¯å¦ä¸ºæ–‡ä»¶ï¼ˆéç›®å½•ï¼‰
            if not path.is_file():
                return {"valid": False, "error": f"è·¯å¾„ä¸æ˜¯æœ‰æ•ˆæ–‡ä»¶: {file_path}"}
            
            # è·å–æ–‡ä»¶æ‰©å±•å
            file_extension = path.suffix.lower()
            
            # æ£€æŸ¥æ–‡ä»¶æ ¼å¼æ˜¯å¦æ”¯æŒ
            if file_extension not in FileValidator.SUPPORTED_FORMATS:
                return {"valid": False, "error": f"ä¸æ”¯æŒçš„æ–‡ä»¶æ ¼å¼: {file_extension}"}
            
            # æ£€æŸ¥æ–‡ä»¶å¤§å°
            file_size_mb = path.stat().st_size / (1024 * 1024)
            max_size = FileValidator.MAX_FILE_SIZE.get(file_extension, 10)
            
            if file_size_mb > max_size:
                return {"valid": False, "error": f"æ–‡ä»¶è¿‡å¤§: {file_size_mb:.1f}MB > {max_size}MB"}
            
            # æ£€æŸ¥æ–‡ä»¶æ˜¯å¦å¯è¯»
            try:
                with open(file_path, 'rb') as f:
                    f.read(1024)  # å°è¯•è¯»å–å‰1KB
            except PermissionError:
                return {"valid": False, "error": "æ–‡ä»¶æ— è¯»å–æƒé™"}
            except Exception as e:
                return {"valid": False, "error": f"æ–‡ä»¶è¯»å–å¤±è´¥: {str(e)}"}
            
            return {
                "valid": True,
                "file_size_mb": file_size_mb,
                "file_extension": file_extension,
                "mime_type": FileValidator.SUPPORTED_FORMATS[file_extension]
            }
            
        except Exception as e:
            return {"valid": False, "error": f"æ–‡ä»¶éªŒè¯å¤±è´¥: {str(e)}"}

class TextProcessor:
    """æ–‡æœ¬å¤„ç†å™¨"""
    
    @staticmethod
    def detect_encoding(file_path: str) -> str:
        """æ£€æµ‹æ–‡ä»¶ç¼–ç """
        try:
            with open(file_path, 'rb') as f:
                raw_data = f.read(10000)  # è¯»å–å‰10KBè¿›è¡Œæ£€æµ‹
                result = chardet.detect(raw_data)
                return result.get('encoding', 'utf-8')
        except Exception:
            return 'utf-8'
    
    @staticmethod
    def clean_text(text: str) -> str:
        """æ¸…ç†æ–‡æœ¬å†…å®¹"""
        if not text:
            return ""
        
        # ç§»é™¤å¤šä½™çš„ç©ºç™½å­—ç¬¦
        text = re.sub(r'\n\s*\n', '\n\n', text)  # å¤šä¸ªè¿ç»­æ¢è¡Œç¬¦åˆå¹¶ä¸ºåŒæ¢è¡Œ
        text = re.sub(r'[ \t]+', ' ', text)      # å¤šä¸ªè¿ç»­ç©ºæ ¼åˆå¹¶ä¸ºå•ä¸ªç©ºæ ¼
        text = text.strip()
        
        return text
    
    @staticmethod
    def truncate_text(text: str, start_index: int = 0, max_length: int = 5000) -> str:
        """å®‰å…¨åœ°æˆªå–æ–‡æœ¬"""
        if not text:
            return ""
        
        start_index = max(0, start_index)
        end_index = min(len(text), start_index + max_length)
        
        return text[start_index:end_index]
    
    @staticmethod
    def get_text_stats(text: str) -> Dict[str, int]:
        """è·å–æ–‡æœ¬ç»Ÿè®¡ä¿¡æ¯"""
        if not text:
            return {"characters": 0, "words": 0, "lines": 0, "paragraphs": 0}
        
        return {
            "characters": len(text),
            "words": len(text.split()),
            "lines": text.count('\n') + 1,
            "paragraphs": len([p for p in text.split('\n\n') if p.strip()])
        }

class PDFParser:
    """PDFè§£æå™¨"""
    
    @staticmethod
    def extract_text(pdf_path: str) -> str:
        """ä»PDFæå–æ–‡æœ¬"""
        try:
            with pdfplumber.open(pdf_path) as pdf:
                text_parts = []
                for page_num, page in enumerate(pdf.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text:
                            text_parts.append(f"=== ç¬¬ {page_num + 1} é¡µ ===\n{page_text}")
                    except Exception as e:
                        logger.warning(f"PDFç¬¬{page_num + 1}é¡µè§£æå¤±è´¥: {e}")
                        text_parts.append(f"=== ç¬¬ {page_num + 1} é¡µ ===\n[é¡µé¢è§£æå¤±è´¥: {str(e)}]")
                
                return "\n\n".join(text_parts)
                
        except Exception as e:
            raise FileParserError(f"PDFè§£æå¤±è´¥: {str(e)}")
    
    @staticmethod
    def get_pdf_info(pdf_path: str) -> Dict[str, Any]:
        """è·å–PDFä¿¡æ¯"""
        try:
            with pdfplumber.open(pdf_path) as pdf:
                return {
                    "pages": len(pdf.pages),
                    "metadata": pdf.metadata or {}
                }
        except Exception as e:
            return {"pages": 0, "metadata": {}}

class OfficeParser:
    """Officeæ–‡æ¡£è§£æå™¨"""
    
    @staticmethod
    def extract_text_from_docx(file_path: str) -> str:
        """ä»DOCXæå–æ–‡æœ¬"""
        try:
            return pypandoc.convert_file(file_path, 'markdown', extra_args=['--extract-media=.'])
        except Exception as e:
            raise FileParserError(f"DOCXè§£æå¤±è´¥: {str(e)}")
    
    @staticmethod
    def extract_text_from_doc(file_path: str) -> str:
        """ä»DOCæå–æ–‡æœ¬ï¼ˆéœ€è¦antiwordï¼‰"""
        try:
            result = subprocess.run(
                ['antiword', file_path], 
                capture_output=True, 
                text=True, 
                timeout=30
            )
            if result.returncode == 0:
                return result.stdout
            else:
                raise FileParserError(f"DOCè½¬æ¢å¤±è´¥: {result.stderr}")
        except subprocess.TimeoutExpired:
            raise FileParserError("DOCè§£æè¶…æ—¶")
        except FileNotFoundError:
            raise FileParserError("æœªæ‰¾åˆ°antiwordå·¥å…·ï¼Œè¯·å®‰è£…: sudo apt-get install antiword")
        except Exception as e:
            raise FileParserError(f"DOCè§£æå¤±è´¥: {str(e)}")
        
    @staticmethod
    def extract_text_from_pptx(file_path: str) -> str:
        """ä»PPTXæå–æ–‡æœ¬"""
        try:
            prs = Presentation(file_path)
            slides_text = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_content = [f"=== å¹»ç¯ç‰‡ {slide_num + 1} ==="]
                
                # æŒ‰ä½ç½®æ’åºå½¢çŠ¶
                shapes = sorted(slide.shapes, key=lambda x: (x.top, x.left))
                
                for shape in shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())
                
                slides_text.append('\n'.join(slide_content))
            
            return '\n\n'.join(slides_text)
            
        except Exception as e:
            raise FileParserError(f"PPTXè§£æå¤±è´¥: {str(e)}")

class ExcelParser:
    """Excelè§£æå™¨"""
    
    @staticmethod
    def extract_text_from_xlsx(file_path: str) -> str:
        """ä»Excelæå–æ–‡æœ¬å¹¶è½¬æ¢ä¸ºMarkdown"""
        try:
            excel_data = ExcelParser._read_excel_to_dict(file_path)
            markdown_tables = []
            
            for sheet_name, sheet_data in excel_data.items():
                # é™åˆ¶è¡Œæ•°
                if len(sheet_data) > 100:
                    sheet_data = sheet_data[:100]
                    
                sheet_md = ExcelParser._sheet_data_to_markdown(sheet_data, sheet_name)
                markdown_tables.append(sheet_md)
            
            return '\n\n'.join(markdown_tables)
            
        except Exception as e:
            raise FileParserError(f"Excelè§£æå¤±è´¥: {str(e)}")
        
    @staticmethod
    def _read_excel_to_dict(file_path: str) -> Dict[str, List[List[str]]]:
        """è¯»å–Excelæ–‡ä»¶åˆ°å­—å…¸"""
        workbook = load_workbook(file_path, data_only=True, read_only=True)
        excel_data = {}

        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_data = []

            # è¯»å–æ•°æ®
            for row in sheet.iter_rows(values_only=True):
                if row is None:
                    continue
                row_data = [
                    str(cell).replace('\n', '\\n') if cell is not None else ''
                    for cell in row
                ]
                sheet_data.append(row_data)
                
            if not sheet_data:
                continue
                
            # ç»Ÿä¸€è¡Œé•¿åº¦
            max_length = max(len(row) for row in sheet_data) if sheet_data else 0
            for i, row in enumerate(sheet_data):
                if len(row) < max_length:
                    sheet_data[i] = row + [''] * (max_length - len(row))
            
            # æ¸…ç†ç©ºè¡Œå’Œç©ºåˆ—
            sheet_data = ExcelParser._clean_empty_rows_cols(sheet_data)

            if sheet_data:
                excel_data[sheet_name] = sheet_data
        
        workbook.close()
        return excel_data
    
    @staticmethod
    def _clean_empty_rows_cols(data: List[List[str]]) -> List[List[str]]:
        """æ¸…ç†ç©ºè¡Œå’Œç©ºåˆ—"""
        if not data:
            return data
        
        # ç§»é™¤ç©ºè¡Œ
        data = [row for row in data if any(cell.strip() for cell in row)]
        
        if not data:
            return data
        
        # ç§»é™¤ç©ºåˆ—
        cols_to_keep = []
        for col_idx in range(len(data[0])):
            if any(row[col_idx].strip() for row in data):
                cols_to_keep.append(col_idx)
        
        if cols_to_keep:
            data = [[row[i] for i in cols_to_keep] for row in data]
        
        return data
    
    @staticmethod
    def _sheet_data_to_markdown(sheet_data: List[List[str]], sheet_name: str) -> str:
        """è½¬æ¢å·¥ä½œè¡¨æ•°æ®ä¸ºMarkdown"""
        if not sheet_data:
            return f'## {sheet_name}\n\n(ç©ºå·¥ä½œè¡¨)'
        
        markdown_lines = [f'## {sheet_name}', '']
        
        # å¦‚æœæœ‰æ•°æ®ï¼Œç¬¬ä¸€è¡Œä½œä¸ºè¡¨å¤´
        if sheet_data:
            # è¡¨å¤´
            header = '| ' + ' | '.join(cell if cell else ' ' for cell in sheet_data[0]) + ' |'
            markdown_lines.append(header)
            
            # åˆ†éš”çº¿
            separator = '| ' + ' | '.join('---' for _ in sheet_data[0]) + ' |'
            markdown_lines.append(separator)
            
            # æ•°æ®è¡Œ
            for row in sheet_data[1:]:
                row_md = '| ' + ' | '.join(cell if cell else ' ' for cell in row) + ' |'
                markdown_lines.append(row_md)
        
        return '\n'.join(markdown_lines)

class WebParser:
    """ç½‘é¡µè§£æå™¨"""
    
    @staticmethod
    def extract_text_from_html(file_path: str) -> str:
        """ä»HTMLæ–‡ä»¶æå–æ–‡æœ¬"""
        try:
            encoding = TextProcessor.detect_encoding(file_path)
            with open(file_path, 'r', encoding=encoding) as file:
                html_content = file.read()
            
            return WebParser._html_to_text(html_content)
            
        except Exception as e:
            raise FileParserError(f"HTMLè§£æå¤±è´¥: {str(e)}")
    
    @staticmethod
    def extract_text_from_url(url: str, timeout: int = 30) -> str:
        """ä»URLæå–æ–‡æœ¬"""
        max_retries = 3
        initial_retry_delay = 2  # åˆå§‹é‡è¯•å»¶è¿Ÿå¢åŠ åˆ°2ç§’
        
        # æ ¹æ®URLç‰¹å¾è°ƒæ•´è¶…æ—¶æ—¶é—´å’Œç­–ç•¥
        adjusted_timeout = timeout
        is_github = 'github.com' in url or 'githubusercontent.com' in url
        
        if is_github:
            adjusted_timeout = max(timeout, 45)  # GitHubè‡³å°‘45ç§’
        elif any(domain in url for domain in ['google.com', 'microsoft.com', 'amazon.com']):
            adjusted_timeout = max(timeout, 40)  # å¤§å‹ç½‘ç«™è‡³å°‘40ç§’
        
        # åˆ›å»ºä¼šè¯å¯¹è±¡å¤ç”¨è¿æ¥
        session = requests.Session()
        
        # é’ˆå¯¹GitHubä½¿ç”¨æ›´å¼ºçš„åæ£€æµ‹ç­–ç•¥
        if is_github:
            session.headers.update({
                'User-Agent': 'Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8',
                'Accept-Language': 'en-US,en;q=0.9',
                'Accept-Encoding': 'gzip, deflate, br',
                'Cache-Control': 'no-cache',
                'Pragma': 'no-cache',
                'Sec-Ch-Ua': '"Not_A Brand";v="8", "Chromium";v="120", "Google Chrome";v="120"',
                'Sec-Ch-Ua-Mobile': '?0',
                'Sec-Ch-Ua-Platform': '"macOS"',
                'Sec-Fetch-Dest': 'document',
                'Sec-Fetch-Mode': 'navigate',
                'Sec-Fetch-Site': 'none',
                'Sec-Fetch-User': '?1',
                'Upgrade-Insecure-Requests': '1'
            })
        else:
            # æ™®é€šç½‘ç«™ä½¿ç”¨æ ‡å‡†è¯·æ±‚å¤´
            session.headers.update({
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8,application/signed-exchange;v=b3;q=0.7',
                'Accept-Language': 'zh-CN,zh;q=0.9,en;q=0.8',
                'Accept-Encoding': 'gzip, deflate, br',
                'DNT': '1',
                'Connection': 'keep-alive',
                'Upgrade-Insecure-Requests': '1',
                'Sec-Fetch-Dest': 'document',
                'Sec-Fetch-Mode': 'navigate',
                'Sec-Fetch-Site': 'none',
                'Sec-Fetch-User': '?1'
            })
        
        try:
            for attempt in range(max_retries):
                retry_delay = initial_retry_delay * (2 ** attempt)  # æŒ‡æ•°é€€é¿ï¼š2, 4, 8ç§’
                
                # å¯¹äºGitHubï¼Œåœ¨é¦–æ¬¡è¯·æ±‚å‰æ·»åŠ éšæœºå»¶è¿Ÿï¼Œæ¨¡æ‹Ÿäººç±»è¡Œä¸º
                if is_github and attempt == 0:
                    import random
                    human_delay = random.uniform(1, 3)  # 1-3ç§’éšæœºå»¶è¿Ÿ
                    logger.debug(f"ğŸ¤– GitHubåæ£€æµ‹ï¼šç­‰å¾…{human_delay:.1f}ç§’æ¨¡æ‹Ÿäººç±»è¡Œä¸º")
                    time.sleep(human_delay)
                
                try:
                    logger.debug(f"ğŸŒ å°è¯•ç¬¬{attempt + 1}æ¬¡è®¿é—®URL: {url} (è¶…æ—¶: {adjusted_timeout}ç§’)")
                    
                    # ä½¿ç”¨ä¼šè¯å‘èµ·è¯·æ±‚ï¼Œè®¾ç½®è¿æ¥å’Œè¯»å–è¶…æ—¶
                    response = session.get(
                        url, 
                        timeout=(10, adjusted_timeout),  # (è¿æ¥è¶…æ—¶, è¯»å–è¶…æ—¶)
                        allow_redirects=True,
                        stream=False  # ä¸ä½¿ç”¨æµå¼ä¸‹è½½ï¼Œç¡®ä¿å®Œæ•´è·å–å†…å®¹
                    )
                    
                    # æ£€æŸ¥å“åº”çŠ¶æ€
                    if response.status_code == 404:
                        raise FileParserError(f"URLä¸å­˜åœ¨ (404): {url}")
                    elif response.status_code == 403:
                        raise FileParserError(f"è®¿é—®è¢«ç¦æ­¢ (403): {url}")
                    elif response.status_code == 500:
                        raise FileParserError(f"æœåŠ¡å™¨å†…éƒ¨é”™è¯¯ (500): {url}")
                    elif response.status_code == 502:
                        # 502é”™è¯¯é€šå¸¸æ˜¯æš‚æ—¶çš„ï¼Œå€¼å¾—é‡è¯•
                        if attempt < max_retries - 1:
                            logger.warning(f"âš ï¸ æœåŠ¡å™¨ç½‘å…³é”™è¯¯ (502)ï¼Œç¬¬{attempt + 1}æ¬¡é‡è¯•ï¼Œ{retry_delay}ç§’åé‡è¯•")
                            time.sleep(retry_delay)
                            continue
                        else:
                            raise FileParserError(f"æœåŠ¡å™¨ç½‘å…³é”™è¯¯ (502): {url}")
                    elif response.status_code == 503:
                        # 503æœåŠ¡ä¸å¯ç”¨ï¼Œå€¼å¾—é‡è¯•
                        if attempt < max_retries - 1:
                            logger.warning(f"âš ï¸ æœåŠ¡æš‚æ—¶ä¸å¯ç”¨ (503)ï¼Œç¬¬{attempt + 1}æ¬¡é‡è¯•ï¼Œ{retry_delay}ç§’åé‡è¯•")
                            time.sleep(retry_delay)
                            continue
                        else:
                            raise FileParserError(f"æœåŠ¡æš‚æ—¶ä¸å¯ç”¨ (503): {url}")
                    elif response.status_code >= 400:
                        # å¯¹äºå…¶ä»–4xxå’Œ5xxé”™è¯¯ï¼Œå¦‚æœä¸æ˜¯æœ€åä¸€æ¬¡å°è¯•ï¼Œåˆ™é‡è¯•
                        if attempt < max_retries - 1:
                            logger.warning(f"âš ï¸ HTTPé”™è¯¯ {response.status_code}ï¼Œç¬¬{attempt + 1}æ¬¡é‡è¯•ï¼Œ{retry_delay}ç§’åé‡è¯•")
                            time.sleep(retry_delay)
                            continue
                        else:
                            raise FileParserError(f"URLè®¿é—®å¤±è´¥ (HTTP {response.status_code}): {url}")
                    
                    response.raise_for_status()
                    
                    # æ£€æŸ¥å†…å®¹é•¿åº¦
                    if len(response.content) == 0:
                        if attempt < max_retries - 1:
                            logger.warning(f"âš ï¸ å“åº”å†…å®¹ä¸ºç©ºï¼Œç¬¬{attempt + 1}æ¬¡é‡è¯•ï¼Œ{retry_delay}ç§’åé‡è¯•")
                            time.sleep(retry_delay)
                            continue
                        else:
                            raise FileParserError(f"å“åº”å†…å®¹ä¸ºç©º: {url}")
                    
                    # æ£€æŸ¥å†…å®¹ç±»å‹
                    content_type = response.headers.get('Content-Type', '').lower()
                    if 'text/html' not in content_type and 'text/plain' not in content_type and 'application/xhtml' not in content_type:
                        logger.warning(f"âš ï¸ æ£€æµ‹åˆ°éæ–‡æœ¬å†…å®¹ç±»å‹: {content_type}ï¼Œä»å°è¯•è§£æ")
                    
                    # æ™ºèƒ½ç¼–ç æ£€æµ‹å’Œå¤„ç†
                    final_text = WebParser._smart_decode_response(response, url)
                    
                    logger.debug(f"âœ… æˆåŠŸè·å–å†…å®¹ï¼Œé•¿åº¦: {len(final_text)} å­—ç¬¦")
                    return WebParser._html_to_text(final_text)
                    
                except requests.exceptions.ConnectTimeout as e:
                    if attempt < max_retries - 1:
                        logger.warning(f"âš ï¸ è¿æ¥è¶…æ—¶ï¼Œç¬¬{attempt + 1}æ¬¡é‡è¯•ï¼Œ{retry_delay}ç§’åé‡è¯•: {str(e)}")
                        time.sleep(retry_delay)
                        continue
                    else:
                        raise FileParserError(f"è¿æ¥è¶…æ—¶: {url} (è¿æ¥è¶…æ—¶: 10ç§’)")
                        
                except requests.exceptions.ReadTimeout as e:
                    if attempt < max_retries - 1:
                        # å¯¹äºè¯»å–è¶…æ—¶ï¼Œå¢åŠ ä¸‹æ¬¡å°è¯•çš„è¶…æ—¶æ—¶é—´
                        adjusted_timeout = min(adjusted_timeout * 1.5, 90)  # æœ€å¤§90ç§’
                        logger.warning(f"âš ï¸ è¯»å–è¶…æ—¶ï¼Œç¬¬{attempt + 1}æ¬¡é‡è¯•ï¼Œ{retry_delay}ç§’åé‡è¯•ï¼Œè°ƒæ•´è¶…æ—¶è‡³{adjusted_timeout:.0f}ç§’: {str(e)}")
                        time.sleep(retry_delay)
                        continue
                    else:
                        raise FileParserError(f"è¯»å–è¶…æ—¶: {url} (è¯»å–è¶…æ—¶: {adjusted_timeout:.0f}ç§’)")
                        
                except requests.exceptions.Timeout as e:
                    if attempt < max_retries - 1:
                        logger.warning(f"âš ï¸ è¯·æ±‚è¶…æ—¶ï¼Œç¬¬{attempt + 1}æ¬¡é‡è¯•ï¼Œ{retry_delay}ç§’åé‡è¯•: {str(e)}")
                        time.sleep(retry_delay)
                        continue
                    else:
                        raise FileParserError(f"è¯·æ±‚è¶…æ—¶: {url} (è¶…æ—¶æ—¶é—´: {adjusted_timeout}ç§’)")
                        
                except requests.exceptions.ConnectionError as e:
                    if attempt < max_retries - 1:
                        logger.warning(f"âš ï¸ è¿æ¥é”™è¯¯ï¼Œç¬¬{attempt + 1}æ¬¡é‡è¯•ï¼Œ{retry_delay}ç§’åé‡è¯•: {str(e)}")
                        time.sleep(retry_delay)
                        continue
                    else:
                        raise FileParserError(f"æ— æ³•è¿æ¥åˆ°URL: {url} - {str(e)}")
                        
                except requests.exceptions.RequestException as e:
                    # å¯¹äºå…¶ä»–è¯·æ±‚å¼‚å¸¸ï¼Œæ ¹æ®é”™è¯¯ç±»å‹å†³å®šæ˜¯å¦é‡è¯•
                    error_str = str(e).lower()
                    if any(keyword in error_str for keyword in ['timeout', 'connection', 'network', 'dns']):
                        if attempt < max_retries - 1:
                            logger.warning(f"âš ï¸ ç½‘ç»œç›¸å…³é”™è¯¯ï¼Œç¬¬{attempt + 1}æ¬¡é‡è¯•ï¼Œ{retry_delay}ç§’åé‡è¯•: {str(e)}")
                            time.sleep(retry_delay)
                            continue
                    raise FileParserError(f"URLè®¿é—®å¤±è´¥: {url} - {str(e)}")
                    
                except Exception as e:
                    # å¯¹äºå…¶ä»–å¼‚å¸¸ï¼Œè®°å½•è¯¦ç»†ä¿¡æ¯
                    logger.error(f"ğŸ’¥ æ„å¤–å¼‚å¸¸: {type(e).__name__}: {str(e)}")
                    raise FileParserError(f"URLè§£æå¤±è´¥: {url} - {str(e)}")
            
            # å¦‚æœæ‰€æœ‰é‡è¯•éƒ½å¤±è´¥äº†ï¼ˆç†è®ºä¸Šä¸ä¼šåˆ°è¾¾è¿™é‡Œï¼‰
            raise FileParserError(f"URLè®¿é—®å¤±è´¥ï¼Œå·²é‡è¯•{max_retries}æ¬¡: {url}")
            
        finally:
            # ç¡®ä¿å…³é—­ä¼šè¯
            session.close()
    
    @staticmethod
    def _smart_decode_response(response, url: str) -> str:
        """æ™ºèƒ½è§£ç HTTPå“åº”å†…å®¹"""
        import chardet
        import re
        
        # 1. é¦–å…ˆå°è¯•ä»Content-Type headerè·å–ç¼–ç 
        content_type = response.headers.get('Content-Type', '')
        charset_match = re.search(r'charset=([^;\s]+)', content_type, re.IGNORECASE)
        header_encoding = None
        if charset_match:
            header_encoding = charset_match.group(1).strip('"\'').lower()
            logger.debug(f"ğŸ” ä»HTTPå¤´è·å–ç¼–ç : {header_encoding}")
        
        # 2. å°è¯•ä»HTMLçš„metaæ ‡ç­¾è·å–ç¼–ç 
        html_encoding = None
        try:
            # ç”¨åŸå§‹å­—èŠ‚æ•°æ®æ£€æŸ¥HTML metaæ ‡ç­¾
            content_preview = response.content[:2048]  # åªæ£€æŸ¥å‰2KB
            content_str = content_preview.decode('ascii', errors='ignore')
            
            # æŸ¥æ‰¾meta charset
            meta_matches = [
                re.search(r'<meta[^>]+charset["\s]*=["\s]*([^">\s]+)', content_str, re.IGNORECASE),
                re.search(r'<meta[^>]+content[^>]*charset=([^">\s;]+)', content_str, re.IGNORECASE)
            ]
            
            for match in meta_matches:
                if match:
                    html_encoding = match.group(1).strip('"\'').lower()
                    logger.debug(f"ğŸ” ä»HTML metaæ ‡ç­¾è·å–ç¼–ç : {html_encoding}")
                    break
        except Exception as e:
            logger.debug(f"âš ï¸ HTMLç¼–ç æ£€æµ‹å¤±è´¥: {e}")
        
        # 3. ä½¿ç”¨chardetè¿›è¡Œè‡ªåŠ¨æ£€æµ‹
        detected_encoding = None
        try:
            detection_result = chardet.detect(response.content)
            if detection_result and detection_result.get('confidence', 0) > 0.7:
                detected_encoding = detection_result['encoding'].lower()
                confidence = detection_result.get('confidence', 0)
                logger.debug(f"ğŸ” chardetæ£€æµ‹ç¼–ç : {detected_encoding} (ç½®ä¿¡åº¦: {confidence:.2f})")
        except Exception as e:
            logger.debug(f"âš ï¸ chardetç¼–ç æ£€æµ‹å¤±è´¥: {e}")
        
        # 4. æ ¹æ®URLåˆ¤æ–­å¯èƒ½çš„ç¼–ç ï¼ˆé’ˆå¯¹ä¸­æ–‡ç½‘ç«™ï¼‰
        url_encoding = None
        if any(domain in url.lower() for domain in ['.cn', '.com.cn', 'baidu', 'sina', 'qq', '163', 'sohu', 'taobao', 'jd']):
            url_encoding = 'utf-8'  # å¤§å¤šæ•°ä¸­æ–‡ç½‘ç«™ä½¿ç”¨UTF-8
            logger.debug(f"ğŸ” æ ¹æ®URLæ¨æµ‹ä¸­æ–‡ç½‘ç«™ç¼–ç : {url_encoding}")
        
        # 5. ç¼–ç ä¼˜å…ˆçº§å’Œå°è¯•é¡ºåº
        encoding_candidates = []
        
        # ä¼˜å…ˆçº§æ’åº
        if header_encoding:
            encoding_candidates.append(header_encoding)
        if html_encoding and html_encoding != header_encoding:
            encoding_candidates.append(html_encoding)
        if detected_encoding and detected_encoding not in encoding_candidates:
            encoding_candidates.append(detected_encoding)
        if url_encoding and url_encoding not in encoding_candidates:
            encoding_candidates.append(url_encoding)
        
        # æ·»åŠ å¸¸è§ç¼–ç ä½œä¸ºåå¤‡
        fallback_encodings = ['utf-8', 'gbk', 'gb2312', 'big5', 'iso-8859-1', 'windows-1252']
        for enc in fallback_encodings:
            if enc not in encoding_candidates:
                encoding_candidates.append(enc)
        
        # 6. é€ä¸ªå°è¯•è§£ç 
        for encoding in encoding_candidates:
            try:
                if encoding:
                    # è§„èŒƒåŒ–ç¼–ç åç§°
                    encoding = encoding.replace('gb2312', 'gbk')  # gb2312æ˜¯gbkçš„å­é›†
                    encoding = encoding.replace('iso-8859-1', 'latin1')
                    
                    decoded_text = response.content.decode(encoding, errors='replace')
                    
                    # ç®€å•éªŒè¯ï¼šæ£€æŸ¥æ˜¯å¦åŒ…å«å¤§é‡æ›¿æ¢å­—ç¬¦
                    replacement_ratio = decoded_text.count('ï¿½') / max(len(decoded_text), 1)
                    if replacement_ratio < 0.1:  # æ›¿æ¢å­—ç¬¦å°‘äº10%
                        logger.debug(f"âœ… æˆåŠŸä½¿ç”¨ç¼–ç  {encoding} è§£ç å†…å®¹")
                        return decoded_text
                    else:
                        logger.debug(f"âš ï¸ ç¼–ç  {encoding} äº§ç”Ÿè¿‡å¤šæ›¿æ¢å­—ç¬¦ ({replacement_ratio:.2%})")
            except (UnicodeDecodeError, LookupError) as e:
                logger.debug(f"âš ï¸ ç¼–ç  {encoding} è§£ç å¤±è´¥: {e}")
                continue
        
        # 7. æœ€åçš„fallbackï¼šä½¿ç”¨UTF-8å¹¶å¿½ç•¥é”™è¯¯
        logger.warning(f"âš ï¸ æ‰€æœ‰ç¼–ç å°è¯•å¤±è´¥ï¼Œä½¿ç”¨UTF-8å¼ºåˆ¶è§£ç : {url}")
        return response.content.decode('utf-8', errors='replace')
    
    @staticmethod
    def _html_to_text(html_content: str) -> str:
        """HTMLè½¬æ–‡æœ¬"""
        h = html2text.HTML2Text()
        h.ignore_links = False
        h.ignore_images = False
        h.body_width = 0  # ä¸é™åˆ¶è¡Œå®½
        return h.handle(html_content)

class PlainTextParser:
    """çº¯æ–‡æœ¬è§£æå™¨"""
    
    @staticmethod
    def extract_text_from_plain_file(file_path: str) -> str:
        """ä»çº¯æ–‡æœ¬æ–‡ä»¶æå–å†…å®¹"""
        try:
            encoding = TextProcessor.detect_encoding(file_path)
            with open(file_path, 'r', encoding=encoding) as file:
                return file.read()
        except Exception as e:
            raise FileParserError(f"æ–‡æœ¬æ–‡ä»¶è§£æå¤±è´¥: {str(e)}")
    
    @staticmethod
    def extract_text_with_pandoc(file_path: str, input_format: str = None) -> str:
        """ä½¿ç”¨Pandocæå–æ–‡æœ¬"""
        try:
            if input_format:
                return pypandoc.convert_file(file_path, 'markdown', format=input_format)
            else:
                return pypandoc.convert_file(file_path, 'markdown')
        except Exception as e:
            raise FileParserError(f"Pandocè§£æå¤±è´¥: {str(e)}")

class FileParserTool(ToolBase):
    """æ–‡ä»¶è§£æå·¥å…·é›†"""
    
    def __init__(self):
        logger.debug("Initializing FileParserTool")
        super().__init__()

    @ToolBase.tool()
    def extract_text_from_file(
        self, 
        input_file_path: str, 
        start_index: int = 0, 
        max_length: int = 5000,
        include_metadata: bool = False
    ) -> Dict[str, Any]:
        """ä»æœ¬åœ°çš„å„ç§æ ¼å¼çš„æ–‡ä»¶ä¸­æå–æ–¹ä¾¿é˜…è¯»çš„markdownæ–‡æœ¬å†…å®¹

        Args:
            input_file_path (str): è¾“å…¥æ–‡ä»¶è·¯å¾„ï¼Œæœ¬åœ°çš„ç»å¯¹è·¯å¾„
            start_index (int): å¼€å§‹æå–çš„å­—ç¬¦ä½ç½®ï¼Œé»˜è®¤0
            max_length (int): æœ€å¤§æå–é•¿åº¦ï¼Œé»˜è®¤5000å­—ç¬¦
            include_metadata (bool): æ˜¯å¦åŒ…å«æ–‡ä»¶å…ƒæ•°æ®ï¼Œé»˜è®¤False

        Returns:
            Dict[str, Any]: åŒ…å«æå–æ–‡æœ¬å’Œç›¸å…³ä¿¡æ¯çš„å­—å…¸
        """
        start_time = time.time()
        operation_id = hashlib.md5(f"extract_{input_file_path}_{time.time()}".encode()).hexdigest()[:8]
        logger.info(f"ğŸ“„ extract_text_from_fileå¼€å§‹æ‰§è¡Œ [{operation_id}] - æ–‡ä»¶: {input_file_path}")
        logger.info(f"ğŸ”§ å‚æ•°: start_index={start_index}, max_length={max_length}, include_metadata={include_metadata}")
        
        try:
            # éªŒè¯æ–‡ä»¶
            logger.debug(f"ğŸ” å¼€å§‹æ–‡ä»¶éªŒè¯")
            validation_result = FileValidator.validate_file(input_file_path)
            
            if not validation_result["valid"]:
                error_time = time.time() - start_time
                logger.error(f"âŒ æ–‡ä»¶éªŒè¯å¤±è´¥ [{operation_id}] - é”™è¯¯: {validation_result['error']}, è€—æ—¶: {error_time:.2f}ç§’")
                return {
                    "success": False,
                    "error": validation_result["error"],
                    "file_path": input_file_path,
                    "execution_time": error_time,
                    "operation_id": operation_id
                }
            
            file_extension = validation_result["file_extension"]
            file_size_mb = validation_result["file_size_mb"]
            logger.info(f"âœ… æ–‡ä»¶éªŒè¯é€šè¿‡ [{operation_id}] - æ ¼å¼: {file_extension}, å¤§å°: {file_size_mb:.2f}MB")
            
            # æ ¹æ®æ–‡ä»¶ç±»å‹é€‰æ‹©è§£æå™¨
            parse_start_time = time.time()
            logger.info(f"ğŸ”§ å¼€å§‹æ–‡ä»¶è§£æ - æ ¼å¼: {file_extension}")
            
            extracted_text = ""
            metadata = {}
            
            try:
                if file_extension == '.pdf':
                    logger.debug(f"ğŸ“• ä½¿ç”¨PDFè§£æå™¨")
                    extracted_text = PDFParser.extract_text(input_file_path)
                    if include_metadata:
                        metadata = PDFParser.get_pdf_info(input_file_path)
                        
                elif file_extension in ['.docx', '.doc']:
                    logger.debug(f"ğŸ“ ä½¿ç”¨Wordè§£æå™¨")
                    if file_extension == '.docx':
                        extracted_text = OfficeParser.extract_text_from_docx(input_file_path)
                    else:
                        extracted_text = OfficeParser.extract_text_from_doc(input_file_path)
                        
                elif file_extension in ['.pptx', '.ppt']:
                    logger.debug(f"ğŸ“Š ä½¿ç”¨PowerPointè§£æå™¨")
                    if file_extension == '.pptx':
                        extracted_text = OfficeParser.extract_text_from_pptx(input_file_path)
                    else:
                        # PPTéœ€è¦é¢å¤–çš„å¤„ç†ï¼Œè¿™é‡Œç®€åŒ–å¤„ç†
                        extracted_text = "PPTæ–‡ä»¶æš‚ä¸æ”¯æŒï¼Œè¯·è½¬æ¢ä¸ºPPTXæ ¼å¼"
                        
                elif file_extension in ['.xlsx', '.xls']:
                    logger.debug(f"ğŸ“ˆ ä½¿ç”¨Excelè§£æå™¨")
                    extracted_text = ExcelParser.extract_text_from_xlsx(input_file_path)
                    
                elif file_extension in ['.html', '.htm']:
                    logger.debug(f"ğŸŒ ä½¿ç”¨HTMLè§£æå™¨")
                    extracted_text = WebParser.extract_text_from_html(input_file_path)
                    
                elif file_extension in ['.txt', '.csv', '.json', '.xml', '.md', '.markdown']:
                    logger.debug(f"ğŸ“„ ä½¿ç”¨çº¯æ–‡æœ¬è§£æå™¨")
                    extracted_text = PlainTextParser.extract_text_from_plain_file(input_file_path)
                    
                else:
                    # å°è¯•ä½¿ç”¨Pandocè§£æ
                    logger.debug(f"ğŸ”§ å°è¯•ä½¿ç”¨Pandocè§£æå™¨")
                    extracted_text = PlainTextParser.extract_text_with_pandoc(input_file_path)
                
                parse_time = time.time() - parse_start_time
                logger.info(f"âœ… æ–‡ä»¶è§£ææˆåŠŸ [{operation_id}] - åŸå§‹æ–‡æœ¬é•¿åº¦: {len(extracted_text)}, è§£æè€—æ—¶: {parse_time:.2f}ç§’")
                
            except Exception as parse_error:
                parse_time = time.time() - parse_start_time
                logger.warning(f"âš ï¸ ä¸»è§£æå™¨å¤±è´¥ [{operation_id}] - é”™è¯¯: {str(parse_error)}, è€—æ—¶: {parse_time:.2f}ç§’")
                logger.debug(f"ğŸ”§ å°è¯•Pandocå¤‡ç”¨è§£æå™¨")
                
                try:
                    extracted_text = PlainTextParser.extract_text_with_pandoc(input_file_path)
                    logger.info(f"âœ… Pandocå¤‡ç”¨è§£ææˆåŠŸ [{operation_id}] - æ–‡æœ¬é•¿åº¦: {len(extracted_text)}")
                except Exception as fallback_error:
                    error_time = time.time() - start_time
                    logger.error(f"âŒ æ‰€æœ‰è§£æå™¨éƒ½å¤±è´¥ [{operation_id}] - ä¸»é”™è¯¯: {str(parse_error)}, å¤‡ç”¨é”™è¯¯: {str(fallback_error)}, è€—æ—¶: {error_time:.2f}ç§’")
                    return {
                        "success": False,
                        "error": f"æ–‡ä»¶è§£æå¤±è´¥ - ä¸»é”™è¯¯: {str(parse_error)}, å¤‡ç”¨é”™è¯¯: {str(fallback_error)}",
                        "file_path": input_file_path,
                        "execution_time": error_time,
                        "operation_id": operation_id
                    }
            
            # æ¸…ç†å’Œå¤„ç†æ–‡æœ¬
            logger.debug(f"ğŸ§¹ å¼€å§‹æ–‡æœ¬æ¸…ç†å’Œå¤„ç†")
            cleaned_text = TextProcessor.clean_text(extracted_text)
            truncated_text = TextProcessor.truncate_text(cleaned_text, start_index, max_length)
            text_stats = TextProcessor.get_text_stats(cleaned_text)
            
            total_time = time.time() - start_time
            
            logger.info(f"âœ… æ–‡æœ¬æå–å®Œæˆ [{operation_id}] - æ¸…ç†åé•¿åº¦: {len(cleaned_text)}, æˆªå–é•¿åº¦: {len(truncated_text)}, æ€»è€—æ—¶: {total_time:.2f}ç§’")
            logger.debug(f"ğŸ“Š æ–‡æœ¬ç»Ÿè®¡: {text_stats}")
            
            # æ„å»ºç»“æœ
            result = {
                "success": True,
                "text": truncated_text,
                "file_info": {
                    "file_path": input_file_path,
                    "file_extension": file_extension,
                    "file_size_mb": round(file_size_mb, 2),
                    "mime_type": validation_result["mime_type"]
                },
                "text_info": {
                    "original_length": len(extracted_text),
                    "cleaned_length": len(cleaned_text),
                    "extracted_length": len(truncated_text),
                    "start_index": start_index,
                    "max_length": max_length,
                    **text_stats
                },
                "execution_time": total_time,
                "operation_id": operation_id
            }
            
            if include_metadata and metadata:
                result["metadata"] = metadata
                logger.debug(f"ğŸ“‹ åŒ…å«å…ƒæ•°æ®: {len(metadata)} é¡¹")
            
            return result
            
        except Exception as e:
            error_time = time.time() - start_time
            logger.error(f"ğŸ’¥ æ–‡æœ¬æå–å¼‚å¸¸ [{operation_id}] - é”™è¯¯: {str(e)}, è€—æ—¶: {error_time:.2f}ç§’")
            logger.error(f"ğŸ” å¼‚å¸¸è¯¦æƒ…: {traceback.format_exc()}")
            
            return {
                "success": False,
                "error": str(e),
                "file_path": input_file_path,
                "execution_time": error_time,
                "operation_id": operation_id
            }

    @ToolBase.tool()
    def extract_text_from_url(
        self,
        url: str,
        start_index: int = 0,
        max_length: int = 5000,
        timeout: int = 30
    ) -> Dict[str, Any]:
        """ä»URLæå–ç½‘é¡µhtmlæ–‡æœ¬å†…å®¹

        Args:
            url (str): ç›®æ ‡URLåœ°å€
            start_index (int): å¼€å§‹æå–çš„å­—ç¬¦ä½ç½®ï¼Œé»˜è®¤0
            max_length (int): æœ€å¤§æå–é•¿åº¦ï¼Œé»˜è®¤5000å­—ç¬¦
            timeout (int): è¯·æ±‚è¶…æ—¶æ—¶é—´ï¼Œé»˜è®¤30ç§’

        Returns:
            Dict[str, Any]: åŒ…å«æå–æ–‡æœ¬å’Œç›¸å…³ä¿¡æ¯çš„å­—å…¸
        """
        start_time = time.time()
        operation_id = hashlib.md5(f"url_{url}_{time.time()}".encode()).hexdigest()[:8]
        logger.info(f"ğŸŒ extract_text_from_urlå¼€å§‹æ‰§è¡Œ [{operation_id}] - URL: {url}")
        logger.info(f"ğŸ”§ å‚æ•°: start_index={start_index}, max_length={max_length}, timeout={timeout}ç§’")
        
        try:
            # éªŒè¯URLæ ¼å¼
            logger.debug(f"ğŸ” éªŒè¯URLæ ¼å¼")
            if not url.startswith(('http://', 'https://')):
                error_time = time.time() - start_time
                logger.error(f"âŒ URLæ ¼å¼æ— æ•ˆ [{operation_id}] - URL: {url}, è€—æ—¶: {error_time:.2f}ç§’")
                return {
                    "success": False,
                    "error": "URLå¿…é¡»ä»¥http://æˆ–https://å¼€å¤´",
                    "error_type": "invalid_url_format",
                    "url": url,
                    "execution_time": error_time,
                    "operation_id": operation_id
                }
            
            # æå–ç½‘é¡µå†…å®¹
            fetch_start_time = time.time()
            logger.info(f"ğŸŒ å¼€å§‹è·å–ç½‘é¡µå†…å®¹")
            
            try:
                extracted_text = WebParser.extract_text_from_url(url, timeout)
                fetch_time = time.time() - fetch_start_time
                logger.info(f"âœ… ç½‘é¡µå†…å®¹è·å–æˆåŠŸ [{operation_id}] - åŸå§‹æ–‡æœ¬é•¿åº¦: {len(extracted_text)}, è·å–è€—æ—¶: {fetch_time:.2f}ç§’")
                
            except FileParserError as e:
                error_msg = str(e)
                fetch_time = time.time() - fetch_start_time
                error_time = time.time() - start_time
                
                # æ ¹æ®é”™è¯¯ç±»å‹æä¾›æ›´å‹å¥½çš„é”™è¯¯ä¿¡æ¯
                error_type = "unknown_error"
                user_friendly_error = error_msg
                
                if "404" in error_msg or "ä¸å­˜åœ¨" in error_msg:
                    error_type = "url_not_found"
                    user_friendly_error = f"URLä¸å­˜åœ¨æˆ–å·²å¤±æ•ˆ: {url}"
                elif "403" in error_msg or "ç¦æ­¢" in error_msg:
                    error_type = "access_forbidden"
                    user_friendly_error = f"è®¿é—®è¢«æ‹’ç»ï¼Œå¯èƒ½éœ€è¦ç™»å½•æˆ–æƒé™: {url}"
                elif "500" in error_msg or "æœåŠ¡å™¨" in error_msg:
                    error_type = "server_error"
                    user_friendly_error = f"ç›®æ ‡æœåŠ¡å™¨å‡ºç°é—®é¢˜: {url}"
                elif "è¿æ¥" in error_msg or "Connection" in error_msg:
                    error_type = "connection_error"
                    user_friendly_error = f"æ— æ³•è¿æ¥åˆ°ç›®æ ‡ç½‘ç«™ï¼Œè¯·æ£€æŸ¥ç½‘ç»œè¿æ¥: {url}"
                elif "è¿æ¥è¶…æ—¶" in error_msg or "ConnectTimeout" in error_msg:
                    error_type = "connect_timeout_error"
                    user_friendly_error = f"è¿æ¥è¶…æ—¶ï¼Œç½‘ç«™å“åº”è¿‡æ…¢ï¼Œè¯·ç¨åé‡è¯•: {url}"
                elif "è¯»å–è¶…æ—¶" in error_msg or "ReadTimeout" in error_msg:
                    error_type = "read_timeout_error"
                    user_friendly_error = f"æ•°æ®è¯»å–è¶…æ—¶ï¼Œç½‘é¡µå†…å®¹è¾ƒå¤§æˆ–ç½‘é€Ÿè¾ƒæ…¢ï¼Œå»ºè®®å¢åŠ è¶…æ—¶æ—¶é—´: {url}"
                elif "è¶…æ—¶" in error_msg or "Timeout" in error_msg:
                    error_type = "timeout_error"
                    user_friendly_error = f"ç½‘é¡µåŠ è½½è¶…æ—¶ï¼Œè¯·ç¨åé‡è¯•æˆ–å¢åŠ è¶…æ—¶æ—¶é—´: {url}"
                
                logger.error(f"âŒ ç½‘é¡µå†…å®¹è·å–å¤±è´¥ [{operation_id}] - é”™è¯¯ç±»å‹: {error_type}, è€—æ—¶: {error_time:.2f}ç§’")
                
                return {
                    "success": False,
                    "error": user_friendly_error,
                    "error_type": error_type,
                    "error_details": error_msg,
                    "url": url,
                    "suggestions": self._get_error_suggestions(error_type, url),
                    "execution_time": error_time,
                    "operation_id": operation_id
                }
            
            # æ¸…ç†å’Œå¤„ç†æ–‡æœ¬
            logger.debug(f"ğŸ§¹ å¼€å§‹æ–‡æœ¬æ¸…ç†å’Œå¤„ç†")
            cleaned_text = TextProcessor.clean_text(extracted_text)
            truncated_text = TextProcessor.truncate_text(cleaned_text, start_index, max_length)
            text_stats = TextProcessor.get_text_stats(cleaned_text)
            
            total_time = time.time() - start_time
            
            logger.info(f"âœ… URLæ–‡æœ¬æå–å®Œæˆ [{operation_id}] - æ¸…ç†åé•¿åº¦: {len(cleaned_text)}, æˆªå–é•¿åº¦: {len(truncated_text)}, æ€»è€—æ—¶: {total_time:.2f}ç§’")
            logger.debug(f"ğŸ“Š æ–‡æœ¬ç»Ÿè®¡: {text_stats}")
            
            return {
                "success": True,
                "text": truncated_text,
                "url_info": {
                    "url": url,
                    "timeout": timeout,
                    "fetch_time": fetch_time
                },
                "text_info": {
                    "original_length": len(extracted_text),
                    "cleaned_length": len(cleaned_text),
                    "extracted_length": len(truncated_text),
                    "start_index": start_index,
                    "max_length": max_length,
                    **text_stats
                },
                "execution_time": total_time,
                "operation_id": operation_id
            }
            
        except Exception as e:
            error_time = time.time() - start_time
            logger.error(f"ğŸ’¥ URLæ–‡æœ¬æå–å¼‚å¸¸ [{operation_id}] - é”™è¯¯: {str(e)}, è€—æ—¶: {error_time:.2f}ç§’")
            logger.error(f"ğŸ” å¼‚å¸¸è¯¦æƒ…: {traceback.format_exc()}")
            
            return {
                "success": False,
                "error": f"å¤„ç†URLæ—¶å‘ç”Ÿæ„å¤–é”™è¯¯: {str(e)}",
                "error_type": "unexpected_error",
                "url": url,
                "execution_time": error_time,
                "operation_id": operation_id
            }
    
    def _get_error_suggestions(self, error_type: str, url: str) -> List[str]:
        """æ ¹æ®é”™è¯¯ç±»å‹æä¾›å»ºè®®"""
        suggestions = []
        
        if error_type == "url_not_found":
            suggestions = [
                "æ£€æŸ¥URLæ˜¯å¦æ­£ç¡®æ‹¼å†™",
                "ç¡®è®¤ç½‘é¡µæ˜¯å¦è¿˜å­˜åœ¨",
                "å°è¯•è®¿é—®ç½‘ç«™é¦–é¡µç¡®è®¤ç½‘ç«™æ˜¯å¦å¯ç”¨"
            ]
        elif error_type == "access_forbidden":
            suggestions = [
                "æ£€æŸ¥æ˜¯å¦éœ€è¦ç™»å½•è´¦æˆ·",
                "ç¡®è®¤æ˜¯å¦æœ‰è®¿é—®æƒé™",
                "å°è¯•åœ¨æµè§ˆå™¨ä¸­æ‰‹åŠ¨è®¿é—®è¯¥URL"
            ]
        elif error_type == "connection_error":
            suggestions = [
                "æ£€æŸ¥ç½‘ç»œè¿æ¥æ˜¯å¦æ­£å¸¸",
                "ç¡®è®¤é˜²ç«å¢™æ˜¯å¦é˜»æ­¢äº†è®¿é—®",
                "å°è¯•è®¿é—®å…¶ä»–ç½‘ç«™ç¡®è®¤ç½‘ç»œçŠ¶æ€"
            ]
        elif error_type == "connect_timeout_error":
            suggestions = [
                "æ£€æŸ¥ç½‘ç»œè¿æ¥æ˜¯å¦ç¨³å®š",
                "ç¨åé‡è¯•",
                "å°è¯•ä½¿ç”¨VPNæˆ–ä»£ç†æœåŠ¡"
            ]
        elif error_type == "read_timeout_error":
            suggestions = [
                "æ˜¾è‘—å¢åŠ è¶…æ—¶æ—¶é—´å‚æ•°ï¼ˆå¦‚è®¾ç½®ä¸º60-90ç§’ï¼‰",
                "æ£€æŸ¥ç½‘ç»œé€Ÿåº¦æ˜¯å¦æ­£å¸¸",
                "åˆ†æ®µè·å–å†…å®¹ï¼Œæˆ–ç¨åé‡è¯•"
            ]
        elif error_type == "timeout_error":
            suggestions = [
                "å¢åŠ è¶…æ—¶æ—¶é—´å‚æ•°",
                "ç¨åé‡è¯•",
                "æ£€æŸ¥ç½‘ç»œé€Ÿåº¦æ˜¯å¦æ­£å¸¸"
            ]
        elif error_type == "server_error":
            suggestions = [
                "ç¨åé‡è¯•",
                "è”ç³»ç½‘ç«™ç®¡ç†å‘˜",
                "å°è¯•è®¿é—®ç½‘ç«™çš„å…¶ä»–é¡µé¢"
            ]
        else:
            suggestions = [
                "æ£€æŸ¥URLæ ¼å¼æ˜¯å¦æ­£ç¡®",
                "ç¡®è®¤ç½‘ç»œè¿æ¥æ­£å¸¸",
                "ç¨åé‡è¯•"
            ]
        
        return suggestions

    @ToolBase.tool()
    def batch_extract_text(
        self,
        file_paths: List[str],
        max_length: int = 3000,
        include_metadata: bool = False
    ) -> Dict[str, Any]:
        """æ‰¹é‡æå–å¤šä¸ªæ–‡ä»¶çš„æ–‡æœ¬

        Args:
            file_paths (List[str]): æ–‡ä»¶è·¯å¾„åˆ—è¡¨
            max_length (int): æ¯ä¸ªæ–‡ä»¶çš„æœ€å¤§æå–é•¿åº¦ï¼Œé»˜è®¤3000
            include_metadata (bool): æ˜¯å¦åŒ…å«å…ƒæ•°æ®ï¼Œé»˜è®¤False

        Returns:
            Dict[str, Any]: æ‰¹é‡å¤„ç†ç»“æœ
        """
        start_time = time.time()
        results = []
        successful = 0
        failed = 0
        
        for file_path in file_paths:
            try:
                result = self.extract_text_from_file(
                    input_file_path=file_path,
                    start_index=0,
                    max_length=max_length,
                    include_metadata=include_metadata
                )
                
                if result["success"]:
                    successful += 1
                else:
                    failed += 1
                    
                results.append({
                    "file_path": file_path,
                    "result": result
                })
                
            except Exception as e:
                failed += 1
                results.append({
                    "file_path": file_path,
                    "result": {
                        "success": False,
                        "error": f"å¤„ç†å¤±è´¥: {str(e)}",
                        "text": "",
                        "execution_time": 0
                    }
                })
        
        processing_time = time.time() - start_time
        
        return {
            "success": True,
            "message": f"æ‰¹é‡å¤„ç†å®Œæˆï¼ŒæˆåŠŸ{successful}ä¸ªï¼Œå¤±è´¥{failed}ä¸ª",
            "summary": {
                "total_files": len(file_paths),
                "successful": successful,
                "failed": failed,
                "processing_time": round(processing_time, 2)
            },
            "results": results
        }

    @ToolBase.tool()
    def get_supported_formats(self) -> Dict[str, Any]:
        """è·å–æ”¯æŒçš„æ–‡ä»¶æ ¼å¼åˆ—è¡¨

        Returns:
            Dict[str, Any]: æ”¯æŒçš„æ–‡ä»¶æ ¼å¼å’Œç›¸å…³ä¿¡æ¯
        """
        return {
            "success": True,
            "supported_formats": FileValidator.SUPPORTED_FORMATS,
            "max_file_sizes": FileValidator.MAX_FILE_SIZE,
            "total_formats": len(FileValidator.SUPPORTED_FORMATS)
        }