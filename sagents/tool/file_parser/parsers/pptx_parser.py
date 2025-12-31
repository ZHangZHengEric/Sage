"""
PPTX文件解析器
支持PowerPoint文件的文本提取和元数据获取
"""

import traceback
from typing import Dict, Any
from pptx import Presentation
from .base_parser import BaseFileParser, ParseResult


class PPTXParser(BaseFileParser):
    """PPTX文件解析器"""
    
    SUPPORTED_EXTENSIONS = ['.pptx']
    SUPPORTED_MIME_TYPES = ['application/vnd.openxmlformats-officedocument.presentationml.presentation']
    
    def parse(self, file_path: str, skip_validation: bool = False) -> ParseResult:
        """
        解析PPTX文件
        
        Args:
            file_path: PPTX文件路径
            skip_validation: 是否跳过文件格式验证（can_parse检查）
            
        Returns:
            ParseResult: 解析结果
        """
        if not self.validate_file(file_path):
            return self.create_error_result(f"文件不存在或无法读取: {file_path}", file_path)
        
        if not skip_validation and not self.can_parse(file_path):
            return self.create_error_result(f"不支持的文件类型: {file_path}", file_path)
        
        try:
            # 加载演示文稿
            prs = Presentation(file_path)
            
            # 提取文本内容
            text_parts = []
            slide_texts = []
            
            for i, slide in enumerate(prs.slides, 1):
                slide_text = f"--- 幻灯片 {i} ---\\n"
                slide_content = []
                
                # 提取幻灯片中的文本
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())
                
                if slide_content:
                    slide_text += "\\n".join(slide_content)
                else:
                    slide_text += "(无文本内容)"
                
                slide_texts.append("\\n".join(slide_content))
                text_parts.append(slide_text)
            
            text = "\\n\\n".join(text_parts)
            
            # 获取基础文件元数据
            base_metadata = self.get_file_metadata(file_path)
            
            # 获取PPTX特定元数据
            pptx_metadata = self._extract_pptx_metadata(prs, slide_texts)
            
            # 合并元数据
            metadata = {**base_metadata, **pptx_metadata}
            
            # 添加文本统计信息
            metadata.update({
                "text_length": len(text),
                "character_count": len(text),
                "word_count": len(text.split()) if text else 0,
                "line_count": text.count('\\n') if text else 0
            })
            
            return ParseResult(
                text=text,
                metadata=metadata,
                success=True
            )
            
        except Exception as e:
            error_msg = f"PPTX解析失败: {str(e)}"
            print(error_msg)
            traceback.print_exc()
            return self.create_error_result(error_msg, file_path)
    
    def _extract_pptx_metadata(self, prs: Any, slide_texts: list) -> Dict[str, Any]:
        """
        提取PPTX特定元数据
        
        Args:
            prs: Presentation对象
            slide_texts: 幻灯片文本列表
            
        Returns:
            Dict[str, Any]: PPTX元数据
        """
        try:
            # 获取文档属性
            core_props = prs.core_properties
            
            # 计算统计信息
            slide_count = len(prs.slides)
            total_shapes = sum(len(slide.shapes) for slide in prs.slides)
            text_shapes = 0
            image_shapes = 0
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_shapes += 1
                    if hasattr(shape, "image"):
                        image_shapes += 1
            
            # 计算每张幻灯片的文本长度
            slide_text_lengths = [len(text) for text in slide_texts]
            
            metadata = {
                # 文档属性
                "title": getattr(core_props, 'title', '') or '',
                "author": getattr(core_props, 'author', '') or '',
                "subject": getattr(core_props, 'subject', '') or '',
                "keywords": getattr(core_props, 'keywords', '') or '',
                "comments": getattr(core_props, 'comments', '') or '',
                "category": getattr(core_props, 'category', '') or '',
                "created": str(getattr(core_props, 'created', '')) if getattr(core_props, 'created', None) else '',
                "modified": str(getattr(core_props, 'modified', '')) if getattr(core_props, 'modified', None) else '',
                "last_modified_by": getattr(core_props, 'last_modified_by', '') or '',
                "revision": getattr(core_props, 'revision', 0) or 0,
                
                # 演示文稿统计
                "slide_count": slide_count,
                "total_shapes": total_shapes,
                "text_shapes": text_shapes,
                "image_shapes": image_shapes,
                "slide_text_lengths": slide_text_lengths,
                "average_text_per_slide": sum(slide_text_lengths) / slide_count if slide_count > 0 else 0,
                
                # 幻灯片尺寸
                "slide_width": prs.slide_width,
                "slide_height": prs.slide_height,
                "slide_size_inches": {
                    "width": prs.slide_width / 914400,  # 转换为英寸
                    "height": prs.slide_height / 914400
                }
            }
            
            return metadata
            
        except Exception as e:
            print(f"提取PPTX元数据时出错: {e}")
            traceback.print_exc()
            return {"metadata_extraction_error": str(e)}