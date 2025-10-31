import os
import shutil
import tempfile
import hashlib
import mimetypes
import logging
import time
from datetime import datetime
from pathlib import Path
import urllib.parse
import requests
import asyncio
import stat
import platform
import zipfile
import tarfile
import re
import chardet
import traceback
from typing import Dict, Any, List, Optional, Union,Tuple
import pdfplumber
import pypandoc
from pptx import Presentation
import html2text
from openpyxl import load_workbook
from openpyxl.styles.numbers import is_date_format
from openpyxl.utils.datetime import from_excel, WINDOWS_EPOCH, MAC_EPOCH
import subprocess

from .tool_base import ToolBase
from sagents.utils.logger import logger

class FileParserError(Exception):
    """文件解析异常"""
    pass

class FileValidator:
    """文件验证器"""
    
    # 支持的文件类型和对应的MIME类型
    SUPPORTED_FORMATS = {
        '.pdf': 'application/pdf',
        '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        '.doc': 'application/msword',
        '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        '.ppt': 'application/vnd.ms-powerpoint',
        '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        '.xls': 'application/vnd.ms-excel',
        '.txt': 'text/plain',
        '.csv': 'text/csv',
        '.json': 'application/json',
        '.xml': 'application/xml',
        '.html': 'text/html',
        '.htm': 'text/html',
        '.md': 'text/markdown',
        '.markdown': 'text/markdown',
        '.rtf': 'application/rtf',
        '.odt': 'application/vnd.oasis.opendocument.text',
        '.epub': 'application/epub+zip',
        '.latex': 'application/x-latex',
        '.tex': 'application/x-tex'
    }
    
    # 文件大小限制 (MB)
    MAX_FILE_SIZE = {
        '.pdf': 50,
        '.docx': 25,
        '.doc': 25,
        '.pptx': 100,
        '.ppt': 100,
        '.xlsx': 25,
        '.xls': 25,
        '.txt': 10,
        '.csv': 50,
        '.json': 10,
        '.xml': 10,
        '.html': 5,
        '.htm': 5,
        '.md': 5,
        '.markdown': 5,
        '.rtf': 10,
        '.odt': 25,
        '.epub': 50,
        '.latex': 10,
        '.tex': 10
    }
    
    @staticmethod
    def validate_file(file_path: str) -> Dict[str, Any]:
        """验证文件的有效性"""
        try:
            path = Path(file_path)
            
            # 检查文件是否存在
            if not path.exists():
                return {"valid": False, "error": f"文件不存在: {file_path}"}
            
            # 检查是否为文件（非目录）
            if not path.is_file():
                return {"valid": False, "error": f"路径不是有效文件: {file_path}"}
            
            # 获取文件扩展名
            file_extension = path.suffix.lower()
            
            # 检查文件格式是否支持
            if file_extension not in FileValidator.SUPPORTED_FORMATS:
                return {"valid": False, "error": f"不支持的文件格式: {file_extension}"}
            
            # 检查文件大小
            file_size_mb = path.stat().st_size / (1024 * 1024)
            max_size = FileValidator.MAX_FILE_SIZE.get(file_extension, 10)
            
            if file_size_mb > max_size:
                return {"valid": False, "error": f"文件过大: {file_size_mb:.1f}MB > {max_size}MB"}
            
            # 检查文件是否可读
            try:
                with open(file_path, 'rb') as f:
                    f.read(1024)  # 尝试读取前1KB
            except PermissionError:
                return {"valid": False, "error": "文件无读取权限"}
            except Exception as e:
                return {"valid": False, "error": f"文件读取失败: {str(e)}"}
            
            return {
                "valid": True,
                "file_size_mb": file_size_mb,
                "file_extension": file_extension,
                "mime_type": FileValidator.SUPPORTED_FORMATS[file_extension]
            }
            
        except Exception as e:
            return {"valid": False, "error": f"文件验证失败: {str(e)}"}

class TextProcessor:
    """文本处理器"""
    
    @staticmethod
    def detect_encoding(file_path: str) -> str:
        """检测文件编码"""
        try:
            with open(file_path, 'rb') as f:
                raw_data = f.read(10000)  # 读取前10KB进行检测
                result = chardet.detect(raw_data)
                return result.get('encoding', 'utf-8')
        except Exception:
            return 'utf-8'
    
    @staticmethod
    def clean_text(text: str) -> str:
        """清理文本内容 - 优化版本：特别针对网页内容"""
        if not text:
            return ""
        
        import re
        
        # 移除常见的网页无关内容
        patterns_to_remove = [
            # 广告和推广内容
            r'广告|advertisement|sponsored|推广|赞助',
            # 导航和菜单
            r'首页|首页|home|导航|menu|导航栏|navbar',
            # 页脚信息
            r'版权所有|copyright|©|all rights reserved|隐私政策|privacy policy|使用条款|terms of use',
            # 社交媒体链接
            r'关注我们|follow us|facebook|twitter|instagram|linkedin|weibo|微信|微博',
            # 订阅和注册
            r'订阅|subscribe|注册|register|登录|login|sign up|sign in',
            # 搜索框
            r'搜索|search|查找|find',
            # 面包屑导航
            r'您的位置|当前位置|breadcrumb|面包屑',
            # 返回顶部
            r'返回顶部|back to top|回到顶部',
            # 分享按钮
            r'分享|share|转发|转发给朋友',
            # 评论相关
            r'评论|comment|留言|发表评论|write a comment',
            # 相关推荐
            r'相关推荐|related|推荐阅读|recommended|你可能还喜欢|you may also like',
            # 热门标签
            r'热门标签|popular tags|标签|tags',
            # 分页
            r'上一页|下一页|previous|next|第\d+页|page \d+',
            # 加载更多
            r'加载更多|load more|查看更多|view more',
            # 阅读量、点赞数等
            r'阅读\s*\d+|views?\s*\d+|点赞\s*\d+|likes?\s*\d+|评论\s*\d+|comments?\s*\d+',
            # 作者信息（保留作者名但移除其他信息）
            r'作者[：:]\s*[^\n]*\n',  # 作者行
            r'by\s+[^\n]*\n',  # 英文作者行
            # 来源信息
            r'来源[：:]\s*[^\n]*\n',  # 来源行
            r'source[：:]\s*[^\n]*\n',  # 英文来源行
        ]
        
        # 应用移除模式
        for pattern in patterns_to_remove:
            text = re.sub(pattern, '', text, flags=re.IGNORECASE)
        
        # 清理HTML标签残留
        text = re.sub(r'<[^>]+>', '', text)  # 移除HTML标签
        
        # 清理多余的空白字符
        text = re.sub(r'\n\s*\n', '\n\n', text)  # 多个连续换行符合并为双换行
        text = re.sub(r'[ \t]+', ' ', text)      # 多个连续空格合并为单个空格
        text = re.sub(r'\n +', '\n', text)       # 行首空格
        text = re.sub(r' +\n', '\n', text)       # 行尾空格
        
        # 清理空行
        text = re.sub(r'\n\s*\n\s*\n', '\n\n', text)
        
        # 清理重复的标点符号
        text = re.sub(r'[。！？，；：]{2,}', lambda m: m.group()[0], text)  # 中文标点
        text = re.sub(r'[.!?,;:]{2,}', lambda m: m.group()[0], text)       # 英文标点
        
        # 清理行首行尾空白
        text = text.strip()
        
        # 如果清理后内容过短，返回原始文本（避免过度清理）
        if len(text.strip()) < 50:
            # 只做基本的空白清理
            text = re.sub(r'\n\s*\n', '\n\n', text)
            text = re.sub(r'[ \t]+', ' ', text)
        text = text.strip()
        
        return text
    
    @staticmethod
    def truncate_text(text: str, start_index: int = 0, max_length: int = 5000) -> str:
        """安全地截取文本"""
        if not text:
            return ""
        
        start_index = max(0, start_index)
        end_index = min(len(text), start_index + max_length)
        
        return text[start_index:end_index]
    
    @staticmethod
    def get_text_stats(text: str) -> Dict[str, int]:
        """获取文本统计信息"""
        if not text:
            return {"characters": 0, "words": 0, "lines": 0, "paragraphs": 0}
        
        return {
            "characters": len(text),
            "words": len(text.split()),
            "lines": text.count('\n') + 1,
            "paragraphs": len([p for p in text.split('\n\n') if p.strip()])
        }

class PDFParser:
    """PDF解析器"""
    
    @staticmethod
    def extract_text(pdf_path: str) -> str:
        """从PDF提取文本"""
        try:
            with pdfplumber.open(pdf_path) as pdf:
                text_parts = []
                for page_num, page in enumerate(pdf.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text:
                            text_parts.append(f"=== 第 {page_num + 1} 页 ===\n{page_text}")
                    except Exception as e:
                        logger.warning(f"PDF第{page_num + 1}页解析失败: {e}")
                        text_parts.append(f"=== 第 {page_num + 1} 页 ===\n[页面解析失败: {str(e)}]")
                
                return "\n\n".join(text_parts)
                
        except Exception as e:
            raise FileParserError(f"PDF解析失败: {str(e)}")
    
    @staticmethod
    def get_pdf_info(pdf_path: str) -> Dict[str, Any]:
        """获取PDF信息"""
        try:
            with pdfplumber.open(pdf_path) as pdf:
                return {
                    "pages": len(pdf.pages),
                    "metadata": pdf.metadata or {}
                }
        except Exception as e:
            return {"pages": 0, "metadata": {}}

class OfficeParser:
    """Office文档解析器"""
    
    @staticmethod
    def extract_text_from_docx(file_path: str) -> str:
        """从DOCX提取文本"""
        try:
            return pypandoc.convert_file(file_path, 'markdown', extra_args=['--extract-media=.'])
        except Exception as e:
            raise FileParserError(f"DOCX解析失败: {str(e)}")
    
    @staticmethod
    def extract_text_from_doc(file_path: str) -> str:
        """从DOC提取文本（需要antiword）"""
        try:
            result = subprocess.run(
                ['antiword', file_path], 
                capture_output=True, 
                text=True, 
                timeout=30
            )
            if result.returncode == 0:
                return result.stdout
            else:
                raise FileParserError(f"DOC转换失败: {result.stderr}")
        except subprocess.TimeoutExpired:
            raise FileParserError("DOC解析超时")
        except FileNotFoundError:
            raise FileParserError("未找到antiword工具，请安装: sudo apt-get install antiword")
        except Exception as e:
            raise FileParserError(f"DOC解析失败: {str(e)}")
        
    @staticmethod
    def extract_text_from_pptx(file_path: str) -> str:
        """从PPTX提取文本"""
        try:
            prs = Presentation(file_path)
            slides_text = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_content = [f"=== 幻灯片 {slide_num + 1} ==="]
                
                # 按位置排序形状
                shapes = sorted(slide.shapes, key=lambda x: (x.top, x.left))
                
                for shape in shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        text_frame = shape.text_frame
                        full_text = []
                        for paragraph in text_frame.paragraphs:
                            para_text = "".join([run.text for run in paragraph.runs])
                            if para_text.strip():
                                full_text.append(para_text.strip())
                        if full_text:
                            slide_content.append('\n'.join(full_text))
                        else:
                            logger.debug(f"Shape {shape.name} has text_frame but no text extracted from paragraphs.")
                    elif hasattr(shape, "image"):
                        logger.debug(f"Shape {shape.name} is an image.")
                    elif shape.has_table:
                        table = shape.table
                        table_text = []
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                table_text.append(' | '.join(row_text))
                        if table_text:
                            slide_content.append('\n'.join(table_text))
                        else:
                            logger.debug(f"Shape {shape.name} is a table but no text extracted.")
                    else:
                        logger.debug(f"Shape {shape.name} is of type {shape.shape_type} and has no text attribute.")
                
                slides_text.append('\n'.join(slide_content))
            
            return '\n\n'.join(slides_text)
            
        except Exception as e:
            raise FileParserError(f"PPTX解析失败: {str(e)}")

class ExcelParser:
    """Excel解析器"""
    
    @staticmethod
    def extract_text_from_xlsx(file_path: str) -> Tuple[str, Dict[str, Any]]:
        """从Excel提取文本并转换为Markdown"""
        try:
            excel_data = ExcelParser._read_excel_to_dict(file_path)
            markdown_tables = []
            metadata = {}
            metadata['sheets'] = []
            for sheet_name, sheet_data in excel_data.items():
                # 限制行数
                # if len(sheet_data) > 100:
                #     sheet_data = sheet_data[:100]
                    
                sheet_md = ExcelParser._sheet_data_to_markdown(sheet_data, sheet_name)
                markdown_tables.append(sheet_md)
                metadata['sheets'].append(sheet_name)
                metadata['sheet_' + sheet_name] = {
                    'rows': len(sheet_data),
                    'columns': len(sheet_data[0]) if sheet_data else 0,
                    # 便于测试：提供表头信息
                    'headers': sheet_data[0] if sheet_data else []
                }
            return '\n\n'.join(markdown_tables), metadata
            
        except Exception as e:
            raise FileParserError(f"Excel解析失败: {str(e)}")
        
    @staticmethod
    def _read_excel_to_dict(file_path: str) -> Dict[str, List[List[str]]]:
        """读取Excel文件到字典，正确处理合并单元格"""
        # 调试：记录入口和文件路径
        try:
            logger.info(f"[xlsx-debug] 开始解析: {file_path}")
        except Exception:
            pass
        # 辅助：识别日期风格（包含中文/自定义格式）
        def _looks_like_date_format(fmt: str) -> bool:
            fmt = (fmt or '').lower()
            try:
                if is_date_format(fmt):
                    return True
            except Exception:
                pass
            # 兼容中文及常见日期标记
            return any(token in fmt for token in ['y', '年', '月', '日', 'h', '时', '分', '秒'])

        # 辅助：无格式时，按“可疑日期序列号”范围做兜底识别
        def _maybe_serial_date(value: float) -> Optional[str]:
            try:
                # 经验范围：Excel日期序列（天）通常在 25000~80000（约1968-2099）
                iv = float(value)
                if 25000 <= iv <= 80000:
                    dt = from_excel(iv, epoch=date_epoch)
                    # 防止误判：限定合理年份
                    if 1900 <= dt.year <= 2100:
                        return dt.strftime('%Y-%m-%d %H:%M:%S')
            except Exception:
                pass
            return None

        # （移除猜测式列识别与兜底转换）
        # 之前的横向“日期序列列”识别和基于范围的兜底转换可能带来误判，按用户要求去掉，
        # 仅使用样式/格式严格识别日期。

        # 将 Excel/WPS 的日期样式转换为输出字符串
        def _format_dt_by_code(dt, fmt_code: Optional[str]) -> str:
            if not fmt_code:
                return dt.strftime('%Y-%m-%d %H:%M:%S')
            fmt = fmt_code.replace('"', '')
            fmt_lower = fmt.lower()
            # 常见中文格式
            if '年' in fmt and '月' in fmt and ('日' not in fmt and 'd' not in fmt_lower):
                return f'{dt.year}年{dt.month}月'
            if '月' in fmt and ('日' in fmt or 'd' in fmt_lower):
                return f'{dt.month}月{dt.day}日'
            # 常见西式日期
            if 'yyyy' in fmt_lower and ('mm' in fmt_lower or 'm' in fmt_lower):
                sep = '-' if '-' in fmt else '/' if '/' in fmt else ' '
                if 'dd' in fmt_lower or 'd' in fmt_lower:
                    return dt.strftime(f'%Y{sep}%m{sep}%d')
                else:
                    return dt.strftime(f'%Y{sep}%m')
            # 时间部分
            if any(t in fmt_lower for t in ['hh', 'h', 'ss']):
                return dt.strftime('%Y-%m-%d %H:%M:%S')
            return dt.strftime('%Y-%m-%d %H:%M:%S')

        # 样式级日期识别：直接读取内置/自定义 numFmtId 判定是否日期
        # openpyxl 对部分东亚内置ID（如57/58）不一定映射到 number_format，但样式仍可用
        def _build_style_numfmt_map(xlsx_path: str):
            try:
                # 直接解析底层 XML，提高兼容性
                import zipfile, xml.etree.ElementTree as ET
                zf = zipfile.ZipFile(xlsx_path)
                styles_xml = zf.read('xl/styles.xml')
                styles = ET.fromstring(styles_xml)
                ns = {'main': 'http://schemas.openxmlformats.org/spreadsheetml/2006/main'}
                # 收集自定义格式
                custom_numfmts = {}
                for nf in styles.findall('.//main:numFmts/main:numFmt', ns):
                    numFmtId = int(nf.get('numFmtId'))
                    fmt = nf.get('formatCode') or ''
                    custom_numfmts[numFmtId] = fmt
                # 样式索引 -> numFmtId
                style_to_numfmt = []
                for xf in styles.findall('.//main:cellXfs/main:xf', ns):
                    numFmtId = xf.get('numFmtId')
                    numFmtId = int(numFmtId) if numFmtId is not None else None
                    style_to_numfmt.append(numFmtId)
                # 标准日期/时间内置ID集合（Excel/WPS常见）
                builtin_date_ids = {
                    14, 15, 16, 17, 18, 19, 20, 21, 22, 45, 46, 47, 57, 58
                }
                builtin_fmt_map = {
                    14: 'm/d/yy',
                    15: 'd-mmm-yy',
                    16: 'd-mmm',
                    17: 'mmm-yy',
                    18: 'h:mm AM/PM',
                    19: 'h:mm:ss AM/PM',
                    20: 'h:mm',
                    21: 'h:mm:ss',
                    22: 'm/d/yy h:mm',
                    45: 'mm:ss',
                    46: '[h]:mm:ss',
                    47: 'mmss.0',
                    57: 'yyyy"年"m"月"',
                    58: 'm"月"d"日"'
                }
                def is_date_numfmt_id(numFmtId: Optional[int]) -> bool:
                    if numFmtId is None:
                        return False
                    if numFmtId in builtin_date_ids:
                        return True
                    # 自定义格式的日期判断
                    fmt = custom_numfmts.get(numFmtId, '')
                    return _looks_like_date_format(fmt)
                def get_numfmt_code(numFmtId: Optional[int]) -> Optional[str]:
                    if numFmtId is None:
                        return None
                    return custom_numfmts.get(numFmtId) or builtin_fmt_map.get(numFmtId)
                return style_to_numfmt, is_date_numfmt_id, get_numfmt_code
            except Exception:
                # 解析失败则返回空映射
                return [], (lambda _id: False), (lambda _id: None)
        # 需要关闭read_only模式才能访问合并单元格信息
        # data_only=True 读取公式的缓存结果；若Excel未计算过，会出现None
        # 同时加载一个 data_only=False 的工作簿用于在缓存缺失时回退到公式文本
        workbook = load_workbook(file_path, data_only=True, read_only=False)
        workbook_formula = load_workbook(file_path, data_only=False, read_only=False)
        # 判断工作簿的日期基准（Windows 1900 或 macOS 1904）
        date1904 = getattr(workbook.properties, 'date1904', False)
        date_epoch = MAC_EPOCH if date1904 else WINDOWS_EPOCH
        try:
            logger.debug(f"[xlsx-debug] 工作簿属性: date1904={date1904}")
        except Exception:
            pass
        excel_data = {}

        # 构建样式到 numFmtId 的映射供单元格样式级判断
        style_to_numfmt, is_date_numfmt_id, get_numfmt_code = _build_style_numfmt_map(file_path)

        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_f = workbook_formula[sheet_name]
            try:
                logger.debug(f"[xlsx-debug] 处理工作表: {sheet_name}, max_row={sheet.max_row}, max_col={sheet.max_column}")
            except Exception:
                pass
            
            # 创建合并单元格值映射
            merged_cell_values = {}
            merged_cell_formats = {}
            merged_cell_is_date = {}
            for merged_range in sheet.merged_cells.ranges:
                # 获取合并单元格左上角的值
                top_left_cell = sheet.cell(merged_range.min_row, merged_range.min_col)
                value = top_left_cell.value
                # 若缓存结果为空，则回退为公式文本（若存在）
                if value is None:
                    try:
                        value = sheet_f.cell(merged_range.min_row, merged_range.min_col).value
                    except Exception:
                        value = None
                
                # 为合并范围内的所有单元格设置相同的值
                for row in range(merged_range.min_row, merged_range.max_row + 1):
                    for col in range(merged_range.min_col, merged_range.max_col + 1):
                        merged_cell_values[(row, col)] = value
                        merged_cell_formats[(row, col)] = getattr(top_left_cell, 'number_format', '')
                        merged_cell_is_date[(row, col)] = getattr(top_left_cell, 'is_date', False)
            
            # 移除猜测列识别，不进行兜底列转换

            # 读取数据，考虑合并单元格
            sheet_data = []
            max_row = sheet.max_row
            max_col = sheet.max_column
            
            if max_row and max_col:
                for row_idx in range(1, max_row + 1):
                    row_data = []
                    for col_idx in range(1, max_col + 1):
                        # 检查是否是合并单元格
                        if (row_idx, col_idx) in merged_cell_values:
                            cell_value = merged_cell_values[(row_idx, col_idx)]
                            fmt = merged_cell_formats.get((row_idx, col_idx), '')
                            is_date_flag = merged_cell_is_date.get((row_idx, col_idx), False)
                        else:
                            cell = sheet.cell(row_idx, col_idx)
                            cell_f = sheet_f.cell(row_idx, col_idx)
                            cell_value = cell.value
                            # 优先使用公式工作簿的样式信息，避免data_only模式下样式缺失/差异
                            fmt = getattr(cell, 'number_format', '')
                            is_date_flag = getattr(cell, 'is_date', False)
                            numFmtId = None
                            try:
                                style_id = getattr(cell_f, 'style_id', getattr(cell, 'style_id', None))
                                if isinstance(style_id, int) and 0 <= style_id < len(style_to_numfmt):
                                    numFmtId = style_to_numfmt[style_id]
                                    if is_date_numfmt_id(numFmtId):
                                        is_date_flag = True
                                        fmt = get_numfmt_code(numFmtId) or fmt
                                    else:
                                        # 非日期格式，但若存在自定义格式码，按其输出
                                        maybe_fmt = get_numfmt_code(numFmtId)
                                        fmt = maybe_fmt or fmt
                            except Exception:
                                pass
                            # 调试：打印识别前的单元格信息（限制前5行以避免日志过多）
                            try:
                                if row_idx <= 5:
                                    logger.debug(f"[xlsx-debug] r={row_idx} c={col_idx} raw='{cell_value}' type={type(cell_value).__name__} style_id={getattr(cell_f, 'style_id', getattr(cell,'style_id', None))} numFmtId={numFmtId} fmt='{fmt}' is_date_flag={is_date_flag}")
                            except Exception:
                                pass
                            # 若缓存结果为空，则回退为公式文本（若存在）
                            if cell_value is None:
                                try:
                                    cell_value = sheet_f.cell(row_idx, col_idx).value
                                except Exception:
                                    cell_value = None
                            # 规范日期类型为可阅读字符串
                            try:
                                from datetime import datetime
                                if isinstance(cell_value, datetime):
                                    # 改为按样式格式化输出
                                    cell_value = _format_dt_by_code(cell_value, fmt)
                                    try:
                                        if row_idx <= 5:
                                            logger.debug(f"[xlsx-debug] r={row_idx} c={col_idx} datetime格式化 => '{cell_value}'")
                                    except Exception:
                                        pass
                                else:
                                    # Excel中的日期常以序列号（float/int）存储，结合单元格格式识别为日期
                                    if isinstance(cell_value, (int, float)):
                                        # 使用单元格的 is_date / 样式 或自定义格式识别
                                        if is_date_flag or _looks_like_date_format(fmt):
                                            try:
                                                dt = from_excel(float(cell_value), epoch=date_epoch)
                                                cell_value = _format_dt_by_code(dt, fmt)
                                                try:
                                                    if row_idx <= 5:
                                                        logger.debug(f"[xlsx-debug] r={row_idx} c={col_idx} serial=>date({float(cell_value) if isinstance(cell_value,(int,float)) else 'n/a'}) fmt='{fmt}' => '{cell_value}'")
                                                except Exception:
                                                    pass
                                            except Exception:
                                                # 转换失败则保留原值
                                                pass
                            except Exception:
                                pass
                        
                        cell_str = str(cell_value).replace('\n', '\\n') if cell_value is not None else ''
                        row_data.append(cell_str)
                    
                    sheet_data.append(row_data)
                
                if not sheet_data:
                    continue
                
                # 清理空行和空列
                sheet_data = ExcelParser._clean_empty_rows_cols(sheet_data)

                if sheet_data:
                    excel_data[sheet_name] = sheet_data
        
        workbook.close()
        workbook_formula.close()
        try:
            logger.info(f"[xlsx-debug] 解析完成: {file_path}")
        except Exception:
            pass
        return excel_data
    
    @staticmethod
    def _clean_empty_rows_cols(data: List[List[str]]) -> List[List[str]]:
        """清理空行和空列"""
        if not data:
            return data
        
        # 移除空行
        data = [row for row in data if any(cell.strip() for cell in row)]
        
        if not data:
            return data
        
        # 移除空列
        cols_to_keep = []
        for col_idx in range(len(data[0])):
            if any(row[col_idx].strip() for row in data):
                cols_to_keep.append(col_idx)
        
        if cols_to_keep:
            data = [[row[i] for i in cols_to_keep] for row in data]
        
        return data
    
    @staticmethod
    def _sheet_data_to_markdown(sheet_data: List[List[str]], sheet_name: str) -> str:
        """转换工作表数据为Markdown"""
        if not sheet_data:
            return f'## {sheet_name}\n\n(空工作表)'
        
        markdown_lines = [f'## {sheet_name}', '']
        
        # 如果有数据，第一行作为表头
        if sheet_data:
            # 表头
            header = '| ' + ' | '.join(cell if cell else ' ' for cell in sheet_data[0]) + ' |'
            markdown_lines.append(header)
            
            # 分隔线
            separator = '| ' + ' | '.join('---' for _ in sheet_data[0]) + ' |'
            markdown_lines.append(separator)
            
            # 数据行
            for row in sheet_data[1:]:
                row_md = '| ' + ' | '.join(cell if cell else ' ' for cell in row) + ' |'
                markdown_lines.append(row_md)
        
        return '\n'.join(markdown_lines)

class WebParser:
    """网页解析器"""
    
    @staticmethod
    def extract_text_from_html(file_path: str) -> str:
        """从HTML文件提取文本"""
        try:
            encoding = TextProcessor.detect_encoding(file_path)
            with open(file_path, 'r', encoding=encoding) as file:
                html_content = file.read()
            
            return WebParser._html_to_text(html_content)
            
        except Exception as e:
            raise FileParserError(f"HTML解析失败: {str(e)}")
    
    @staticmethod
    def extract_text_from_url(url: str, timeout: int = 30) -> str:
        """从URL提取文本"""
        max_retries = 3
        initial_retry_delay = 2  # 初始重试延迟增加到2秒
        
        # 根据URL特征调整超时时间和策略
        adjusted_timeout = timeout
        is_github = 'github.com' in url or 'githubusercontent.com' in url
        
        if is_github:
            adjusted_timeout = max(timeout, 45)  # GitHub至少45秒
        elif any(domain in url for domain in ['google.com', 'microsoft.com', 'amazon.com']):
            adjusted_timeout = max(timeout, 40)  # 大型网站至少40秒
        
        # 创建会话对象复用连接
        session = requests.Session()
        
        # 针对GitHub使用更强的反检测策略
        if is_github:
            session.headers.update({
                'User-Agent': 'Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8',
                'Accept-Language': 'en-US,en;q=0.9',
                'Accept-Encoding': 'gzip, deflate, br',
                'Cache-Control': 'no-cache',
                'Pragma': 'no-cache',
                'Sec-Ch-Ua': '"Not_A Brand";v="8", "Chromium";v="120", "Google Chrome";v="120"',
                'Sec-Ch-Ua-Mobile': '?0',
                'Sec-Ch-Ua-Platform': '"macOS"',
                'Sec-Fetch-Dest': 'document',
                'Sec-Fetch-Mode': 'navigate',
                'Sec-Fetch-Site': 'none',
                'Sec-Fetch-User': '?1',
                'Upgrade-Insecure-Requests': '1'
            })
        else:
            # 普通网站使用标准请求头
            session.headers.update({
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8,application/signed-exchange;v=b3;q=0.7',
                'Accept-Language': 'zh-CN,zh;q=0.9,en;q=0.8',
                'Accept-Encoding': 'gzip, deflate, br',
                'DNT': '1',
                'Connection': 'keep-alive',
                'Upgrade-Insecure-Requests': '1',
                'Sec-Fetch-Dest': 'document',
                'Sec-Fetch-Mode': 'navigate',
                'Sec-Fetch-Site': 'none',
                'Sec-Fetch-User': '?1'
            })
        
        try:
            for attempt in range(max_retries):
                retry_delay = initial_retry_delay * (2 ** attempt)  # 指数退避：2, 4, 8秒
                
                # 对于GitHub，在首次请求前添加随机延迟，模拟人类行为
                if is_github and attempt == 0:
                    import random
                    human_delay = random.uniform(1, 3)  # 1-3秒随机延迟
                    logger.debug(f"🤖 GitHub反检测：等待{human_delay:.1f}秒模拟人类行为")
                    time.sleep(human_delay)
                
                try:
                    logger.debug(f"🌐 尝试第{attempt + 1}次访问URL: {url} (超时: {adjusted_timeout}秒)")
                    
                    # 使用会话发起请求，设置连接和读取超时
                    response = session.get(
                        url, 
                        timeout=(10, adjusted_timeout),  # (连接超时, 读取超时)
                        allow_redirects=True,
                        stream=False  # 不使用流式下载，确保完整获取内容
                    )
                    
                    # 检查响应状态
                    if response.status_code == 404:
                        raise FileParserError(f"URL不存在 (404): {url}")
                    elif response.status_code == 403:
                        raise FileParserError(f"访问被禁止 (403): {url}")
                    elif response.status_code == 500:
                        raise FileParserError(f"服务器内部错误 (500): {url}")
                    elif response.status_code == 502:
                        # 502错误通常是暂时的，值得重试
                        if attempt < max_retries - 1:
                            logger.warning(f"⚠️ 服务器网关错误 (502)，第{attempt + 1}次重试，{retry_delay}秒后重试")
                            time.sleep(retry_delay)
                            continue
                        else:
                            raise FileParserError(f"服务器网关错误 (502): {url}")
                    elif response.status_code == 503:
                        # 503服务不可用，值得重试
                        if attempt < max_retries - 1:
                            logger.warning(f"⚠️ 服务暂时不可用 (503)，第{attempt + 1}次重试，{retry_delay}秒后重试")
                            time.sleep(retry_delay)
                            continue
                        else:
                            raise FileParserError(f"服务暂时不可用 (503): {url}")
                    elif response.status_code >= 400:
                        # 对于其他4xx和5xx错误，如果不是最后一次尝试，则重试
                        if attempt < max_retries - 1:
                            logger.warning(f"⚠️ HTTP错误 {response.status_code}，第{attempt + 1}次重试，{retry_delay}秒后重试")
                            time.sleep(retry_delay)
                            continue
                        else:
                            raise FileParserError(f"URL访问失败 (HTTP {response.status_code}): {url}")
                    
                    response.raise_for_status()
                    
                    # 检查内容长度
                    if len(response.content) == 0:
                        if attempt < max_retries - 1:
                            logger.warning(f"⚠️ 响应内容为空，第{attempt + 1}次重试，{retry_delay}秒后重试")
                            time.sleep(retry_delay)
                            continue
                        else:
                            raise FileParserError(f"响应内容为空: {url}")
                    
                    # 检查内容类型
                    content_type = response.headers.get('Content-Type', '').lower()
                    if 'text/html' not in content_type and 'text/plain' not in content_type and 'application/xhtml' not in content_type:
                        logger.warning(f"⚠️ 检测到非文本内容类型: {content_type}，仍尝试解析")
                    
                    # 智能编码检测和处理
                    final_text = WebParser._smart_decode_response(response, url)
                    
                    logger.debug(f"✅ 成功获取内容，长度: {len(final_text)} 字符")
                    return WebParser._html_to_text(final_text)
                    
                except requests.exceptions.ConnectTimeout as e:
                    if attempt < max_retries - 1:
                        logger.warning(f"⚠️ 连接超时，第{attempt + 1}次重试，{retry_delay}秒后重试: {str(e)}")
                        time.sleep(retry_delay)
                        continue
                    else:
                        raise FileParserError(f"连接超时: {url} (连接超时: 10秒)")
                        
                except requests.exceptions.ReadTimeout as e:
                    if attempt < max_retries - 1:
                        # 对于读取超时，增加下次尝试的超时时间
                        adjusted_timeout = min(adjusted_timeout * 1.5, 90)  # 最大90秒
                        logger.warning(f"⚠️ 读取超时，第{attempt + 1}次重试，{retry_delay}秒后重试，调整超时至{adjusted_timeout:.0f}秒: {str(e)}")
                        time.sleep(retry_delay)
                        continue
                    else:
                        raise FileParserError(f"读取超时: {url} (读取超时: {adjusted_timeout:.0f}秒)")
                        
                except requests.exceptions.Timeout as e:
                    if attempt < max_retries - 1:
                        logger.warning(f"⚠️ 请求超时，第{attempt + 1}次重试，{retry_delay}秒后重试: {str(e)}")
                        time.sleep(retry_delay)
                        continue
                    else:
                        raise FileParserError(f"请求超时: {url} (超时时间: {adjusted_timeout}秒)")
                        
                except requests.exceptions.ConnectionError as e:
                    if attempt < max_retries - 1:
                        logger.warning(f"⚠️ 连接错误，第{attempt + 1}次重试，{retry_delay}秒后重试: {str(e)}")
                        time.sleep(retry_delay)
                        continue
                    else:
                        raise FileParserError(f"无法连接到URL: {url} - {str(e)}")
                        
                except requests.exceptions.RequestException as e:
                    # 对于其他请求异常，根据错误类型决定是否重试
                    error_str = str(e).lower()
                    if any(keyword in error_str for keyword in ['timeout', 'connection', 'network', 'dns']):
                        if attempt < max_retries - 1:
                            logger.warning(f"⚠️ 网络相关错误，第{attempt + 1}次重试，{retry_delay}秒后重试: {str(e)}")
                            time.sleep(retry_delay)
                            continue
                    raise FileParserError(f"URL访问失败: {url} - {str(e)}")
                    
                except Exception as e:
                    # 对于其他异常，记录详细信息
                    logger.error(f"💥 意外异常: {type(e).__name__}: {str(e)}")
                    raise FileParserError(f"URL解析失败: {url} - {str(e)}")
            
            # 如果所有重试都失败了（理论上不会到达这里）
            raise FileParserError(f"URL访问失败，已重试{max_retries}次: {url}")
            
        finally:
            # 确保关闭会话
            session.close()
    
    @staticmethod
    def _smart_decode_response(response, url: str) -> str:
        """智能解码HTTP响应内容"""
        import chardet
        import re
        
        # 1. 首先尝试从Content-Type header获取编码
        content_type = response.headers.get('Content-Type', '')
        charset_match = re.search(r'charset=([^;\s]+)', content_type, re.IGNORECASE)
        header_encoding = None
        if charset_match:
            header_encoding = charset_match.group(1).strip('"\'').lower()
            logger.debug(f"🔍 从HTTP头获取编码: {header_encoding}")
        
        # 2. 尝试从HTML的meta标签获取编码
        html_encoding = None
        try:
            # 用原始字节数据检查HTML meta标签
            content_preview = response.content[:2048]  # 只检查前2KB
            content_str = content_preview.decode('ascii', errors='ignore')
            
            # 查找meta charset
            meta_matches = [
                re.search(r'<meta[^>]+charset["\s]*=["\s]*([^">\s]+)', content_str, re.IGNORECASE),
                re.search(r'<meta[^>]+content[^>]*charset=([^">\s;]+)', content_str, re.IGNORECASE)
            ]
            
            for match in meta_matches:
                if match:
                    html_encoding = match.group(1).strip('"\'').lower()
                    logger.debug(f"🔍 从HTML meta标签获取编码: {html_encoding}")
                    break
        except Exception as e:
            logger.debug(f"⚠️ HTML编码检测失败: {e}")
        
        # 3. 使用chardet进行自动检测
        detected_encoding = None
        try:
            detection_result = chardet.detect(response.content)
            if detection_result and detection_result.get('confidence', 0) > 0.7:
                detected_encoding = detection_result['encoding'].lower()
                confidence = detection_result.get('confidence', 0)
                logger.debug(f"🔍 chardet检测编码: {detected_encoding} (置信度: {confidence:.2f})")
        except Exception as e:
            logger.debug(f"⚠️ chardet编码检测失败: {e}")
        
        # 4. 根据URL判断可能的编码（针对中文网站）
        url_encoding = None
        if any(domain in url.lower() for domain in ['.cn', '.com.cn', 'baidu', 'sina', 'qq', '163', 'sohu', 'taobao', 'jd']):
            url_encoding = 'utf-8'  # 大多数中文网站使用UTF-8
            logger.debug(f"🔍 根据URL推测中文网站编码: {url_encoding}")
        
        # 5. 编码优先级和尝试顺序
        encoding_candidates = []
        
        # 优先级排序
        if header_encoding:
            encoding_candidates.append(header_encoding)
        if html_encoding and html_encoding != header_encoding:
            encoding_candidates.append(html_encoding)
        if detected_encoding and detected_encoding not in encoding_candidates:
            encoding_candidates.append(detected_encoding)
        if url_encoding and url_encoding not in encoding_candidates:
            encoding_candidates.append(url_encoding)
        
        # 添加常见编码作为后备
        fallback_encodings = ['utf-8', 'gbk', 'gb2312', 'big5', 'iso-8859-1', 'windows-1252']
        for enc in fallback_encodings:
            if enc not in encoding_candidates:
                encoding_candidates.append(enc)
        
        # 6. 逐个尝试解码
        for encoding in encoding_candidates:
            try:
                if encoding:
                    # 规范化编码名称
                    encoding = encoding.replace('gb2312', 'gbk')  # gb2312是gbk的子集
                    encoding = encoding.replace('iso-8859-1', 'latin1')
                    
                    decoded_text = response.content.decode(encoding, errors='replace')
                    
                    # 简单验证：检查是否包含大量替换字符
                    replacement_ratio = decoded_text.count('�') / max(len(decoded_text), 1)
                    if replacement_ratio < 0.1:  # 替换字符少于10%
                        logger.debug(f"✅ 成功使用编码 {encoding} 解码内容")
                        return decoded_text
                    else:
                        logger.debug(f"⚠️ 编码 {encoding} 产生过多替换字符 ({replacement_ratio:.2%})")
            except (UnicodeDecodeError, LookupError) as e:
                logger.debug(f"⚠️ 编码 {encoding} 解码失败: {e}")
                continue
        
        # 7. 最后的fallback：使用UTF-8并忽略错误
        logger.warning(f"⚠️ 所有编码尝试失败，使用UTF-8强制解码: {url}")
        return response.content.decode('utf-8', errors='replace')
    
    @staticmethod
    def _html_to_text(html_content: str) -> str:
        """HTML转文本 - 优化版本：保留纯文字和超链接，忽略图片"""
        h = html2text.HTML2Text()
        h.ignore_links = False  # 保留超链接
        h.ignore_images = True  # 忽略图片
        h.ignore_emphasis = False  # 保留强调格式（粗体、斜体等）
        h.body_width = 0  # 不限制行宽
        h.unicode_snob = True  # 使用Unicode字符
        h.escape_snob = True  # 转义特殊字符
        h.skip_internal_links = False  # 不跳过内部链接
        h.inline_links = True  # 使用内联链接格式
        h.protect_links = True  # 保护链接不被破坏
        h.mark_code = True  # 标记代码块
        
        # 转换HTML为文本
        text = h.handle(html_content)
        
        # 后处理：清理多余的空白和格式化
        import re
        
        # 清理多余的换行符
        text = re.sub(r'\n{3,}', '\n\n', text)
        
        # 清理图片相关的残留文本（即使ignore_images=True，有时仍会有残留）
        text = re.sub(r'!\[.*?\]\(.*?\)', '', text)  # 移除markdown图片语法
        text = re.sub(r'<img[^>]*>', '', text)  # 移除HTML img标签
        text = re.sub(r'图片|image|img', '', text, flags=re.IGNORECASE)  # 移除图片相关文字
        
        # 清理多余的空白字符
        text = re.sub(r'[ \t]+', ' ', text)  # 多个空格合并为一个
        text = re.sub(r'\n +', '\n', text)  # 行首空格
        text = re.sub(r' +\n', '\n', text)  # 行尾空格
        
        # 清理空行
        text = re.sub(r'\n\s*\n\s*\n', '\n\n', text)
        
        # 确保文本开头和结尾没有多余空白
        text = text.strip()
        
        return text

class PlainTextParser:
    """纯文本解析器"""
    
    @staticmethod
    def extract_text_from_plain_file(file_path: str) -> str:
        """从纯文本文件提取内容"""
        try:
            encoding = TextProcessor.detect_encoding(file_path)
            with open(file_path, 'r', encoding=encoding) as file:
                return file.read()
        except Exception as e:
            raise FileParserError(f"文本文件解析失败: {str(e)}")
    
    @staticmethod
    def extract_text_with_pandoc(file_path: str, input_format: str = None) -> str:
        """使用Pandoc提取文本"""
        try:
            if input_format:
                return pypandoc.convert_file(file_path, 'markdown', format=input_format)
            else:
                return pypandoc.convert_file(file_path, 'markdown')
        except Exception as e:
            raise FileParserError(f"Pandoc解析失败: {str(e)}")

class FileParserTool(ToolBase):
    """文件解析工具集"""
    
    def __init__(self):
        logger.debug("Initializing FileParserTool")
        super().__init__()


    @ToolBase.tool()
    def extract_text_from_non_text_file(
        self, 
        input_file_path: str, 
        start_index: int = 0, 
        max_length: int = 5000,
        include_metadata: bool = True
    ) -> Dict[str, Any]:
        """读取本地存储下的非文本文件，例如pdf，docx，doc，ppt，pptx，xlsx，xls，csv等文件，返回Markdown的文本数据

        Args:
            input_file_path (str): 输入文件路径，本地的绝对路径
            start_index (int): 开始提取的字符位置，默认0
            max_length (int): 单次最大提取长度，默认5000字符，最大5000字符
            include_metadata (bool): 是否包含文件元数据，默认True

        Returns:
            Dict[str, Any]: 包含提取文本和相关信息的字典
        """
        if max_length > 5000:
            max_length = 5000
        start_time = time.time()
        operation_id = hashlib.md5(f"extract_{input_file_path}_{time.time()}".encode()).hexdigest()[:8]
        logger.info(f"📄 extract_text_from_file开始执行 [{operation_id}] - 文件: {input_file_path}")
        logger.info(f"🔧 参数: start_index={start_index}, max_length={max_length}, include_metadata={include_metadata}")
        
        try:
            # 验证文件
            logger.debug(f"🔍 开始文件验证")
            validation_result = FileValidator.validate_file(input_file_path)
            
            if not validation_result["valid"]:
                error_time = time.time() - start_time
                logger.error(f"❌ 文件验证失败 [{operation_id}] - 错误: {validation_result['error']}, 耗时: {error_time:.2f}秒")
                return {
                    "success": False,
                    "error": validation_result["error"],
                    "file_path": input_file_path,
                    "execution_time": error_time,
                    "operation_id": operation_id
                }
            
            file_extension = validation_result["file_extension"]
            file_size_mb = validation_result["file_size_mb"]
            logger.info(f"✅ 文件验证通过 [{operation_id}] - 格式: {file_extension}, 大小: {file_size_mb:.2f}MB")
            
            # 根据文件类型选择解析器
            parse_start_time = time.time()
            logger.info(f"🔧 开始文件解析 - 格式: {file_extension}")
            
            extracted_text = ""
            metadata = {}
            
            try:
                if file_extension == '.pdf':
                    logger.debug(f"📕 使用PDF解析器")
                    extracted_text = PDFParser.extract_text(input_file_path)
                    if include_metadata:
                        metadata = PDFParser.get_pdf_info(input_file_path)
                        
                elif file_extension in ['.docx', '.doc']:
                    logger.debug(f"📝 使用Word解析器")
                    if file_extension == '.docx':
                        extracted_text = OfficeParser.extract_text_from_docx(input_file_path)
                    else:
                        extracted_text = OfficeParser.extract_text_from_doc(input_file_path)
                        
                elif file_extension in ['.pptx', '.ppt']:
                    logger.debug(f"📊 使用PowerPoint解析器")
                    if file_extension == '.pptx':
                        extracted_text = OfficeParser.extract_text_from_pptx(input_file_path)
                    else:
                        # PPT需要额外的处理，尝试使用LibreOffice转换为PPTX
                        pptx_output_path = input_file_path + 'x'
                        try:
                            # 检查libreoffice是否安装
                            subprocess.run(['libreoffice', '--version'], check=True, capture_output=True)
                            logger.info(f"尝试使用LibreOffice将PPT转换为PPTX: {input_file_path} -> {pptx_output_path}")
                            command = [
                                'libreoffice',
                                '--headless',
                                '--convert-to', 'pptx',
                                '--outdir', os.path.dirname(input_file_path),
                                input_file_path
                            ]
                            result = subprocess.run(command, capture_output=True, text=True, check=True)
                            logger.info(f"LibreOffice转换输出: {result.stdout}")
                            logger.error(f"LibreOffice转换错误输出: {result.stderr}")

                            if os.path.exists(pptx_output_path):
                                extracted_text = OfficeParser.extract_text_from_pptx(pptx_output_path)
                                logger.info(f"PPT成功转换为PPTX并提取内容: {pptx_output_path}")
                                # 转换成功后删除临时生成的pptx文件
                                os.remove(pptx_output_path)
                            else:
                                extracted_text = "PPT文件转换失败，无法生成PPTX文件。请手动转换为PPTX格式。"
                                logger.error(f"PPT文件转换失败，未找到生成的PPTX文件: {pptx_output_path}")
                        except FileNotFoundError:
                            extracted_text = "LibreOffice未安装或不在PATH中，无法自动转换PPT文件。请安装LibreOffice或手动转换为PPTX格式。"
                            logger.error("LibreOffice未安装或不在PATH中，无法自动转换PPT文件。")
                        except subprocess.CalledProcessError as e:
                            extracted_text = f"LibreOffice转换PPT文件时出错: {e.stderr}"
                            logger.error(f"LibreOffice转换PPT文件时出错: {e.stderr}")
                        except Exception as e:
                            extracted_text = f"处理PPT文件时发生未知错误: {e}"
                            logger.error(f"处理PPT文件时发生未知错误: {e}")
                        
                elif file_extension in ['.xlsx', '.xls']:
                    logger.debug(f"📈 使用Excel解析器")
                    extracted_text, excel_metadata = ExcelParser.extract_text_from_xlsx(input_file_path)
                    if include_metadata:
                        metadata.update(excel_metadata)
                elif file_extension in ['.html', '.htm']:
                    logger.debug(f"🌐 使用HTML解析器")
                    extracted_text = WebParser.extract_text_from_html(input_file_path)
                    
                elif file_extension in ['.txt', '.csv', '.json', '.xml', '.md', '.markdown']:
                    logger.debug(f"📄 使用纯文本解析器")
                    extracted_text = PlainTextParser.extract_text_from_plain_file(input_file_path)
                    
                else:
                    # 尝试使用Pandoc解析
                    logger.debug(f"🔧 尝试使用Pandoc解析器")
                    extracted_text = PlainTextParser.extract_text_with_pandoc(input_file_path)
                
                parse_time = time.time() - parse_start_time
                logger.info(f"✅ 文件解析成功 [{operation_id}] - 原始文本长度: {len(extracted_text)}, 解析耗时: {parse_time:.2f}秒")
                
            except Exception as parse_error:
                parse_time = time.time() - parse_start_time
                logger.warning(f"⚠️ 主解析器失败 [{operation_id}] - 错误: {str(parse_error)}, 耗时: {parse_time:.2f}秒")
                logger.debug(f"🔧 尝试Pandoc备用解析器")
                
                try:
                    extracted_text = PlainTextParser.extract_text_with_pandoc(input_file_path)
                    logger.info(f"✅ Pandoc备用解析成功 [{operation_id}] - 文本长度: {len(extracted_text)}")
                except Exception as fallback_error:
                    error_time = time.time() - start_time
                    logger.error(f"❌ 所有解析器都失败 [{operation_id}] - 主错误: {str(parse_error)}, 备用错误: {str(fallback_error)}, 耗时: {error_time:.2f}秒")
                    return {
                        "success": False,
                        "error": f"文件解析失败 - 主错误: {str(parse_error)}, 备用错误: {str(fallback_error)}",
                        "file_path": input_file_path,
                        "execution_time": error_time,
                        "operation_id": operation_id
                    }
            
            # 清理和处理文本
            logger.debug(f"🧹 开始文本清理和处理")
            if file_extension in ['.xlsx', '.xls']:
                cleaned_text = extracted_text
                truncated_text = extracted_text
                text_stats = TextProcessor.get_text_stats(cleaned_text)
            else:
                cleaned_text = TextProcessor.clean_text(extracted_text)
                truncated_text = TextProcessor.truncate_text(cleaned_text, start_index, max_length)
                text_stats = TextProcessor.get_text_stats(cleaned_text)
                logger.info(f"✅ 文本提取完成 [{operation_id}] - 清理后长度: {len(cleaned_text)}, 截取长度: {len(truncated_text)}, 总耗时: {total_time:.2f}秒")
                logger.debug(f"📊 文本统计: {text_stats}")
            total_time = time.time() - start_time
            
            
            # 构建结果
            result = {
                "success": True,
                "文件的文本信息": {
                    "文件全部的文本长度": len(cleaned_text),
                    "本次读取的长度": len(truncated_text),
                    "剩余未读取的长度": len(cleaned_text) - (start_index + len(truncated_text)),
                    "本次读取文本的开始位置": start_index,
                    "本次读取文本的结束位置": start_index + len(truncated_text),
                }
            }
            
            # 如果是全部文本，没有截断，key 使用text ，否则使用部分text
            if start_index == 0 and max_length >= len(truncated_text):
                result["本次读取的文本"] = truncated_text
            else:
                result["本次读取的部分文本"] = truncated_text

            if include_metadata and metadata:
                result["文件的metadata"] = metadata
                logger.debug(f"📋 包含元数据: {len(metadata)} 项")
            return result
            
        except Exception as e:
            error_time = time.time() - start_time
            logger.error(f"💥 文本提取异常 [{operation_id}] - 错误: {str(e)}, 耗时: {error_time:.2f}秒")
            logger.error(f"🔍 异常详情: {traceback.format_exc()}")
            
            return {
                "success": False,
                "error": str(e),
                "file_path": input_file_path,
                "execution_time": error_time,
                "operation_id": operation_id
            }

    @ToolBase.tool()
    def extract_text_from_url(
        self,
        url: str,
        start_index: int = 0,
        max_length: int = 5000,
        timeout: int = 30
    ) -> Dict[str, Any]:
        """从URL提取网页html文本内容

        Args:
            url (str): 目标URL地址
            start_index (int): 开始提取的字符位置，默认0
            max_length (int): 最大提取长度，默认5000字符
            timeout (int): 请求超时时间，默认30秒

        Returns:
            Dict[str, Any]: 包含提取文本和相关信息的字典
        """
        start_time = time.time()
        operation_id = hashlib.md5(f"url_{url}_{time.time()}".encode()).hexdigest()[:8]
        logger.info(f"🌐 extract_text_from_url开始执行 [{operation_id}] - URL: {url}")
        logger.info(f"🔧 参数: start_index={start_index}, max_length={max_length}, timeout={timeout}秒")
        
        try:
            # 验证URL格式
            logger.debug(f"🔍 验证URL格式")
            if not url.startswith(('http://', 'https://')):
                error_time = time.time() - start_time
                logger.error(f"❌ URL格式无效 [{operation_id}] - URL: {url}, 耗时: {error_time:.2f}秒")
                return {
                    "success": False,
                    "error": "URL必须以http://或https://开头",
                    "error_type": "invalid_url_format",
                    "url": url,
                    "execution_time": error_time,
                    "operation_id": operation_id
                }
            
            # 提取网页内容
            fetch_start_time = time.time()
            logger.info(f"🌐 开始获取网页内容")
            
            try:
                extracted_text = WebParser.extract_text_from_url(url, timeout)
                fetch_time = time.time() - fetch_start_time
                logger.info(f"✅ 网页内容获取成功 [{operation_id}] - 原始文本长度: {len(extracted_text)}, 获取耗时: {fetch_time:.2f}秒")
                
            except FileParserError as e:
                error_msg = str(e)
                fetch_time = time.time() - fetch_start_time
                error_time = time.time() - start_time
                
                # 根据错误类型提供更友好的错误信息
                error_type = "unknown_error"
                user_friendly_error = error_msg
                
                if "404" in error_msg or "不存在" in error_msg:
                    error_type = "url_not_found"
                    user_friendly_error = f"URL不存在或已失效: {url}"
                elif "403" in error_msg or "禁止" in error_msg:
                    error_type = "access_forbidden"
                    user_friendly_error = f"访问被拒绝，可能需要登录或权限: {url}"
                elif "500" in error_msg or "服务器" in error_msg:
                    error_type = "server_error"
                    user_friendly_error = f"目标服务器出现问题: {url}"
                elif "连接" in error_msg or "Connection" in error_msg:
                    error_type = "connection_error"
                    user_friendly_error = f"无法连接到目标网站，请检查网络连接: {url}"
                elif "连接超时" in error_msg or "ConnectTimeout" in error_msg:
                    error_type = "connect_timeout_error"
                    user_friendly_error = f"连接超时，网站响应过慢，请稍后重试: {url}"
                elif "读取超时" in error_msg or "ReadTimeout" in error_msg:
                    error_type = "read_timeout_error"
                    user_friendly_error = f"数据读取超时，网页内容较大或网速较慢，建议增加超时时间: {url}"
                elif "超时" in error_msg or "Timeout" in error_msg:
                    error_type = "timeout_error"
                    user_friendly_error = f"网页加载超时，请稍后重试或增加超时时间: {url}"
                
                logger.error(f"❌ 网页内容获取失败 [{operation_id}] - 错误类型: {error_type}, 耗时: {error_time:.2f}秒")
                
                return {
                    "success": False,
                    "error": user_friendly_error,
                    "error_type": error_type,
                    "error_details": error_msg,
                    "url": url,
                    "suggestions": self._get_error_suggestions(error_type, url),
                    "execution_time": error_time,
                    "operation_id": operation_id
                }
            
            # 清理和处理文本
            logger.debug(f"🧹 开始文本清理和处理")
            cleaned_text = TextProcessor.clean_text(extracted_text)
            truncated_text = TextProcessor.truncate_text(cleaned_text, start_index, max_length)
            text_stats = TextProcessor.get_text_stats(cleaned_text)
            
            total_time = time.time() - start_time
            
            logger.info(f"✅ URL文本提取完成 [{operation_id}] - 清理后长度: {len(cleaned_text)}, 截取长度: {len(truncated_text)}, 总耗时: {total_time:.2f}秒")
            logger.debug(f"📊 文本统计: {text_stats}")
            
            return {
                "success": True,
                "text": truncated_text,
                "url_info": {
                    "url": url,
                    "timeout": timeout,
                    "fetch_time": fetch_time
                },
                "text_info": {
                    "original_length": len(extracted_text),
                    "cleaned_length": len(cleaned_text),
                    "extracted_length": len(truncated_text),
                    "start_index": start_index,
                    "max_length": max_length,
                    **text_stats
                },
                "execution_time": total_time,
                "operation_id": operation_id
            }
            
        except Exception as e:
            error_time = time.time() - start_time
            logger.error(f"💥 URL文本提取异常 [{operation_id}] - 错误: {str(e)}, 耗时: {error_time:.2f}秒")
            logger.error(f"🔍 异常详情: {traceback.format_exc()}")
            
            return {
                "success": False,
                "error": f"处理URL时发生意外错误: {str(e)}",
                "error_type": "unexpected_error",
                "url": url,
                "execution_time": error_time,
                "operation_id": operation_id
            }
    
    def _get_error_suggestions(self, error_type: str, url: str) -> List[str]:
        """根据错误类型提供建议"""
        suggestions = []
        
        if error_type == "url_not_found":
            suggestions = [
                "检查URL是否正确拼写",
                "确认网页是否还存在",
                "尝试访问网站首页确认网站是否可用"
            ]
        elif error_type == "access_forbidden":
            suggestions = [
                "检查是否需要登录账户",
                "确认是否有访问权限",
                "尝试在浏览器中手动访问该URL"
            ]
        elif error_type == "connection_error":
            suggestions = [
                "检查网络连接是否正常",
                "确认防火墙是否阻止了访问",
                "尝试访问其他网站确认网络状态"
            ]
        elif error_type == "connect_timeout_error":
            suggestions = [
                "检查网络连接是否稳定",
                "稍后重试",
                "尝试使用VPN或代理服务"
            ]
        elif error_type == "read_timeout_error":
            suggestions = [
                "显著增加超时时间参数（如设置为60-90秒）",
                "检查网络速度是否正常",
                "分段获取内容，或稍后重试"
            ]
        elif error_type == "timeout_error":
            suggestions = [
                "增加超时时间参数",
                "稍后重试",
                "检查网络速度是否正常"
            ]
        elif error_type == "server_error":
            suggestions = [
                "稍后重试",
                "联系网站管理员",
                "尝试访问网站的其他页面"
            ]
        else:
            suggestions = [
                "检查URL格式是否正确",
                "确认网络连接正常",
                "稍后重试"
            ]
        
        return suggestions

    @ToolBase.tool()


    @ToolBase.tool()
    def get_supported_formats(self) -> Dict[str, Any]:
        """获取支持的文件格式列表

        Returns:
            Dict[str, Any]: 支持的文件格式和相关信息
        """
        return {
            "success": True,
            "supported_formats": FileValidator.SUPPORTED_FORMATS,
            "max_file_sizes": FileValidator.MAX_FILE_SIZE,
            "total_formats": len(FileValidator.SUPPORTED_FORMATS)
        }