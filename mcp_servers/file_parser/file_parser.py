import argparse
import hashlib
import logging
import os
import re
import subprocess
import tempfile
import time
import traceback
from pathlib import Path
from typing import Any, Dict, List

import aspose.slides as slides
import chardet
import html2text
import pdfplumber
import pypandoc
import requests
import uvicorn
from mcp.server.fastmcp import FastMCP
from openpyxl import load_workbook
from pptx import Presentation
from starlette.applications import Starlette
from starlette.routing import Mount

# 设置日志
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

mcp = FastMCP("Advanced File Parser")

parser = argparse.ArgumentParser(description='启动高级文件解析 MCP Server')
args = parser.parse_args()


class FileParserError(Exception):
    """文件解析异常"""
    pass


class FileValidator:
    """文件验证器"""

    # 支持的文件类型和对应的MIME类型
    SUPPORTED_FORMATS = {
        '.pdf': 'application/pdf',
        '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        '.doc': 'application/msword',
        '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        '.ppt': 'application/vnd.ms-powerpoint',
        '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        '.xls': 'application/vnd.ms-excel',
        '.txt': 'text/plain',
        '.csv': 'text/csv',
        '.json': 'application/json',
        '.xml': 'application/xml',
        '.html': 'text/html',
        '.htm': 'text/html',
        '.md': 'text/markdown',
        '.markdown': 'text/markdown',
        '.rtf': 'application/rtf',
        '.odt': 'application/vnd.oasis.opendocument.text',
        '.epub': 'application/epub+zip',
        '.latex': 'application/x-latex',
        '.tex': 'application/x-tex'
    }

    # 文件大小限制 (MB)
    MAX_FILE_SIZE = {
        '.pdf': 50,
        '.docx': 25,
        '.doc': 25,
        '.pptx': 100,
        '.ppt': 100,
        '.xlsx': 25,
        '.xls': 25,
        '.txt': 10,
        '.csv': 50,
        '.json': 10,
        '.xml': 10,
        '.html': 5,
        '.htm': 5,
        '.md': 5,
        '.markdown': 5,
        '.rtf': 10,
        '.odt': 25,
        '.epub': 50,
        '.latex': 10,
        '.tex': 10
    }

    @staticmethod
    def validate_file(file_path: str) -> Dict[str, Any]:
        """验证文件的有效性"""
        try:
            path = Path(file_path)

            # 检查文件是否存在
            if not path.exists():
                return {"valid": False, "error": f"文件不存在: {file_path}"}

            # 检查是否为文件（非目录）
            if not path.is_file():
                return {"valid": False, "error": f"路径不是有效文件: {file_path}"}

            # 获取文件扩展名
            file_extension = path.suffix.lower()

            # 检查文件格式是否支持
            if file_extension not in FileValidator.SUPPORTED_FORMATS:
                return {"valid": False, "error": f"不支持的文件格式: {file_extension}"}

            # 检查文件大小
            file_size_mb = path.stat().st_size / (1024 * 1024)
            max_size = FileValidator.MAX_FILE_SIZE.get(file_extension, 10)

            if file_size_mb > max_size:
                return {"valid": False, "error": f"文件过大: {file_size_mb:.1f}MB > {max_size}MB"}

            # 检查文件是否可读
            try:
                with open(file_path, 'rb') as f:
                    f.read(1024)  # 尝试读取前1KB
            except PermissionError:
                return {"valid": False, "error": "文件无读取权限"}
            except Exception as e:
                return {"valid": False, "error": f"文件读取失败: {str(e)}"}

            return {
                "valid": True,
                "file_size_mb": file_size_mb,
                "file_extension": file_extension,
                "mime_type": FileValidator.SUPPORTED_FORMATS[file_extension]
            }

        except Exception as e:
            return {"valid": False, "error": f"文件验证失败: {str(e)}"}


class TextProcessor:
    """文本处理器"""

    @staticmethod
    def detect_encoding(file_path: str) -> str:
        """检测文件编码"""
        try:
            with open(file_path, 'rb') as f:
                raw_data = f.read(10000)  # 读取前10KB进行检测
                result = chardet.detect(raw_data)
                return result.get('encoding', 'utf-8')
        except Exception:
            return 'utf-8'

    @staticmethod
    def clean_text(text: str) -> str:
        """清理文本内容"""
        if not text:
            return ""

        # 移除多余的空白字符
        text = re.sub(r'\n\s*\n', '\n\n', text)  # 多个连续换行符合并为双换行
        text = re.sub(r'[ \t]+', ' ', text)      # 多个连续空格合并为单个空格
        text = text.strip()

        return text

    @staticmethod
    def truncate_text(text: str, start_index: int = 0, max_length: int = 5000) -> str:
        """安全地截取文本"""
        if not text:
            return ""

        start_index = max(0, start_index)
        end_index = min(len(text), start_index + max_length)

        return text[start_index:end_index]

    @staticmethod
    def get_text_stats(text: str) -> Dict[str, int]:
        """获取文本统计信息"""
        if not text:
            return {"characters": 0, "words": 0, "lines": 0, "paragraphs": 0}

        return {
            "characters": len(text),
            "words": len(text.split()),
            "lines": text.count('\n') + 1,
            "paragraphs": len([p for p in text.split('\n\n') if p.strip()])
        }


class PDFParser:
    """PDF解析器"""

    @staticmethod
    def extract_text(pdf_path: str) -> str:
        """从PDF提取文本"""
        try:
            with pdfplumber.open(pdf_path) as pdf:
                text_parts = []
                for page_num, page in enumerate(pdf.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text:
                            text_parts.append(f"=== 第 {page_num + 1} 页 ===\n{page_text}")
                    except Exception as e:
                        logger.warning(f"PDF第{page_num + 1}页解析失败: {e}")
                        text_parts.append(f"=== 第 {page_num + 1} 页 ===\n[页面解析失败: {str(e)}]")

                return "\n\n".join(text_parts)

        except Exception as e:
            raise FileParserError(f"PDF解析失败: {str(e)}")

    @staticmethod
    def get_pdf_info(pdf_path: str) -> Dict[str, Any]:
        """获取PDF信息"""
        try:
            with pdfplumber.open(pdf_path) as pdf:
                return {
                    "pages": len(pdf.pages),
                    "metadata": pdf.metadata or {}
                }
        except Exception:
            return {"pages": 0, "metadata": {}}


class OfficeParser:
    """Office文档解析器"""

    @staticmethod
    def extract_text_from_docx(file_path: str) -> str:
        """从DOCX提取文本"""
        try:
            return pypandoc.convert_file(file_path, 'markdown', extra_args=['--extract-media=.'])
        except Exception as e:
            raise FileParserError(f"DOCX解析失败: {str(e)}")

    @staticmethod
    def extract_text_from_doc(file_path: str) -> str:
        """从DOC提取文本（需要antiword）"""
        try:
            result = subprocess.run(
                ['antiword', file_path],
                capture_output=True,
                text=True,
                timeout=30
            )
            if result.returncode == 0:
                return result.stdout
            else:
                raise FileParserError(f"DOC转换失败: {result.stderr}")
        except subprocess.TimeoutExpired:
            raise FileParserError("DOC解析超时")
        except FileNotFoundError:
            raise FileParserError("未找到antiword工具，请安装: sudo apt-get install antiword")
        except Exception as e:
            raise FileParserError(f"DOC解析失败: {str(e)}")

    @staticmethod
    def extract_text_from_pptx(file_path: str) -> str:
        """从PPTX提取文本"""
        try:
            prs = Presentation(file_path)
            slides_text = []

            for slide_num, slide in enumerate(prs.slides):
                slide_content = [f"=== 幻灯片 {slide_num + 1} ==="]

                # 按位置排序形状
            shapes = sorted(slide.shapes, key=lambda x: (x.top, x.left))

            for shape in shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())

                slides_text.append('\n'.join(slide_content))

            return '\n\n'.join(slides_text)

        except Exception as e:
            raise FileParserError(f"PPTX解析失败: {str(e)}")

    @staticmethod
    def extract_text_from_ppt(file_path: str) -> str:
        """从PPT提取文本（需要aspose.slides）"""
        try:
            with tempfile.TemporaryDirectory() as temp_dir:
                temp_pptx = os.path.join(temp_dir, "temp.pptx")

                with slides.Presentation(file_path) as presentation:
                    presentation.save(temp_pptx, slides.export.SaveFormat.PPTX)

                return OfficeParser.extract_text_from_pptx(temp_pptx)

        except Exception as e:
            raise FileParserError(f"PPT解析失败: {str(e)}")


class ExcelParser:
    """Excel解析器"""

    @staticmethod
    def extract_text_from_xlsx(file_path: str) -> str:
        """从Excel提取文本并转换为Markdown"""
        try:
            excel_data = ExcelParser._read_excel_to_dict(file_path)
            markdown_tables = []

            for sheet_name, sheet_data in excel_data.items():
                # 限制行数
                if len(sheet_data) > 100:
                    sheet_data = sheet_data[:100]

                sheet_md = ExcelParser._sheet_data_to_markdown(sheet_data, sheet_name)
                markdown_tables.append(sheet_md)

            return '\n\n'.join(markdown_tables)

        except Exception as e:
            raise FileParserError(f"Excel解析失败: {str(e)}")

    @staticmethod
    def _read_excel_to_dict(file_path: str) -> Dict[str, List[List[str]]]:
        """读取Excel文件到字典，正确处理合并单元格"""
        # 需要关闭read_only模式才能访问合并单元格信息
        workbook = load_workbook(file_path, data_only=True, read_only=False)
        excel_data = {}

        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]

            # 创建合并单元格值映射
            merged_cell_values = {}
            for merged_range in sheet.merged_cells.ranges:
                # 获取合并单元格左上角的值
                top_left_cell = sheet.cell(merged_range.min_row, merged_range.min_col)
                value = top_left_cell.value

                # 为合并范围内的所有单元格设置相同的值
                for row in range(merged_range.min_row, merged_range.max_row + 1):
                    for col in range(merged_range.min_col, merged_range.max_col + 1):
                        merged_cell_values[(row, col)] = value

            # 读取数据，考虑合并单元格
            sheet_data = []
            max_row = sheet.max_row
            max_col = sheet.max_column

            if max_row and max_col:
                for row_idx in range(1, max_row + 1):
                    row_data = []
                    for col_idx in range(1, max_col + 1):
                        # 检查是否是合并单元格
                        if (row_idx, col_idx) in merged_cell_values:
                            cell_value = merged_cell_values[(row_idx, col_idx)]
                        else:
                            cell = sheet.cell(row_idx, col_idx)
                            cell_value = cell.value

                        cell_str = str(cell_value).replace('\n', '\\n') if cell_value is not None else ''
                        row_data.append(cell_str)

                    sheet_data.append(row_data)

                if not sheet_data:
                    continue

                # 清理空行和空列
                sheet_data = ExcelParser._clean_empty_rows_cols(sheet_data)

                if sheet_data:
                    excel_data[sheet_name] = sheet_data

        workbook.close()
        return excel_data

    @staticmethod
    def _clean_empty_rows_cols(data: List[List[str]]) -> List[List[str]]:
        """清理空行和空列"""
        if not data:
            return data

        # 移除空行
        data = [row for row in data if any(cell.strip() for cell in row)]

        if not data:
            return data

        # 移除空列
        cols_to_keep = []
        for col_idx in range(len(data[0])):
            if any(row[col_idx].strip() for row in data):
                cols_to_keep.append(col_idx)

        if cols_to_keep:
            data = [[row[i] for i in cols_to_keep] for row in data]

        return data

    @staticmethod
    def _sheet_data_to_markdown(sheet_data: List[List[str]], sheet_name: str) -> str:
        """转换工作表数据为Markdown"""
        if not sheet_data:
            return f'## {sheet_name}\n\n(空工作表)'

        markdown_lines = [f'## {sheet_name}', '']

        # 如果有数据，第一行作为表头
        if sheet_data:
            # 表头
            header = '| ' + ' | '.join(cell if cell else ' ' for cell in sheet_data[0]) + ' |'
            markdown_lines.append(header)

            # 分隔线
            separator = '| ' + ' | '.join('---' for _ in sheet_data[0]) + ' |'
            markdown_lines.append(separator)

            # 数据行
            for row in sheet_data[1:]:
                row_md = '| ' + ' | '.join(cell if cell else ' ' for cell in row) + ' |'
                markdown_lines.append(row_md)

        return '\n'.join(markdown_lines)


class WebParser:
    """网页解析器"""

    @staticmethod
    def extract_text_from_html(file_path: str) -> str:
        """从HTML文件提取文本"""
        try:
            encoding = TextProcessor.detect_encoding(file_path)
            with open(file_path, 'r', encoding=encoding) as file:
                html_content = file.read()

            return WebParser._html_to_text(html_content)

        except Exception as e:
            raise FileParserError(f"HTML解析失败: {str(e)}")

    @staticmethod
    def extract_text_from_url(url: str, timeout: int = 30) -> str:
        """从URL提取文本"""
        try:
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
            }
            response = requests.get(url, headers=headers, timeout=timeout)
            response.raise_for_status()

            return WebParser._html_to_text(response.text)

        except requests.RequestException as e:
            raise FileParserError(f"URL访问失败: {str(e)}")
        except Exception as e:
            raise FileParserError(f"URL解析失败: {str(e)}")

    @staticmethod
    def _html_to_text(html_content: str) -> str:
        """HTML转文本"""
        h = html2text.HTML2Text()
        h.ignore_links = False
        h.ignore_images = False
        h.body_width = 0  # 不限制行宽
        return h.handle(html_content)


class PlainTextParser:
    """纯文本解析器"""

    @staticmethod
    def extract_text_from_plain_file(file_path: str) -> str:
        """从纯文本文件提取内容"""
        try:
            encoding = TextProcessor.detect_encoding(file_path)
            with open(file_path, 'r', encoding=encoding) as file:
                return file.read()
        except Exception as e:
            raise FileParserError(f"文本文件解析失败: {str(e)}")

    @staticmethod
    def extract_text_with_pandoc(file_path: str, input_format: str = None) -> str:
        """使用Pandoc提取文本"""
        try:
            if input_format:
                return pypandoc.convert_file(file_path, 'markdown', format=input_format)
            else:
                return pypandoc.convert_file(file_path, 'markdown')
        except Exception as e:
            raise FileParserError(f"Pandoc解析失败: {str(e)}")

# ==================== MCP 工具函数 ====================


@mcp.tool()
async def extract_text_from_file(
    input_file_path: str,
    start_index: int = 0,
    max_length: int = 5000,
    include_metadata: bool = False
) -> Dict[str, Any]:
    """
    从本地的各种格式的文件中提取方便阅读的markdown文本内容，如果是网络上的文件，请先进行下载到本地磁盘。

    支持的格式：PDF, DOCX, DOC, PPTX, PPT, XLSX, XLS, TXT, CSV, JSON, XML, HTML, MD等

    Args:
        input_file_path: 输入文件路径，本地的绝对路径
        start_index: 开始提取的字符位置（默认0）
        max_length: 最大提取长度（默认5000字符）
        include_metadata: 是否包含文件元数据（默认False）

    Returns:
        包含提取文本和相关信息的字典
    """
    start_time = time.time()
    operation_id = hashlib.md5(f"extract_{input_file_path}_{time.time()}".encode()).hexdigest()[:8]
    logger.info(f"📄 extract_text_from_file开始执行 [{operation_id}] - 文件: {input_file_path}")
    logger.info(f"🔧 参数: start_index={start_index}, max_length={max_length}, include_metadata={include_metadata}")

    try:
        # 验证文件
        logger.debug("🔍 开始文件验证")
        validation_result = FileValidator.validate_file(input_file_path)

        if not validation_result["valid"]:
            error_time = time.time() - start_time
            logger.error(f"❌ 文件验证失败 [{operation_id}] - 错误: {validation_result['error']}, 耗时: {error_time:.2f}秒")
            return {
                "success": False,
                "error": validation_result["error"],
                "file_path": input_file_path,
                "execution_time": error_time,
                "operation_id": operation_id
            }

        file_extension = validation_result["file_extension"]
        file_size_mb = validation_result["file_size_mb"]
        logger.info(f"✅ 文件验证通过 [{operation_id}] - 格式: {file_extension}, 大小: {file_size_mb:.2f}MB")

        # 根据文件类型选择解析器
        parse_start_time = time.time()
        logger.info(f"🔧 开始文件解析 - 格式: {file_extension}")

        extracted_text = ""
        metadata = {}

        try:
            if file_extension == '.pdf':
                logger.debug("📕 使用PDF解析器")
                extracted_text = PDFParser.extract_text(input_file_path)
                if include_metadata:
                    metadata = PDFParser.get_pdf_info(input_file_path)

            elif file_extension in ['.docx', '.doc']:
                logger.debug("📝 使用Word解析器")
                if file_extension == '.docx':
                    extracted_text = OfficeParser.extract_text_from_docx(input_file_path)
                else:
                    extracted_text = OfficeParser.extract_text_from_doc(input_file_path)

            elif file_extension in ['.pptx', '.ppt']:
                logger.debug("📊 使用PowerPoint解析器")
                if file_extension == '.pptx':
                    extracted_text = OfficeParser.extract_text_from_pptx(input_file_path)
                else:
                    extracted_text = OfficeParser.extract_text_from_ppt(input_file_path)

            elif file_extension in ['.xlsx', '.xls']:
                logger.debug("📈 使用Excel解析器")
                extracted_text = ExcelParser.extract_text_from_xlsx(input_file_path)

            elif file_extension in ['.html', '.htm']:
                logger.debug("🌐 使用HTML解析器")
                extracted_text = WebParser.extract_text_from_html(input_file_path)

            elif file_extension in ['.txt', '.csv', '.json', '.xml', '.md', '.markdown']:
                logger.debug("📄 使用纯文本解析器")
                extracted_text = PlainTextParser.extract_text_from_plain_file(input_file_path)

            else:
                # 尝试使用Pandoc解析
                logger.debug("🔧 尝试使用Pandoc解析器")
                extracted_text = PlainTextParser.extract_text_with_pandoc(input_file_path)

            parse_time = time.time() - parse_start_time
            logger.info(f"✅ 文件解析成功 [{operation_id}] - 原始文本长度: {len(extracted_text)}, 解析耗时: {parse_time:.2f}秒")

        except Exception as parse_error:
            parse_time = time.time() - parse_start_time
            logger.warning(f"⚠️ 主解析器失败 [{operation_id}] - 错误: {str(parse_error)}, 耗时: {parse_time:.2f}秒")
            logger.debug("🔧 尝试Pandoc备用解析器")

            try:
                extracted_text = PlainTextParser.extract_text_with_pandoc(input_file_path)
                logger.info(f"✅ Pandoc备用解析成功 [{operation_id}] - 文本长度: {len(extracted_text)}")
            except Exception as fallback_error:
                error_time = time.time() - start_time
                logger.error(f"❌ 所有解析器都失败 [{operation_id}] - 主错误: {str(parse_error)}, 备用错误: {str(fallback_error)}, 耗时: {error_time:.2f}秒")
                return {
                    "success": False,
                    "error": f"文件解析失败 - 主错误: {str(parse_error)}, 备用错误: {str(fallback_error)}",
                    "file_path": input_file_path,
                    "execution_time": error_time,
                    "operation_id": operation_id
                }

        # 清理和处理文本
        logger.debug("🧹 开始文本清理和处理")
        cleaned_text = TextProcessor.clean_text(extracted_text)
        truncated_text = TextProcessor.truncate_text(cleaned_text, start_index, max_length)
        text_stats = TextProcessor.get_text_stats(cleaned_text)

        total_time = time.time() - start_time

        logger.info(f"✅ 文本提取完成 [{operation_id}] - 清理后长度: {len(cleaned_text)}, 截取长度: {len(truncated_text)}, 总耗时: {total_time:.2f}秒")
        logger.debug(f"📊 文本统计: {text_stats}")

        # 构建结果
        result = {
            "success": True,
            "text": truncated_text,
            "file_info": {
                "file_path": input_file_path,
                "file_extension": file_extension,
                "file_size_mb": round(file_size_mb, 2),
                "mime_type": validation_result["mime_type"]
            },
            "text_info": {
                "original_length": len(extracted_text),
                "cleaned_length": len(cleaned_text),
                "extracted_length": len(truncated_text),
                "start_index": start_index,
                "max_length": max_length,
                **text_stats
            },
            "execution_time": total_time,
            "operation_id": operation_id
        }

        if include_metadata and metadata:
            result["metadata"] = metadata
            logger.debug(f"📋 包含元数据: {len(metadata)} 项")

        return result

    except Exception as e:
        error_time = time.time() - start_time
        logger.error(f"💥 文本提取异常 [{operation_id}] - 错误: {str(e)}, 耗时: {error_time:.2f}秒")
        logger.error(f"🔍 异常详情: {traceback.format_exc()}")

        return {
            "success": False,
            "error": str(e),
            "file_path": input_file_path,
            "execution_time": error_time,
            "operation_id": operation_id
        }


@mcp.tool()
async def extract_text_from_url(
    url: str,
    start_index: int = 0,
    max_length: int = 5000,
    timeout: int = 30
) -> Dict[str, Any]:
    """
    从URL提取网页html文本内容

    Args:
        url: 目标URL地址
        start_index: 开始提取的字符位置（默认0）
        max_length: 最大提取长度（默认5000字符）
        timeout: 请求超时时间（默认30秒）

    Returns:
        包含提取文本和相关信息的字典
    """
    start_time = time.time()
    operation_id = hashlib.md5(f"url_{url}_{time.time()}".encode()).hexdigest()[:8]
    logger.info(f"🌐 extract_text_from_url开始执行 [{operation_id}] - URL: {url}")
    logger.info(f"🔧 参数: start_index={start_index}, max_length={max_length}, timeout={timeout}秒")

    try:
        # 验证URL格式
        logger.debug("🔍 验证URL格式")
        if not url.startswith(('http://', 'https://')):
            error_time = time.time() - start_time
            logger.error(f"❌ URL格式无效 [{operation_id}] - URL: {url}, 耗时: {error_time:.2f}秒")
            return {
                "success": False,
                "error": "URL必须以http://或https://开头",
                "url": url,
                "execution_time": error_time,
                "operation_id": operation_id
            }

        # 提取网页内容
        fetch_start_time = time.time()
        logger.info("🌐 开始获取网页内容")

        extracted_text = WebParser.extract_text_from_url(url, timeout)

        fetch_time = time.time() - fetch_start_time
        logger.info(f"✅ 网页内容获取成功 [{operation_id}] - 原始文本长度: {len(extracted_text)}, 获取耗时: {fetch_time:.2f}秒")

        # 清理和处理文本
        logger.debug("🧹 开始文本清理和处理")
        cleaned_text = TextProcessor.clean_text(extracted_text)
        truncated_text = TextProcessor.truncate_text(cleaned_text, start_index, max_length)
        text_stats = TextProcessor.get_text_stats(cleaned_text)

        total_time = time.time() - start_time

        logger.info(f"✅ URL文本提取完成 [{operation_id}] - 清理后长度: {len(cleaned_text)}, 截取长度: {len(truncated_text)}, 总耗时: {total_time:.2f}秒")
        logger.debug(f"📊 文本统计: {text_stats}")

        return {
            "success": True,
            "text": truncated_text,
            "url_info": {
                "url": url,
                "timeout": timeout,
                "fetch_time": fetch_time
            },
            "text_info": {
                "original_length": len(extracted_text),
                "cleaned_length": len(cleaned_text),
                "extracted_length": len(truncated_text),
                "start_index": start_index,
                "max_length": max_length,
                **text_stats
            },
            "execution_time": total_time,
            "operation_id": operation_id
        }

    except Exception as e:
        error_time = time.time() - start_time
        logger.error(f"💥 URL文本提取异常 [{operation_id}] - 错误: {str(e)}, 耗时: {error_time:.2f}秒")
        logger.error(f"🔍 异常详情: {traceback.format_exc()}")

        return {
            "success": False,
            "error": str(e),
            "url": url,
            "execution_time": error_time,
            "operation_id": operation_id
        }


@mcp.tool()
async def batch_extract_text(
    file_paths: List[str],
    max_length: int = 3000,
    include_metadata: bool = False
) -> Dict[str, Any]:
    """批量提取多个文件的文本

    Args:
        file_paths (List[str]): 文件路径列表
        max_length (int): 每个文件的最大提取长度，默认3000
        include_metadata (bool): 是否包含元数据，默认False

    Returns:
        Dict: 批量处理结果
    """
    start_time = time.time()
    results = []
    successful = 0
    failed = 0

    for file_path in file_paths:
        try:
            result = await extract_text_from_file(
                input_file_path=file_path,
                start_index=0,
                max_length=max_length,
                include_metadata=include_metadata
            )

            if result["status"] == "success":
                successful += 1
            else:
                failed += 1

            results.append({
                "file_path": file_path,
                "result": result
            })

        except Exception as e:
            failed += 1
            results.append({
                "file_path": file_path,
                "result": {
                    "status": "error",
                    "message": f"处理失败: {str(e)}",
                    "text": "",
                    "length": 0
                }
            })

    processing_time = time.time() - start_time

    return {
        "status": "success",
        "message": f"批量处理完成，成功{successful}个，失败{failed}个",
        "summary": {
            "total_files": len(file_paths),
            "successful": successful,
            "failed": failed,
            "processing_time": round(processing_time, 2)
        },
        "results": results
    }

if __name__ == "__main__":
    app = Starlette(
        routes=[
            Mount('/', app=mcp.sse_app()),
        ]
    )

    uvicorn.run(app, host="0.0.0.0", port=34001)
