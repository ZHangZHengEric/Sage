#!/usr/bin/env python3

import argparse
import os
import shutil
import traceback
import sys
import xml.etree.ElementTree as ET
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE
from pptx.util import Length

def _format_inch(val) -> str:
    # Convert EMUs to Inches with reasonable precision
    if val is None:
        return "0"
    inches = val / 914400.0
    s = f"{inches:.3f}".rstrip("0").rstrip(".")
    if s == "":
        return "0"
    return s

def _format_color(color_format) -> str:
    if not color_format:
        return ""
    try:
        if color_format.type == 1: # RGB
             return str(color_format.rgb)
        # Theme colors are complex, skip for now or use placeholder
        return "" 
    except:
        return ""

class PptToProjectConverter:
    def __init__(self, ppt_path: Path, output_dir: Path):
        self.ppt_path = ppt_path
        self.output_dir = output_dir
        self.slides_dir = output_dir / "slides"
        self.assets_dir = output_dir / "assets"
        self.images_dir = self.assets_dir / "images"
        self.prs = None

    def convert(self):
        print(f"Loading PPT: {self.ppt_path}...")
        self.prs = Presentation(str(self.ppt_path))
        
        # Prepare dirs
        self.slides_dir.mkdir(parents=True, exist_ok=True)
        self.images_dir.mkdir(parents=True, exist_ok=True)
        
        # Process slides
        for i, slide in enumerate(self.prs.slides):
            self._process_slide(i, slide)

        # Copy template (the original PPT) to output root as 'template.pptx'
        template_path = self.output_dir / "template.pptx"
        try:
            shutil.copy2(self.ppt_path, template_path)
            print(f"Saved template to: {template_path}")
        except Exception as e:
            print(f"Warning: Could not save template: {e}")

        print(f"✓ Project created at: {self.output_dir}")
            
    def _process_slide(self, index: int, slide):
        slide_num = index + 1
        filename = f"{slide_num:02d}-slide.xml"
        output_path = self.slides_dir / filename
        
        root = ET.Element("ppt-slide")
        # Set basic attributes
        root.set("width", _format_inch(self.prs.slide_width))
        root.set("height", _format_inch(self.prs.slide_height))
        
        # Try to detect background color if solid
        try:
            bg = slide.background
            if bg and bg.fill.type == 1: # solid
                 root.set("bg-color", _format_color(bg.fill.fore_color))
        except:
            pass
            
        # Process shapes
        for shape in slide.shapes:
            self._process_shape(shape, root)
            
        # Write XML
        # Pretty print using minidom or just indent manually
        # ET.indent is available in Python 3.9+
        if hasattr(ET, "indent"):
            ET.indent(root, space="  ", level=0)
            
        tree = ET.ElementTree(root)
        with open(output_path, "wb") as f:
            tree.write(f, encoding="utf-8", xml_declaration=False)
        print(f"  Generated: {filename}")

    def _process_shape(self, shape, parent):
        try:
            x = _format_inch(shape.left)
            y = _format_inch(shape.top)
            w = _format_inch(shape.width)
            h = _format_inch(shape.height)
            
            tag = None
            attrs = {"x": x, "y": y, "w": w, "h": h}
            text = ""
            
            shape_type = shape.shape_type
            
            # TEXT BOX
            if shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                tag = "ppt-text"
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    # Extract font properties
                    if shape.text_frame.paragraphs:
                        p = shape.text_frame.paragraphs[0]
                        if p.runs:
                            run = p.runs[0]
                            if run.font.size:
                                attrs["font-size"] = str(int(run.font.size.pt))
                            if run.font.bold:
                                attrs["bold"] = "true"
                            if run.font.italic:
                                attrs["italic"] = "true"
                            if run.font.color and run.font.color.type == 1:
                                attrs["color"] = _format_color(run.font.color)
                            if run.font.name:
                                attrs["font-name"] = run.font.name
                        # Extract alignment
                        if p.alignment:
                             # Map PP_ALIGN to string
                             # This is tricky as PP_ALIGN is an enum. 
                             # Simpler to just skip or do rough mapping if needed.
                             pass

            # AUTO SHAPE (RECT, ETC)
            elif shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                tag = "ppt-rect"
                attrs["shape"] = "rect" # Default
                
                # Try to map shape type to name
                # MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE -> "rounded"
                # This requires importing MSO_AUTO_SHAPE_TYPE and checking
                if hasattr(shape, "auto_shape_type"):
                     # Just a heuristic
                     if shape.auto_shape_type == MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE:
                         attrs["shape"] = "rounded"
                         # Try to extract radius? Hard in python-pptx
                
                # Check if it has text, if so, maybe treat as text box with fill?
                if shape.has_text_frame and shape.text_frame.text.strip():
                     tag = "ppt-text"
                     text = shape.text_frame.text
                     # Extract fill color if any
                     try:
                         if shape.fill.type == 1:
                             attrs["fill"] = _format_color(shape.fill.fore_color)
                     except:
                         pass
                     # Extract font props for text in shape
                     if shape.text_frame.paragraphs:
                        p = shape.text_frame.paragraphs[0]
                        if p.runs:
                            run = p.runs[0]
                            if run.font.size:
                                attrs["font-size"] = str(int(run.font.size.pt))
                            if run.font.color and run.font.color.type == 1:
                                attrs["color"] = _format_color(run.font.color)
                else:
                     # Pure shape
                     try:
                         if shape.fill.type == 1:
                             attrs["fill"] = _format_color(shape.fill.fore_color)
                     except:
                         pass
                     try:
                         if shape.line.fill.type == 1:
                             attrs["line"] = _format_color(shape.line.color)
                             attrs["line-width"] = str(int(shape.line.width.pt)) if shape.line.width else "1"
                     except:
                         pass

            # PICTURE
            elif shape_type == MSO_SHAPE_TYPE.PICTURE:
                tag = "ppt-image"
                # Extract image
                try:
                    image = shape.image
                    image_bytes = image.blob
                    ext = image.ext
                    # Use sha1 or just sequential naming
                    image_filename = f"image_{shape.shape_id}.{ext}"
                    image_path = self.images_dir / image_filename
                    with open(image_path, "wb") as f:
                        f.write(image_bytes)
                    attrs["src"] = f"assets/images/{image_filename}"
                except Exception as e:
                    print(f"Warning: Failed to extract image from shape {shape.name}: {e}")
                    return

            # TABLE
            elif shape_type == MSO_SHAPE_TYPE.TABLE:
                tag = "ppt-table"
                self._process_table(shape.table, parent, attrs)
                return # Table adds its own element in _process_table helper

            if tag:
                el = ET.SubElement(parent, tag, attrs)
                if text:
                    el.text = text

        except Exception as e:
            print(f"    Skipping shape {shape.name}: {e}")

    def _process_table(self, table, parent, attrs):
        table_el = ET.SubElement(parent, "ppt-table", attrs)
        for row in table.rows:
            row_el = ET.SubElement(table_el, "ppt-row")
            for cell in row.cells:
                cell_el = ET.SubElement(row_el, "ppt-cell")
                if cell.text_frame.text:
                    cell_el.text = cell.text_frame.text


def main():
    parser = argparse.ArgumentParser(description="Convert existing PPTX to XML project structure")
    parser.add_argument("ppt_path", help="Path to existing .pptx file")
    parser.add_argument("output_dir", help="Directory to output the project")
    
    args = parser.parse_args()
    
    ppt_path = Path(args.ppt_path)
    if not ppt_path.exists():
        print(f"Error: File not found: {ppt_path}")
        sys.exit(1)
        
    output_dir = Path(args.output_dir)
    
    try:
        converter = PptToProjectConverter(ppt_path, output_dir)
        converter.convert()
    except Exception:
        traceback.print_exc()
        sys.exit(1)

if __name__ == "__main__":
    main()
