#!/usr/bin/env python3
from __future__ import annotations

import argparse
import math
import csv
import io
import json
import os
import re
import sys
import traceback
import xml.etree.ElementTree as ET
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR as MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


BUILTIN_THEMES = {
    "tech-dark": {
        "page_bg": "#0B1220",
        "text": {
            "title": {"font-size": 40, "bold": True, "color": "#E5E7EB", "align": "left", "valign": "top"},
            "subtitle": {"font-size": 18, "color": "#94A3B8", "align": "left", "valign": "top"},
            "body": {"font-size": 16, "color": "#CBD5E1", "align": "left", "valign": "top"},
            "muted": {"font-size": 14, "color": "#94A3B8", "align": "left", "valign": "top"},
        },
        "rect": {
            "card": {"fill": "#111827", "line": "#1F2937", "line-width": 1, "shape": "rounded", "radius": 0.18},
            "outline": {"fill": None, "line": "#1F2937", "line-width": 1, "shape": "rect"},
        },
        "table": {
            "default": {
                "font-size": 12,
                "header": True,
                "header-fill": "#111827",
                "header-color": "#E5E7EB",
                "cell-fill": "#0F172A",
                "cell-color": "#CBD5E1",
            }
        },
    },
    "corporate": {
        "page_bg": "#FFFFFF",
        "text": {
            "title": {"font-size": 40, "bold": True, "color": "#111827", "align": "left", "valign": "top"},
            "subtitle": {"font-size": 18, "color": "#475569", "align": "left", "valign": "top"},
            "body": {"font-size": 16, "color": "#1F2937", "align": "left", "valign": "top"},
            "muted": {"font-size": 14, "color": "#64748B", "align": "left", "valign": "top"},
        },
        "rect": {
            "card": {"fill": "#F8FAFC", "line": "#E2E8F0", "line-width": 1, "shape": "rounded", "radius": 0.18},
            "outline": {"fill": None, "line": "#CBD5E1", "line-width": 1, "shape": "rect"},
        },
        "table": {
            "default": {
                "font-size": 12,
                "header": True,
                "header-fill": "#F1F5F9",
                "header-color": "#0F172A",
                "cell-fill": "#FFFFFF",
                "cell-color": "#0F172A",
            }
        },
    },
}


def _parse_bool(value: str, default: bool = False) -> bool:
    if value is None:
        return default
    v = str(value).strip().lower()
    if v in {"1", "true", "yes", "y", "on"}:
        return True
    if v in {"0", "false", "no", "n", "off"}:
        return False
    return default


def _parse_float(value: str, default: float | None = None) -> float | None:
    if value is None:
        return default
    try:
        return float(str(value).strip())
    except Exception:
        traceback.print_exc()
        return default


def _parse_color(value: str | None, default: str = "111827") -> RGBColor:
    if not value:
        value = default
    v = str(value).strip()
    if v.startswith("#"):
        v = v[1:]
    if re.fullmatch(r"[0-9a-fA-F]{6}", v) is None:
        v = default
    r = int(v[0:2], 16)
    g = int(v[2:4], 16)
    b = int(v[4:6], 16)
    return RGBColor(r, g, b)


def _parse_inch_list(value: str | None) -> list[float] | None:
    if not value:
        return None
    parts = [p.strip() for p in str(value).split(",") if p.strip()]
    out: list[float] = []
    for p in parts:
        try:
            out.append(float(p))
        except Exception:
            traceback.print_exc()
            raise
            return None
    return out if out else None


def _apply_table_cell_border(cell, color: str | None, width: float | None) -> None:
    if width is None:
        return
    try:
        width_value = float(width)
        width_pt = width_value * 72 if width_value > 0 else 0
        color_rgb = _parse_color(color, default="CBD5E1")
        color_hex = f"{color_rgb[0]:02X}{color_rgb[1]:02X}{color_rgb[2]:02X}"
        tc = cell._tc
        tc_pr = tc.get_or_add_tcPr()
        for edge in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
            ln = tc_pr.find(qn(edge))
            if ln is None:
                ln = OxmlElement(edge)
                tc_pr.append(ln)
            for child in list(ln):
                ln.remove(child)
            ln.set("w", str(int(width_pt * 12700)))
            solid_fill = OxmlElement("a:solidFill")
            srgb_clr = OxmlElement("a:srgbClr")
            srgb_clr.set("val", color_hex)
            solid_fill.append(srgb_clr)
            ln.append(solid_fill)
            prst_dash = OxmlElement("a:prstDash")
            prst_dash.set("val", "solid")
            ln.append(prst_dash)
    except Exception:
        traceback.print_exc()


def _merge_style(base: dict, override: dict) -> dict:
    out = dict(base or {})
    for k, v in (override or {}).items():
        out[k] = v
    return out


def _extract_slide_fragment(text: str) -> str | None:
    m = re.search(r"<ppt-slide\b[\s\S]*?</ppt-slide>", text, flags=re.IGNORECASE)
    if not m:
        return None
    return m.group(0)


def _parse_css_vars(css_text: str) -> dict[str, str]:
    out: dict[str, str] = {}
    for m in re.finditer(r"--([A-Za-z0-9_-]+)\s*:\s*([^;]+)\s*;", css_text):
        key = m.group(1).strip()
        val = m.group(2).strip()
        if key:
            out[key] = val
    return out


def _css_vars_to_theme(vars_map: dict[str, str]) -> dict:
    def v(name: str) -> str | None:
        raw = vars_map.get(name)
        return raw.strip() if isinstance(raw, str) and raw.strip() else None

    title_size = _parse_float(v("title-size"), default=40) or 40
    subtitle_size = _parse_float(v("subtitle-size"), default=18) or 18
    body_size = _parse_float(v("body-size"), default=16) or 16
    muted_size = _parse_float(v("muted-size"), default=14) or 14
    card_radius = _parse_float(v("card-radius"), default=0.18) or 0.18
    table_size = _parse_float(v("table-font-size"), default=12) or 12

    return {
        "page_bg": v("page-bg"),
        "text": {
            "title": {"font-size": title_size, "bold": True, "color": v("title-color"), "align": "left", "valign": "top"},
            "subtitle": {"font-size": subtitle_size, "color": v("subtitle-color"), "align": "left", "valign": "top"},
            "body": {"font-size": body_size, "color": v("body-color"), "align": "left", "valign": "top"},
            "muted": {"font-size": muted_size, "color": v("muted-color"), "align": "left", "valign": "top"},
        },
        "rect": {
            "card": {
                "fill": v("card-fill"),
                "line": v("card-line"),
                "line-width": _parse_float(v("card-line-width"), default=1) or 1,
                "shape": "rounded",
                "radius": card_radius,
            },
            "outline": {"fill": None, "line": v("outline-line"), "line-width": _parse_float(v("outline-line-width"), default=1) or 1, "shape": "rect"},
        },
        "table": {
            "default": {
                "font-size": table_size,
                "header": True,
                "header-fill": v("table-header-fill"),
                "header-color": v("table-header-color"),
                "cell-fill": v("table-cell-fill"),
                "cell-color": v("table-cell-color"),
            }
        },
    }


def _iter_slide_files(slides_dir: Path) -> list[Path]:
    if not slides_dir.exists() or not slides_dir.is_dir():
        return []
    files = []
    for p in slides_dir.iterdir():
        if p.is_file() and p.suffix.lower() in {".xml", ".html", ".htm"}:
            files.append(p)
    return sorted(files, key=lambda x: x.name)


def validate_and_fix_slide_xml(xml_content: str, theme_name: str = "tech-dark") -> tuple[bool, list[str], str]:
    """
    验证并自动修复 slide XML
    
    Args:
        xml_content: XML 内容字符串
        theme_name: 主题名称，用于获取默认颜色等
    
    Returns:
        (是否有效, 修复信息列表, 修复后的XML内容)
        如果无法修复，返回的XML内容为空字符串
    """
    fixes = []
    
    # 提取 slide 片段
    fragment = _extract_slide_fragment(xml_content)
    if fragment is None:
        return False, ["未找到 <ppt-slide> 根节点"], ""
    
    # 解析 XML
    try:
        slide_el = ET.fromstring(fragment)
    except ET.ParseError as e:
        return False, [f"XML 解析失败: {e}"], ""
    
    # 获取或设置默认尺寸
    width = _parse_float(slide_el.attrib.get("width"))
    height = _parse_float(slide_el.attrib.get("height"))
    
    if width is None:
        slide_el.set("width", "13.333")
        width = 13.333
        fixes.append("已修复: 添加默认 width=13.333")
    
    if height is None:
        slide_el.set("height", "7.5")
        height = 7.5
        fixes.append("已修复: 添加默认 height=7.5")
    
    # 获取主题颜色
    theme = BUILTIN_THEMES.get(theme_name, BUILTIN_THEMES["tech-dark"])
    default_text_color = theme.get("text", {}).get("body", {}).get("color", "#CBD5E1")
    
    # 检查内容元素
    content_tags = {"ppt-text", "ppt-rect", "ppt-table", "ppt-chart", "ppt-image"}
    has_content = False
    
    for child in slide_el:
        tag = (child.tag or "").strip().lower()
        if tag not in content_tags:
            continue
        
        has_content = True
        
        # 获取坐标和尺寸
        x = _parse_float(child.attrib.get("x"))
        y = _parse_float(child.attrib.get("y"))
        w = _parse_float(child.attrib.get("w"))
        h = _parse_float(child.attrib.get("h"))
        
        if x is None or y is None or w is None or h is None:
            continue
        
        # 修复负数坐标
        if x < 0:
            child.set("x", "0")
            fixes.append(f"已修复: {tag} 的 x 坐标从 {x} 调整为 0")
            x = 0
        
        if y < 0:
            child.set("y", "0")
            fixes.append(f"已修复: {tag} 的 y 坐标从 {y} 调整为 0")
            y = 0
        
        # 修复超出边界的坐标
        if x > width - 0.5:
            new_x = max(0, width - 2)
            child.set("x", str(new_x))
            fixes.append(f"已修复: {tag} 的 x 坐标从 {x} 调整为 {new_x}（避免超出幻灯片）")
            x = new_x
        
        if y > height - 0.5:
            new_y = max(0, height - 1)
            child.set("y", str(new_y))
            fixes.append(f"已修复: {tag} 的 y 坐标从 {y} 调整为 {new_y}（避免超出幻灯片）")
            y = new_y
        
        # 修复尺寸溢出
        if x + w > width:
            new_w = max(0.1, width - x - 0.1)
            child.set("w", str(round(new_w, 3)))
            fixes.append(f"已修复: {tag} 的宽度从 {w} 调整为 {new_w:.3f}（避免超出幻灯片）")
            w = new_w
        
        if y + h > height:
            new_h = max(0.1, height - y - 0.1)
            child.set("h", str(round(new_h, 3)))
            fixes.append(f"已修复: {tag} 的高度从 {h} 调整为 {new_h:.3f}（避免超出幻灯片）")
            h = new_h
        
        # 修复文本颜色
        if tag == "ppt-text":
            color = child.attrib.get("color")
            if not color:
                child.set("color", default_text_color)
                fixes.append(f"已修复: 添加默认文本颜色 {default_text_color}")
    
    if not has_content:
        return False, ["页面缺少内容元素 (ppt-text/ppt-rect/ppt-table/ppt-chart/ppt-image)"], ""
    
    # 生成修复后的 XML
    attrs_str = "".join(f' {k}="{v}"' for k, v in slide_el.attrib.items())
    fixed_xml = f'<ppt-slide{attrs_str}>\n'
    for child in slide_el:
        child_str = ET.tostring(child, encoding='unicode')
        fixed_xml += f"  {child_str}\n"
    fixed_xml += '</ppt-slide>'
    
    return True, fixes, fixed_xml


class HtmlPptConverter:
    def __init__(self, project_dir: str | Path, template_path: Path | None = None, verbose: bool = True):
        self.project_dir = Path(project_dir).resolve()
        self.verbose = verbose
        self.debug_skip_images = False
        self.theme_overrides = self._load_theme_overrides()
        self.themes = self._load_css_themes()
        self.default_theme_name = self._load_default_theme_name()
        self.current_slide_path = None
        self.current_slide_width = 13.333
        self.current_slide_height = 7.5
        self.current_layout_items = []
        self.current_element_index = {}
        self.validation_errors = []

        self.slides_dir = self.project_dir / "slides"
        self.assets_dir = self.project_dir / "assets"
        
        # Try to use a base template if available
        # This can help with MacOS preview issues if the template is "blessed"
        if template_path and template_path.exists():
            if self.verbose:
                print(f"Using template: {template_path}")
            self.prs = Presentation(str(template_path))
            self.using_template = True
        else:
            self.prs = Presentation() # Default empty
            self.prs.slide_width = Inches(13.333)
            self.prs.slide_height = Inches(7.5)
            self.using_template = False

    def convert(self, output_file: Path) -> Path:
        prs = self.prs
        
        original_slide_count = len(prs.slides)
        
        # If not using template, ensure slide size is set
        if not self.using_template:
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

        # Set core properties (crucial for MacOS preview)
        if hasattr(prs, "core_properties"):
            prs.core_properties.title = self.project_dir.name
            prs.core_properties.author = "Sage Skills"
            prs.core_properties.subject = "Generated by Sage Skills"
            prs.core_properties.keywords = "presentation, generated"
            # Setting revision helps some parsers
            prs.core_properties.revision = 1

        slides_dir = self.slides_dir
        slide_files = sorted([f for f in slides_dir.glob("*.xml")])
        if not slide_files:
            raise RuntimeError(f"未找到slides文件: {slides_dir}")

        for slide_path in slide_files:
            try:
                self._add_slide_from_file(prs, slide_path)
            except Exception as e:
                self._add_error(f"Error: 幻灯片生成失败: {slide_path} | {e}")

        # If using template, remove the original slides *after* adding new ones
        # This is safer than clearing first, as it avoids "empty presentation" states
        if self.using_template and original_slide_count > 0:
            self._remove_first_n_slides(prs, original_slide_count)

        if self.validation_errors:
            print("Error: 发现布局问题，已阻止生成。", file=sys.stderr)
            for err in self.validation_errors:
                print(err, file=sys.stderr)
            raise RuntimeError("发现布局问题，已阻止生成。")

        output_file = Path(output_file)
        output_file.parent.mkdir(parents=True, exist_ok=True)
        try:
            prs.save(str(output_file))
        except Exception as e:
            print(f"Error saving PPT: {e}")
            raise e
            
        if self.verbose:
            print(f"✓ PPT生成成功: {output_file}")
        return output_file

    def _clear_existing_slides(self, prs) -> None:
        """
        Remove all slides from the presentation while keeping layouts and masters.
        This allows using a template PPTX purely for its properties/theme/layouts
        without keeping its sample content.
        """
        # python-pptx doesn't expose a public API to remove slides easily.
        # We need to access the internal XML element _sldIdLst.
        try:
            slides = prs.slides
            # The underlying xml element is slides._sldIdLst
            if hasattr(slides, "_sldIdLst"):
                xml_slides = slides._sldIdLst
                # Create a list of children to remove to avoid iterator invalidation
                children = list(xml_slides)
                for child in children:
                    xml_slides.remove(child)
                if self.verbose:
                    print("✓ 已清空模板原有幻灯片")
            else:
                 print("Warning: Template slides could not be cleared (structure mismatch)")
        except Exception as e:
            print(f"Warning: Failed to clear template slides: {e}")
            traceback.print_exc()

    def _remove_first_n_slides(self, prs, n: int) -> None:
        """Remove the first n slides from the presentation"""
        try:
            slides = prs.slides
            if hasattr(slides, "_sldIdLst"):
                xml_slides = slides._sldIdLst
                children = list(xml_slides)[:n]
                for child in children:
                    xml_slides.remove(child)
                if self.verbose:
                    print(f"✓ 已移除模板原有{len(children)}张幻灯片")
            else:
                print("Warning: Template slides could not be removed (structure mismatch)")
        except Exception as e:
            print(f"Warning: Failed to remove template slides: {e}")
            traceback.print_exc()

    def _load_theme_overrides(self) -> dict:
        theme_path = self.project_dir / "theme.json"
        if not theme_path.exists():
            return {}
        try:
            data = json.loads(theme_path.read_text(encoding="utf-8"))
            return data if isinstance(data, dict) else {}
        except Exception:
            traceback.print_exc()
            raise
            return {}

    def _load_default_theme_name(self) -> str:
        theme_txt = self.project_dir / "theme.txt"
        if not theme_txt.exists():
            return "tech-dark"
        try:
            raw = theme_txt.read_text(encoding="utf-8")
            for line in raw.splitlines():
                s = line.strip()
                if not s or s.startswith("#"):
                    continue
                if s.lower().startswith("theme="):
                    return s.split("=", 1)[1].strip().lower() or "tech-dark"
        except Exception:
            traceback.print_exc()
        return "tech-dark"

    def _load_css_themes(self) -> dict[str, dict]:
        themes: dict[str, dict] = {}

        skill_dir = Path(__file__).resolve().parent.parent
        builtin_dir = skill_dir / "themes"
        project_dir = self.project_dir / "themes"

        for d in [builtin_dir, project_dir]:
            if not d.exists() or not d.is_dir():
                continue
            for p in d.iterdir():
                if not p.is_file() or p.suffix.lower() != ".css":
                    continue
                name = p.stem.strip().lower()
                try:
                    css_text = p.read_text(encoding="utf-8")
                    vars_map = _parse_css_vars(css_text)
                    themes[name] = _css_vars_to_theme(vars_map)
                except Exception:
                    traceback.print_exc()

        return themes

    def _get_theme(self, name: str | None) -> dict:
        n = (name or "").strip().lower() or self.default_theme_name or "tech-dark"
        if n in {"none", "off", "false"}:
            return {}
        base = self.themes.get(n) or BUILTIN_THEMES.get(n, {})
        override = self.theme_overrides.get(n, {}) if isinstance(self.theme_overrides, dict) else {}
        if not isinstance(override, dict):
            override = {}
        merged = dict(base)
        for k, v in override.items():
            if isinstance(v, dict) and isinstance(merged.get(k), dict):
                merged[k] = _merge_style(merged.get(k, {}), v)
            else:
                merged[k] = v
        return merged

    def _add_slide_from_file(self, prs: Presentation, slide_path: Path) -> None:
        raw = slide_path.read_text(encoding="utf-8")
        fragment = _extract_slide_fragment(raw)
        if fragment is None:
            raise RuntimeError(f"未找到<ppt-slide>根节点: {slide_path}")

        try:
            slide_el = ET.fromstring(fragment)
        except Exception:
            traceback.print_exc()
            raise RuntimeError(f"解析XML失败(请确保<ppt-slide>内是合法XML): {slide_path}")

        width = _parse_float(slide_el.attrib.get("width"), default=13.333)
        height = _parse_float(slide_el.attrib.get("height"), default=7.5)
        if width and height:
            prs.slide_width = Inches(width)
            prs.slide_height = Inches(height)
        self.current_slide_path = slide_path
        self.current_slide_width = width or 13.333
        self.current_slide_height = height or 7.5
        self.current_layout_items = self._collect_layout_items(slide_el)
        self.current_element_index = {item["el"]: item["idx"] for item in self.current_layout_items}

        theme = self._get_theme(slide_el.attrib.get("theme"))
        auto_bg = _parse_bool(slide_el.attrib.get("auto-bg"), default=True)
        notes_text = slide_el.attrib.get("notes")
        bg_src = slide_el.attrib.get("bg-src") or slide_el.attrib.get("background-src")
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        for child in list(slide_el):
            tag = (child.tag or "").strip().lower()
            if tag == "ppt-notes" and (notes_text is None or str(notes_text).strip() == ""):
                notes_text = "".join(child.itertext())
            if tag == "ppt-bg" and not bg_src:
                bg_src = (child.attrib.get("src") or "").strip() or None

        if bg_src and auto_bg and not self.debug_skip_images:
            try:
                bg_path = (slide_path.parent / bg_src).resolve()
                if not bg_path.exists():
                    bg_path = (self.project_dir / bg_src).resolve()
                if bg_path.exists():
                    # Validate image before adding
                    if bg_path.stat().st_size > 0:
                        slide.shapes.add_picture(str(bg_path), Inches(0), Inches(0), prs.slide_width, prs.slide_height)
                    else:
                        print(f"Skipping empty image: {bg_path}")
            except Exception:
                traceback.print_exc()

        if theme.get("page_bg") and auto_bg:
            bg = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(0),
                Inches(0),
                prs.slide_width,
                prs.slide_height,
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = _parse_color(theme.get("page_bg"), default="0B1220")
            try:
                bg.line.fill.background()
            except Exception:
                traceback.print_exc()

        if notes_text is not None and str(notes_text).strip() != "":
            try:
                n = slide.notes_slide
                n.notes_text_frame.text = str(notes_text).strip()
            except Exception:
                traceback.print_exc()

        for child in list(slide_el):
            tag = (child.tag or "").strip().lower()
            if tag == "ppt-text":
                self._add_text(slide, child, theme)
            elif tag == "ppt-rect":
                self._add_rect(slide, child, theme)
            elif tag == "ppt-table":
                self._add_table(slide, child, theme)
            elif tag == "ppt-image":
                self._add_image(slide, child, slide_path)
            elif tag == "ppt-chart":
                self._add_chart(slide, child, theme)
            elif tag in {"ppt-notes", "ppt-bg"}:
                continue

    def _resolve_preset_style(self, tag: str, theme: dict, style_name: str | None) -> dict:
        style_key = (style_name or "").strip()
        if not style_key:
            return {}
        group = theme.get(tag, {})
        if isinstance(group, dict) and style_key in group and isinstance(group[style_key], dict):
            return dict(group[style_key])
        return {}

    def _add_error(self, message: str) -> None:
        if message and message not in self.validation_errors:
            self.validation_errors.append(message)

    def _color_to_rgb_tuple(self, value: str | RGBColor | None, default: str = "111827") -> tuple[int, int, int]:
        if isinstance(value, RGBColor):
            return int(value[0]), int(value[1]), int(value[2])
        rgb = _parse_color(value, default=default)
        return int(rgb[0]), int(rgb[1]), int(rgb[2])

    def _relative_luminance(self, rgb: tuple[int, int, int]) -> float:
        def channel(v: int) -> float:
            c = v / 255.0
            return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4

        r, g, b = rgb
        return 0.2126 * channel(r) + 0.7152 * channel(g) + 0.0722 * channel(b)

    def _contrast_ratio(self, fg: tuple[int, int, int], bg: tuple[int, int, int]) -> float:
        l1 = self._relative_luminance(fg)
        l2 = self._relative_luminance(bg)
        hi, lo = (l1, l2) if l1 >= l2 else (l2, l1)
        return (hi + 0.05) / (lo + 0.05)

    def _resolve_rect_fill(self, el: ET.Element, theme: dict) -> str | None:
        preset = self._resolve_preset_style("rect", theme, el.attrib.get("class") or el.attrib.get("style"))
        fill = el.attrib.get("fill", preset.get("fill"))
        if not fill:
            return None
        if str(fill).lower() in {"none", "transparent", "no"}:
            return None
        return str(fill).strip()

    def _resolve_chart_bg_color(self, el: ET.Element, theme: dict) -> str:
        bg_color = None
        idx = self.current_element_index.get(el)
        if idx is not None:
            x = _parse_float(el.attrib.get("x"))
            y = _parse_float(el.attrib.get("y"))
            w = _parse_float(el.attrib.get("w"))
            h = _parse_float(el.attrib.get("h"))
            if None not in {x, y, w, h}:
                chart_bbox = (x, y, w, h)
                for item in self.current_layout_items:
                    if item["tag"] != "ppt-rect":
                        continue
                    if item["idx"] >= idx:
                        continue
                    rect_bbox = item["bbox"]
                    if not self._bbox_contains(rect_bbox, chart_bbox):
                        continue
                    rect_fill = self._resolve_rect_fill(item["el"], theme)
                    if rect_fill:
                        bg_color = rect_fill
        return bg_color or theme.get("page_bg") or "#FFFFFF"

    def _pick_contrast_color(
        self,
        bg_color: str | RGBColor | None,
        preferred: str | RGBColor | None,
        min_ratio: float = 4.5,
    ) -> tuple[RGBColor, float]:
        bg_tuple = self._color_to_rgb_tuple(bg_color, default="FFFFFF")
        candidates: list[tuple[tuple[int, int, int], float]] = []
        if preferred:
            preferred_tuple = self._color_to_rgb_tuple(preferred, default="111827")
            candidates.append((preferred_tuple, self._contrast_ratio(preferred_tuple, bg_tuple)))
        dark_tuple = self._color_to_rgb_tuple("#111827", default="111827")
        light_tuple = self._color_to_rgb_tuple("#FFFFFF", default="FFFFFF")
        candidates.append((dark_tuple, self._contrast_ratio(dark_tuple, bg_tuple)))
        candidates.append((light_tuple, self._contrast_ratio(light_tuple, bg_tuple)))
        best_tuple, best_ratio = max(candidates, key=lambda item: item[1])
        return RGBColor(*best_tuple), best_ratio

    def _is_placeholder_text(self, text: str) -> bool:
        if not text:
            return False
        stripped = text.strip()
        if not stripped:
            return False
        lower = stripped.lower()
        if re.search(r"\blorem ipsum\b", lower):
            return True
        if re.search(r"\b(placeholder|todo|tbd|draft)\b", lower):
            return True
        if re.fullmatch(r"[xX]{3,}", stripped):
            return True
        if re.fullmatch(r"[-_]{3,}", stripped):
            return True
        keywords = [
            "占位",
            "占位符",
            "待补充",
            "待填写",
            "请填写",
            "请输入",
            "这里是标题",
            "这里是内容",
            "示例文本",
        ]
        return any(k in stripped for k in keywords)

    def _is_placeholder_image(self, src: str) -> bool:
        if not src:
            return False
        lower = src.lower()
        keywords = ["placeholder", "dummy", "sample", "temp", "todo", "tbd", "test", "占位", "示例"]
        return any(k in lower for k in keywords)

    def _contains_cjk(self, text: str) -> bool:
        return re.search(r"[\u4e00-\u9fff]", text or "") is not None

    def _bbox_overlap(self, a: tuple[float, float, float, float], b: tuple[float, float, float, float], eps: float = 1e-4) -> bool:
        ax, ay, aw, ah = a
        bx, by, bw, bh = b
        return not (ax + aw <= bx + eps or bx + bw <= ax + eps or ay + ah <= by + eps or by + bh <= ay + eps)

    def _bbox_contains(self, outer: tuple[float, float, float, float], inner: tuple[float, float, float, float], eps: float = 1e-4) -> bool:
        ox, oy, ow, oh = outer
        ix, iy, iw, ih = inner
        return ix + eps >= ox and iy + eps >= oy and (ix + iw) <= (ox + ow + eps) and (iy + ih) <= (oy + oh + eps)

    def _estimate_text_layout(
        self,
        text: str,
        font_size: float,
        line_spacing: float | None,
        box_w: float,
        pad_left: float | None,
        pad_right: float | None,
    ) -> dict:
        lines = text.split("\n") if text else [""]
        spacing = line_spacing if line_spacing and line_spacing > 0 else 1.0
        bullet_bias = 0.3 if re.search(r'^\s*[-*•]\s+', text or "", flags=re.MULTILINE) else 0.0
        effective_w = max(0.0, box_w - (pad_left or 0) - (pad_right or 0) - bullet_bias)
        if font_size <= 0 or effective_w <= 0:
            return {
                "line_count": len(lines) or 1,
                "text_w": 0.0,
                "text_h": 0.0,
                "max_chars_per_line": 1,
            }
        char_factor = 1.1 if self._contains_cjk(text) else 0.6
        max_chars_per_line = max(1, int((effective_w * 72.0) / (font_size * char_factor)))
        total_lines = 0
        max_line_chars = 0
        for line in lines:
            line_content = (line or "").strip()
            line_len = max(1, len(line_content))
            wrapped = int(math.ceil(line_len / max_chars_per_line))
            total_lines += max(1, wrapped)
            max_line_chars = max(max_line_chars, min(max_chars_per_line, line_len))
        text_w = min(effective_w, (max_line_chars * font_size * char_factor) / 72.0)
        text_h = (total_lines * font_size * spacing) / 72.0
        return {
            "line_count": total_lines,
            "text_w": text_w,
            "text_h": text_h,
            "max_chars_per_line": max_chars_per_line,
        }

    def _fit_font_to_box(
        self,
        text: str,
        font_size: float,
        line_spacing: float | None,
        w: float,
        h: float,
        pad_left: float | None,
        pad_right: float | None,
        pad_top: float | None,
        pad_bottom: float | None,
        autofit: bool = False,
    ) -> tuple[float, dict, bool]:
        """Fit font to box. Returns (font_size, layout, fits).
        If fits is False, the text cannot fit even at minimum font size (6pt)."""
        effective_h = (h - (pad_top or 0) - (pad_bottom or 0)) * 0.9
        if effective_h <= 0:
            layout = self._estimate_text_layout(text, font_size, line_spacing, w, pad_left, pad_right)
            return font_size, layout, True
        layout = self._estimate_text_layout(text, font_size, line_spacing, w, pad_left, pad_right)
        
        # Step 1: If text is too tall, shrink the font
        for _ in range(20):  # Increased iterations
            if layout["text_h"] <= effective_h + 1e-3:
                break
            ratio = effective_h / max(layout["text_h"], 1e-6)
            new_font = max(6.0, font_size * ratio * 0.95)  # More aggressive shrinking
            if abs(new_font - font_size) < 0.1:
                font_size = new_font
                break
            font_size = new_font
            layout = self._estimate_text_layout(text, font_size, line_spacing, w, pad_left, pad_right)
        
        # Check if text still doesn't fit at minimum font size
        fits = layout["text_h"] <= effective_h + 1e-3
        
        # Step 2: If autofit is enabled and text is too small compared to box, grow the font
        if autofit and effective_h > 0:
            # Calculate how much of the box is used
            usage_ratio = layout["text_h"] / effective_h
            # If less than 50% of the box is used, try to grow
            if usage_ratio < 0.5 and font_size < 200:
                for _ in range(12):
                    new_font = font_size * 1.05
                    if new_font > 200:
                        break
                    new_layout = self._estimate_text_layout(text, new_font, line_spacing, w, pad_left, pad_right)
                    # Check if new size still fits
                    if new_layout["text_h"] <= effective_h + 1e-3:
                        font_size = new_font
                        layout = new_layout
                    else:
                        break
        
        layout = self._estimate_text_layout(text, font_size, line_spacing, w, pad_left, pad_right)
        return font_size, layout, fits

    def _estimate_text_bbox(
        self,
        text: str,
        font_size: float,
        line_spacing: float | None,
        align: str,
        valign: str,
        x: float,
        y: float,
        w: float,
        h: float,
        pad_left: float | None,
        pad_right: float | None,
        pad_top: float | None,
        pad_bottom: float | None,
    ) -> tuple[tuple[float, float, float, float], dict]:
        layout = self._estimate_text_layout(text, font_size, line_spacing, w, pad_left, pad_right)
        text_w = layout["text_w"]
        text_h = layout["text_h"]
        effective_h = max(0.0, h - (pad_top or 0) - (pad_bottom or 0))
        if align == "center":
            text_x = x + max(0.0, (w - text_w) / 2.0)
        elif align == "right":
            text_x = x + w - (pad_right or 0) - text_w
        else:
            text_x = x + (pad_left or 0)
        if valign in ("middle", "center"):
            text_y = y + (pad_top or 0) + max(0.0, (effective_h - text_h) / 2.0)
        elif valign == "bottom":
            text_y = y + h - (pad_bottom or 0) - text_h
        else:
            text_y = y + (pad_top or 0)
        return (text_x, text_y, text_w, text_h), layout

    def _get_text_style(self, el: ET.Element, theme: dict) -> dict:
        preset = self._resolve_preset_style("text", theme, el.attrib.get("class") or el.attrib.get("style"))
        text_parts = []
        for child in el:
            if (child.tag or "").strip().lower() == "br":
                text_parts.append("\n")
            else:
                child_text = "".join(child.itertext()).strip()
                if child_text:
                    text_parts.append(child_text)
        
        if text_parts:
            text = "".join(text_parts)
        else:
            text = "".join(el.itertext()).strip()
        
        # Convert literal \n (backslash + n) to actual newlines
        text = text.replace('\\n', '\n')
        text = re.sub(r'\n\s*\n', '\n', text)
        text = text.strip("\n")
        
        icon = (el.attrib.get("icon") or "").strip()
        if icon:
            text = f"{icon} {text}".strip()
        font_size = _parse_float(el.attrib.get("font-size"), default=preset.get("font-size", 14)) or 14
        align = (el.attrib.get("align") or preset.get("align") or "left").strip().lower()
        valign = (el.attrib.get("valign") or preset.get("valign") or "top").strip().lower()
        pad = _parse_float(el.attrib.get("pad"), default=None)
        pad_left = _parse_float(el.attrib.get("pad-left"), default=pad)
        pad_right = _parse_float(el.attrib.get("pad-right"), default=pad)
        pad_top = _parse_float(el.attrib.get("pad-top"), default=pad)
        pad_bottom = _parse_float(el.attrib.get("pad-bottom"), default=pad)
        line_spacing = _parse_float(el.attrib.get("line-spacing"), default=preset.get("line-spacing"))
        autofit = _parse_bool(el.attrib.get("autofit"), default=_parse_bool(preset.get("autofit"), default=True))
        return {
            "text": text,
            "font_size": font_size,
            "align": align,
            "valign": valign,
            "pad_left": pad_left,
            "pad_right": pad_right,
            "pad_top": pad_top,
            "pad_bottom": pad_bottom,
            "line_spacing": line_spacing,
            "autofit": autofit,
        }

    def _collect_layout_items(self, slide_el: ET.Element) -> list[dict]:
        items = []
        for idx, child in enumerate(list(slide_el)):
            tag = (child.tag or "").strip().lower()
            if tag not in {"ppt-text", "ppt-rect", "ppt-image", "ppt-chart", "ppt-table"}:
                continue
            x = _parse_float(child.attrib.get("x"))
            y = _parse_float(child.attrib.get("y"))
            w = _parse_float(child.attrib.get("w"))
            h = _parse_float(child.attrib.get("h"))
            if None in {x, y, w, h}:
                continue
            items.append({"idx": idx, "tag": tag, "bbox": (x, y, w, h), "el": child})
        return items

    def _check_element_within_slide(self, tag: str, x: float, y: float, w: float, h: float) -> None:
        slide_w = self.current_slide_width or 13.333
        slide_h = self.current_slide_height or 7.5
        overflow = x < 0 or y < 0 or (x + w) > (slide_w + 1e-3) or (y + h) > (slide_h + 1e-3)
        if overflow:
            slide_info = str(self.current_slide_path) if self.current_slide_path else "Unknown XML"
            self._add_error(
                f"Error: 元素超出幻灯片范围，必须修改XML: {slide_info} | tag={tag} | x={x}, y={y}, w={w}, h={h}, slide=({slide_w}, {slide_h})"
            )

    def _add_rich_text(self, p, text: str, default_font_props: dict) -> None:
        p.text = ""  # Clear default text
        if not text:
            return

        # Split by **bold**
        parts = re.split(r'(\*\*[^*]+\*\*)', text)
        for part in parts:
            if not part:
                continue
            is_bold = part.startswith("**") and part.endswith("**") and len(part) > 4
            content = part[2:-2] if is_bold else part

            # Split by *italic*
            subparts = re.split(r'(\*[^*]+\*)', content)
            for subpart in subparts:
                if not subpart:
                    continue
                is_italic = subpart.startswith("*") and subpart.endswith("*") and len(subpart) > 2
                subcontent = subpart[1:-1] if is_italic else subpart
                
                if not subcontent:
                    continue

                run = p.add_run()
                run.text = subcontent
                
                # Apply font properties
                run.font.size = Pt(default_font_props.get("size", 14))
                if default_font_props.get("color"):
                    run.font.color.rgb = default_font_props.get("color")
                if default_font_props.get("name"):
                    run.font.name = default_font_props.get("name")
                
                run.font.bold = True if is_bold else default_font_props.get("bold", False)
                run.font.italic = True if is_italic else default_font_props.get("italic", False)
                run.font.underline = default_font_props.get("underline", False)

    def _add_text(self, slide, el: ET.Element, theme: dict) -> None:
        x = _parse_float(el.attrib.get("x"))
        y = _parse_float(el.attrib.get("y"))
        w = _parse_float(el.attrib.get("w"))
        h = _parse_float(el.attrib.get("h"))
        if None in {x, y, w, h}:
            return

        preset = self._resolve_preset_style("text", theme, el.attrib.get("class") or el.attrib.get("style"))
        style = self._get_text_style(el, theme)
        text = style["text"]
            
        # print(f"Adding text '{text[:20]}...' at ({x}, {y}) class={el.attrib.get('class')}")
        font_size = style["font_size"]
        bold = _parse_bool(el.attrib.get("bold"), default=_parse_bool(preset.get("bold"), default=False))
        italic = _parse_bool(el.attrib.get("italic"), default=_parse_bool(preset.get("italic"), default=False))
        underline = _parse_bool(el.attrib.get("underline"), default=_parse_bool(preset.get("underline"), default=False))
        color = _parse_color(el.attrib.get("color") or preset.get("color"))
        font_name = (el.attrib.get("font-name") or preset.get("font-name") or "").strip() or None
        align = style["align"]
        valign = style["valign"]
        pad_left = style["pad_left"]
        pad_right = style["pad_right"]
        pad_top = style["pad_top"]
        pad_bottom = style["pad_bottom"]
        line_spacing = style["line_spacing"]
        autofit = style.get("autofit", False)
        if self._is_placeholder_text(text):
            slide_info = str(self.current_slide_path) if self.current_slide_path else "Unknown XML"
            self._add_error(
                f"Error: 不允许占位文字，必须修改XML: {slide_info} | text='{text[:80]}'"
            )
        if align == "right":
            slide_info = str(self.current_slide_path) if self.current_slide_path else "Unknown XML"
            self._add_error(
                f"Error: 不接受右对齐文本，必须修改XML: {slide_info} | text='{text[:80]}'"
            )
        if re.search(r'^\s*[-*•]\s+', text or "", flags=re.MULTILINE) and align != "left":
            slide_info = str(self.current_slide_path) if self.current_slide_path else "Unknown XML"
            self._add_error(
                f"Error: 列表文字必须左对齐，必须修改XML: {slide_info} | text='{text[:80]}'"
            )

        spacing = line_spacing if line_spacing and line_spacing > 0 else 1.0
        effective_h = h - (pad_top or 0) - (pad_bottom or 0)
        font_size, layout, fits = self._fit_font_to_box(
            text,
            font_size,
            line_spacing,
            w,
            h,
            pad_left,
            pad_right,
            pad_top,
            pad_bottom,
            autofit,
        )
        if not fits:
            slide_info = str(self.current_slide_path) if self.current_slide_path else "Unknown XML"
            self._add_error(
                f"Error: 文本内容超出文本框且无法通过autofit矫正(已缩至最小字体6pt)，必须修改XML: {slide_info} | text='{text[:80]}' | text_h={layout['text_h']:.2f} | box_h={effective_h:.2f}"
            )
        text_bbox, layout = self._estimate_text_bbox(
            text,
            font_size,
            line_spacing,
            align,
            valign,
            x,
            y,
            w,
            h,
            pad_left,
            pad_right,
            pad_top,
            pad_bottom,
        )
        slide_w = self.current_slide_width or 13.333
        slide_h = self.current_slide_height or 7.5
        if text_bbox[0] < 0 or text_bbox[1] < 0 or (text_bbox[0] + text_bbox[2]) > (slide_w + 1e-3) or (text_bbox[1] + text_bbox[3]) > (slide_h + 1e-3):
            slide_info = str(self.current_slide_path) if self.current_slide_path else "Unknown XML"
            self._add_error(
                f"Error: 文本超出幻灯片范围，必须修改XML: {slide_info} | text='{text[:80]}' | text_bbox={text_bbox} | slide=({slide_w}, {slide_h})"
            )

        idx = self.current_element_index.get(el)
        bg_color = None
        if idx is not None:
            for item in self.current_layout_items:
                if item["el"] is el:
                    continue
                other_bbox = item["bbox"]
                if item["tag"] == "ppt-text":
                    other_style = self._get_text_style(item["el"], theme)
                    other_autofit = other_style.get("autofit", False)
                    other_font, _, _ = self._fit_font_to_box(
                        other_style["text"],
                        other_style["font_size"],
                        other_style["line_spacing"],
                        other_bbox[2],
                        other_bbox[3],
                        other_style["pad_left"],
                        other_style["pad_right"],
                        other_style["pad_top"],
                        other_style["pad_bottom"],
                        other_autofit,
                    )
                    other_text_bbox, _ = self._estimate_text_bbox(
                        other_style["text"],
                        other_font,
                        other_style["line_spacing"],
                        other_style["align"],
                        other_style["valign"],
                        other_bbox[0],
                        other_bbox[1],
                        other_bbox[2],
                        other_bbox[3],
                        other_style["pad_left"],
                        other_style["pad_right"],
                        other_style["pad_top"],
                        other_style["pad_bottom"],
                    )
                    if item["idx"] > idx and self._bbox_overlap(text_bbox, other_text_bbox):
                        slide_info = str(self.current_slide_path) if self.current_slide_path else "Unknown XML"
                        other_text = other_style["text"][:80]
                        self._add_error(
                            f"Error: 文本重叠，必须修改XML: {slide_info} | text='{text[:80]}' | overlap_with='{other_text}'"
                        )
                else:
                    if item["idx"] > idx and self._bbox_overlap(text_bbox, other_bbox):
                        slide_info = str(self.current_slide_path) if self.current_slide_path else "Unknown XML"
                        self._add_error(
                            f"Error: 文本被遮挡，必须修改XML: {slide_info} | text='{text[:80]}' | cover_by={item['tag']} | bbox={other_bbox}"
                        )
                    elif item["tag"] == "ppt-rect" and self._bbox_overlap(text_bbox, other_bbox):
                        if not self._bbox_contains(other_bbox, text_bbox):
                            slide_info = str(self.current_slide_path) if self.current_slide_path else "Unknown XML"
                            self._add_error(
                                f"Error: 文本超出背景矩形边框，必须修改XML: {slide_info} | text='{text[:80]}' | rect_bbox={other_bbox} | text_bbox={text_bbox}"
                            )
                        elif item["idx"] < idx:
                            rect_fill = self._resolve_rect_fill(item["el"], theme)
                            if rect_fill:
                                bg_color = rect_fill

        if bg_color is None:
            bg_color = theme.get("page_bg") or "#FFFFFF"
        contrast = self._contrast_ratio(
            self._color_to_rgb_tuple(color, default="111827"),
            self._color_to_rgb_tuple(bg_color, default="FFFFFF"),
        )
        if contrast < 4.5:
            slide_info = str(self.current_slide_path) if self.current_slide_path else "Unknown XML"
            self._add_error(
                f"Error: 文本与背景对比度不足，必须修改XML: {slide_info} | text='{text[:80]}' | ratio={contrast:.2f}"
            )

        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        # Use TEXT_TO_FIT_SHAPE to automatically resize text if it overflows
        # This addresses user feedback about text exceeding boundaries
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.margin_left = Inches(pad_left) if pad_left is not None else 0
        tf.margin_right = Inches(pad_right) if pad_right is not None else 0
        tf.margin_top = Inches(pad_top) if pad_top is not None else 0
        tf.margin_bottom = Inches(pad_bottom) if pad_bottom is not None else 0

        if valign in ("middle", "center"):
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        elif valign == "bottom":
            tf.vertical_anchor = MSO_ANCHOR.BOTTOM
        else:
            tf.vertical_anchor = MSO_ANCHOR.TOP

        p = tf.paragraphs[0]
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT
        
        if line_spacing is not None:
            try:
                # Use multiplier for line spacing if it's a small float (common case),
                # otherwise treat as points if > 10? No, standard CSS line-height is unitless (multiplier).
                # We assume the input is a multiplier (e.g. 1.5).
                # python-pptx: float -> multiplier, Length -> exact points.
                p.line_spacing = line_spacing
            except Exception:
                traceback.print_exc()

        font_props = {
            "size": font_size,
            "bold": bold,
            "italic": italic,
            "underline": underline,
            "color": color,
            "name": font_name
        }
        
        # Split text by newlines and add paragraphs
        # This supports single \n as line break and basic markdown lists
        lines = text.split('\n')
        
        # Clear default paragraph content and ensure clean state
        # tf.clear() leaves one empty paragraph
        if tf.paragraphs:
            tf.paragraphs[0].text = ""
        tf.clear()
        
        # When using TEXT_TO_FIT_SHAPE, explicit font sizes can be problematic if too large,
        # but we still want to respect user intent as a starting point.
        # The auto-fit logic will shrink it if necessary.
        
        for i, line in enumerate(lines):
            # Check for list markers directly on the line content before stripping all whitespace
            # This preserves meaning if needed, though usually we trim.
            # But the regex above handles leading whitespace.
            line_content = line.strip() 
            
            # For the first line, use the existing first paragraph
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            # Fix indentation issues: explicitly reset indent and level
            # This ensures all paragraphs (reused or new) behave consistently
            
            # Check for list markers
            is_list_item = False
            # Match common list markers: "-", "*", "•" followed by space
            # Note: "•" is U+2022
            match = re.match(r'^(\s*)([-*•])\s+(.*)', line_content)
            if match:
                is_list_item = True
                prefix_spaces = match.group(1)
                marker = match.group(2)
                rest = match.group(3)
                
                # Manually add bullet point and use hanging indent
                # This ensures consistent rendering regardless of template styles
                line_content = f"• {rest}"
                
                # Set indentation for hanging indent
                p.level = 0
                try:
                    # Indent text body by 0.25 inches
                    p.paragraph_format.left_indent = Inches(0.25)
                    # Pull first line (bullet) back by 0.25 inches
                    p.paragraph_format.first_line_indent = Inches(-0.25)
                    p.paragraph_format.space_before = Pt(0)
                    p.paragraph_format.space_after = Pt(0)
                except Exception:
                    pass
            else:
                p.level = 0
                try:
                    p.paragraph_format.left_indent = Inches(0)
                    p.paragraph_format.first_line_indent = Inches(0)
                    p.paragraph_format.space_before = Pt(0)
                    # Also reset space after to avoid huge gaps if defaults are weird
                    p.paragraph_format.space_after = Pt(0) 
                except Exception:
                    pass
            
            if align == "center":
                p.alignment = PP_ALIGN.CENTER
            elif align == "right":
                p.alignment = PP_ALIGN.RIGHT
            else:
                p.alignment = PP_ALIGN.LEFT
                
            if line_spacing is not None:
                try:
                    # Use multiplier (float) directly, not Pt()
                    # Revert to p.line_spacing as p.paragraph_format seems to be missing in this env
                    p.line_spacing = line_spacing
                except Exception:
                    traceback.print_exc()

            self._add_rich_text(p, line_content, font_props)

    def _add_rect(self, slide, el: ET.Element, theme: dict) -> None:
        x = _parse_float(el.attrib.get("x"))
        y = _parse_float(el.attrib.get("y"))
        w = _parse_float(el.attrib.get("w"))
        h = _parse_float(el.attrib.get("h"))
        if None in {x, y, w, h}:
            return
        self._check_element_within_slide("ppt-rect", x, y, w, h)

        preset = self._resolve_preset_style("rect", theme, el.attrib.get("class") or el.attrib.get("style"))
        fill = el.attrib.get("fill", preset.get("fill"))
        line = el.attrib.get("line") or el.attrib.get("line-color") or preset.get("line")
        line_width = _parse_float(el.attrib.get("line-width"), default=preset.get("line-width", 1)) or 1
        shape_kind = (el.attrib.get("shape") or preset.get("shape") or "rect").strip().lower()
        radius = _parse_float(el.attrib.get("radius"), default=preset.get("radius"))
        fill_opacity = _parse_float(el.attrib.get("fill-opacity"), default=preset.get("fill-opacity"))
        line_opacity = _parse_float(el.attrib.get("line-opacity"), default=preset.get("line-opacity"))

        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
            if shape_kind in {"rounded", "round", "rounded-rect", "rounded_rectangle"}
            else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(h),
        )
        if radius is not None:
            try:
                if hasattr(shape, "adjustments") and shape.adjustments and 0 <= float(radius) <= 1:
                    shape.adjustments[0] = float(radius)
            except Exception:
                traceback.print_exc()
        
        if fill and str(fill).lower() in {"none", "transparent", "no"}:
            shape.fill.background()
        elif fill:
            shape.fill.solid()
            shape.fill.fore_color.rgb = _parse_color(fill, default="0B1220")
            if fill_opacity is not None:
                try:
                    shape.fill.fore_color.transparency = max(0.0, min(1.0, float(fill_opacity)))
                except Exception:
                    traceback.print_exc()
        else:
            shape.fill.background()

        if line and str(line).lower() not in {"none", "transparent", "no"}:
            shape.line.color.rgb = _parse_color(line, default="1F2937")
            shape.line.width = Pt(line_width)
            if line_opacity is not None:
                try:
                    shape.line.color.transparency = max(0.0, min(1.0, float(line_opacity)))
                except Exception:
                    traceback.print_exc()
        else:
            shape.line.fill.background()

    def _add_table(self, slide, el: ET.Element, theme: dict) -> None:
        x = _parse_float(el.attrib.get("x"))
        y = _parse_float(el.attrib.get("y"))
        w = _parse_float(el.attrib.get("w"))
        h = _parse_float(el.attrib.get("h"))
        if None in {x, y, w, h}:
            return
        self._check_element_within_slide("ppt-table", x, y, w, h)

        preset = self._resolve_preset_style("table", theme, el.attrib.get("class") or el.attrib.get("style"))
        font_size = _parse_float(el.attrib.get("font-size"), default=preset.get("font-size", 12)) or 12
        line_spacing = _parse_float(el.attrib.get("line-spacing"), default=preset.get("line-spacing")) or 1.2
        header = _parse_bool(el.attrib.get("header"), default=_parse_bool(preset.get("header"), default=False))
        header_fill = el.attrib.get("header-fill") or preset.get("header-fill")
        header_color = el.attrib.get("header-color") or preset.get("header-color")
        cell_fill = el.attrib.get("cell-fill") or preset.get("cell-fill")
        cell_color = el.attrib.get("cell-color") or preset.get("cell-color")
        border_width = _parse_float(el.attrib.get("border-width"))
        border_color = el.attrib.get("border-color")
        col_widths = _parse_inch_list(el.attrib.get("col-widths"))
        row_heights = _parse_inch_list(el.attrib.get("row-heights"))

        rows_el = [c for c in list(el) if (c.tag or "").strip().lower() in {"ppt-row", "ppt-tr"}]
        rows = []
        for r in rows_el:
            cells = []
            for c in list(r):
                if (c.tag or "").strip().lower() not in {"ppt-cell", "ppt-td", "ppt-th"}:
                    continue
                cells.append("".join(c.itertext()).strip())
            if cells:
                rows.append(cells)
        if not rows:
            return

        col_count = max(len(r) for r in rows)
        normalized = [r + [""] * (col_count - len(r)) for r in rows]
        if col_widths and len(col_widths) == col_count:
            total_col = sum(col_widths)
            if total_col > 0 and total_col > w + 1e-3:
                scale = w / total_col
                col_widths = [cw * scale for cw in col_widths]
            col_widths_eff = col_widths
        else:
            col_widths_eff = [w / col_count] * col_count
        if row_heights and len(row_heights) == len(normalized):
            if sum(row_heights) > h + 1e-3:
                slide_info = str(self.current_slide_path) if self.current_slide_path else "Unknown XML"
                self._add_error(
                    f"Error: 表格行高总和超过表格高度，必须修改XML: {slide_info} | table=({x}, {y}, {w}, {h}) | rows={sum(row_heights):.3f}"
                )
        else:
            total_required = 0.0
            for r_idx, row in enumerate(normalized):
                max_lines = 1
                for c_idx, value in enumerate(row):
                    cell_text = (value or "").strip()
                    layout = self._estimate_text_layout(
                        cell_text,
                        font_size,
                        line_spacing,
                        col_widths_eff[c_idx],
                        0,
                        0,
                    )
                    max_lines = max(max_lines, layout["line_count"])
                row_h = (max_lines * font_size * line_spacing) / 72.0
                total_required += row_h
            if total_required > h + 1e-3:
                slide_info = str(self.current_slide_path) if self.current_slide_path else "Unknown XML"
                self._add_error(
                    f"Error: 表格内容超出表格高度，必须修改XML: {slide_info} | table=({x}, {y}, {w}, {h}) | need={total_required:.3f}"
                )

        table_shape = slide.shapes.add_table(
            len(normalized),
            col_count,
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(h),
        )
        table = table_shape.table
        if col_widths and len(col_widths) == col_count:
            try:
                for i, cw in enumerate(col_widths):
                    table.columns[i].width = Inches(cw)
            except Exception:
                traceback.print_exc()
        if row_heights and len(row_heights) == len(normalized):
            try:
                for i, rh in enumerate(row_heights):
                    table.rows[i].height = Inches(rh)
            except Exception:
                traceback.print_exc()

        for r_idx, row in enumerate(normalized):
            for c_idx, value in enumerate(row):
                cell = table.cell(r_idx, c_idx)
                tf = cell.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                # p.text = value # Don't set text directly, use rich text if needed, but for table usually simple.
                # However, to support newlines in table cells too:
                if "\n" in value:
                    tf.clear()
                    lines = value.split('\n')
                    for i, v_line in enumerate(lines):
                        if i == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        p.text = v_line.strip()
                        p.font.size = Pt(font_size)
                        # Re-apply styles
                        if header and r_idx == 0:
                            p.font.bold = True
                            if header_color:
                                try:
                                    p.font.color.rgb = _parse_color(header_color, default="E5E7EB")
                                except Exception:
                                    traceback.print_exc()
                        else:
                            if cell_color:
                                try:
                                    p.font.color.rgb = _parse_color(cell_color, default="CBD5E1")
                                except Exception:
                                    traceback.print_exc()
                else:
                    p.text = value
                    p.font.size = Pt(font_size)
                    if header and r_idx == 0:
                        p.font.bold = True
                        if header_color:
                            try:
                                p.font.color.rgb = _parse_color(header_color, default="E5E7EB")
                            except Exception:
                                traceback.print_exc()
                        if header_fill:
                            try:
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = _parse_color(header_fill, default="111827")
                            except Exception:
                                traceback.print_exc()
                    else:
                        if cell_color:
                            try:
                                p.font.color.rgb = _parse_color(cell_color, default="CBD5E1")
                            except Exception:
                                traceback.print_exc()
                        if cell_fill:
                            try:
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = _parse_color(cell_fill, default="0F172A")
                            except Exception:
                                traceback.print_exc()
                if border_width:
                    _apply_table_cell_border(cell, border_color, border_width)

        self._check_table_contrast(slide, table_shape, el, theme)

    def _check_table_contrast(self, slide, table_shape, el: ET.Element, theme: dict) -> None:
        """检查表格文字与背景的对比度，确保文字清晰可见"""
        preset = self._resolve_preset_style("table", theme, el.attrib.get("class") or el.attrib.get("style"))

        # 获取实际使用的颜色（XML属性优先，其次是主题预设）
        header_fill = el.attrib.get("header-fill") or preset.get("header-fill") or "#111827"
        cell_fill = el.attrib.get("cell-fill") or preset.get("cell-fill") or "#0F172A"
        header_color = el.attrib.get("header-color") or preset.get("header-color") or "#E5E7EB"
        cell_color = el.attrib.get("cell-color") or preset.get("cell-color") or "#CBD5E1"

        slide_info = str(self.current_slide_path) if self.current_slide_path else "Unknown XML"

        # 检查表头对比度
        self._validate_contrast_pair(
            slide_info=slide_info,
            bg_color=header_fill,
            text_color=header_color,
            element_name="表格表头",
            color_attrs="header-fill, header-color"
        )

        # 检查单元格对比度
        self._validate_contrast_pair(
            slide_info=slide_info,
            bg_color=cell_fill,
            text_color=cell_color,
            element_name="表格单元格",
            color_attrs="cell-fill, cell-color"
        )

    def _validate_contrast_pair(
        self,
        slide_info: str,
        bg_color: str,
        text_color: str,
        element_name: str,
        color_attrs: str
    ) -> None:
        """验证一对背景色和文字色的对比度"""
        bg_rgb = self._color_to_rgb_tuple(bg_color, default="FFFFFF")
        text_rgb = self._color_to_rgb_tuple(text_color, default="000000")
        contrast = self._contrast_ratio(text_rgb, bg_rgb)

        if contrast < 4.5:
            # 计算推荐的颜色
            recommended_color = self._recommend_text_color(bg_color)

            self._add_error(
                f"Error: {element_name}对比度不足({contrast:.2f})，必须修改XML: {slide_info} | "
                f"{color_attrs} | 当前: {text_color} on {bg_color} | "
                f"建议: 使用 {recommended_color} 作为文字色"
            )

    def _recommend_text_color(self, bg_color: str) -> str:
        """根据背景色推荐合适的文字颜色（黑色或白色）"""
        bg_rgb = self._color_to_rgb_tuple(bg_color, default="FFFFFF")

        # 计算与黑色和白色的对比度
        black_rgb = (0, 0, 0)
        white_rgb = (255, 255, 255)
        black_contrast = self._contrast_ratio(black_rgb, bg_rgb)
        white_contrast = self._contrast_ratio(white_rgb, bg_rgb)

        # 返回对比度更高的颜色
        if white_contrast > black_contrast:
            return "#FFFFFF" if white_contrast >= 4.5 else f"#FFFFFF(对比度{white_contrast:.1f}仍不足)"
        else:
            return "#000000" if black_contrast >= 4.5 else f"#000000(对比度{black_contrast:.1f}仍不足)"

    def _add_image(self, slide, el: ET.Element, slide_path: Path) -> None:
        if self.debug_skip_images:
            return

        x = _parse_float(el.attrib.get("x"))
        y = _parse_float(el.attrib.get("y"))
        w = _parse_float(el.attrib.get("w"))
        h = _parse_float(el.attrib.get("h"))
        src = (el.attrib.get("src") or "").strip()
        if None in {x, y, w, h} or not src:
            return
        if self._is_placeholder_image(src):
            slide_info = str(self.current_slide_path) if self.current_slide_path else "Unknown XML"
            self._add_error(
                f"Error: 不允许占位图片，必须修改XML: {slide_info} | src='{src}'"
            )
        self._check_element_within_slide("ppt-image", x, y, w, h)
        # Try to find the image
        img_path = Path(src)
        if not img_path.is_absolute():
            # Try relative to assets dir
            check_path = self.assets_dir / src
            if check_path.exists():
                img_path = check_path
            else:
                # Try relative to slide file
                check_path = slide_path.parent / src
                if check_path.exists():
                    img_path = check_path
                else:
                    print(f"DEBUG: Image not found: {src}")
                    print(f"DEBUG: Checked assets_dir: {self.assets_dir / src}")
                    print(f"DEBUG: Checked slide_dir: {slide_path.parent / src}")
                    raise RuntimeError(f"图片不存在: {src}")
        slide.shapes.add_picture(str(img_path), Inches(x), Inches(y), Inches(w), Inches(h))

    def _parse_chart_data(self, el: ET.Element) -> CategoryChartData | None:
        """Parse CSV-like content from element text into CategoryChartData"""
        text = "".join(el.itertext()).strip()
        
        # Also check 'data' attribute
        if not text:
            data_attr = el.attrib.get("data", "").strip()
            if data_attr:
                text = data_attr

        # specific default based on chart type if possible, but generic is fine
        if not text or text.startswith("{{"):
             text = "Category, Series 1\nItem A, 10\nItem B, 20\nItem C, 15"

        if not text:
            return None

        # Use csv module to handle quoting properly
        f = io.StringIO(text)
        reader = csv.reader(f, skipinitialspace=True)
        rows = list(reader)

        if len(rows) < 2:
            return None

        # First row is headers: [CategoryLabel, Series1, Series2, ...]
        headers = rows[0]
        if not headers:
            return None

        # First column is categories
        categories = []
        series_data = [[] for _ in range(len(headers) - 1)]

        for row in rows[1:]:
            if not row:
                continue
            # Ensure row has enough columns
            if len(row) < 1:
                continue
            
            categories.append(row[0])
            for i in range(len(headers) - 1):
                val_str = row[i + 1] if (i + 1) < len(row) else "0"
                try:
                    val = float(val_str)
                except ValueError:
                    val = 0.0
                series_data[i].append(val)

        chart_data = CategoryChartData()
        chart_data.categories = categories
        
        for i, series_vals in enumerate(series_data):
            series_name = headers[i + 1]
            chart_data.add_series(series_name, tuple(series_vals))

        return chart_data

    def _add_chart(self, slide, el: ET.Element, theme: dict) -> None:
        x = _parse_float(el.attrib.get("x"))
        y = _parse_float(el.attrib.get("y"))
        w = _parse_float(el.attrib.get("w"))
        h = _parse_float(el.attrib.get("h"))
        
        if None in {x, y, w, h}:
            return
        self._check_element_within_slide("ppt-chart", x, y, w, h)

        chart_type_str = (el.attrib.get("type") or "column").strip().lower()
        title = el.attrib.get("title")
        
        # Map type string to XL_CHART_TYPE
        type_map = {
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "col": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE,
            "radar": XL_CHART_TYPE.RADAR,
            "area": XL_CHART_TYPE.AREA,
            "scatter": XL_CHART_TYPE.XY_SCATTER,
        }
        
        chart_type = type_map.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)
        
        chart_data = self._parse_chart_data(el)
        if not chart_data:
            print(f"Warning: Empty or invalid data for chart at ({x}, {y})")
            return

        try:
            graphic_frame = slide.shapes.add_chart(
                chart_type, Inches(x), Inches(y), Inches(w), Inches(h), chart_data
            )
            chart = graphic_frame.chart
            bg_color = self._resolve_chart_bg_color(el, theme)
            theme_text = theme.get("text", {}) if isinstance(theme.get("text"), dict) else {}
            title_color = None
            body_color = None
            if isinstance(theme_text, dict):
                title_color = (theme_text.get("title") or {}).get("color")
                body_color = (theme_text.get("body") or {}).get("color")
            title_rgb, title_ratio = self._pick_contrast_color(bg_color, title_color)
            body_rgb, body_ratio = self._pick_contrast_color(bg_color, body_color)
            
            if title:
                chart.has_title = True
                chart.chart_title.text_frame.text = title
                try:
                    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = title_rgb
                except Exception:
                    traceback.print_exc()
            
            # Basic styling
            if chart_type in (XL_CHART_TYPE.PIE,):
                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            if chart.has_legend and chart.legend:
                try:
                    chart.legend.font.color.rgb = body_rgb
                except Exception:
                    traceback.print_exc()
            try:
                # 饼图等某些图表类型没有 category_axis，需要特殊处理
                try:
                    cat_axis = chart.category_axis
                    if cat_axis and hasattr(cat_axis, 'has_tick_labels') and cat_axis.has_tick_labels:
                        cat_axis.tick_labels.font.color.rgb = body_rgb
                except ValueError:
                    # 饼图等没有 category axis，忽略
                    pass
                try:
                    val_axis = chart.value_axis
                    if val_axis and hasattr(val_axis, 'has_tick_labels') and val_axis.has_tick_labels:
                        val_axis.tick_labels.font.color.rgb = body_rgb
                except ValueError:
                    # 某些图表没有 value axis，忽略
                    pass
            except Exception:
                traceback.print_exc()
            try:
                for series in chart.series:
                    if hasattr(series, 'has_data_labels') and series.has_data_labels:
                        series.data_labels.font.color.rgb = body_rgb
            except Exception:
                traceback.print_exc()
            if title_ratio < 4.5 or body_ratio < 4.5:
                slide_info = str(self.current_slide_path) if self.current_slide_path else "Unknown XML"
                self._add_error(
                    f"Error: 图表文字与背景对比度不足，必须修改XML: {slide_info} | ratio_title={title_ratio:.2f} | ratio_body={body_ratio:.2f}"
                )

        except Exception:
            traceback.print_exc()
            print(f"Error creating chart type '{chart_type_str}'")


def main() -> None:
    parser = argparse.ArgumentParser(description="Convert constrained HTML/XML slides into editable PPTX")
    parser.add_argument("project_dir", help="Project directory containing slides/")
    parser.add_argument("--out", dest="out", default=None, help="Output PPTX path")
    parser.add_argument("--template", dest="template", default=None, help="Base PPTX template to use")
    parser.add_argument("--debug-no-images", action="store_true", help="Skip adding images for debugging")
    args = parser.parse_args()

    try:
        project_dir = Path(os.path.abspath(args.project_dir))
        out = Path(os.path.abspath(args.out)) if args.out else (project_dir / "presentation.pptx")
        template_path = Path(os.path.abspath(args.template)) if args.template else None
        
        # If no template is provided, try to use the default built-in template
        if not template_path:
            script_dir = Path(__file__).resolve().parent
            default_template = script_dir.parent / "templates" / "default.pptx"
            if default_template.exists():
                print(f"Using default built-in template: {default_template}")
                template_path = default_template

        converter = HtmlPptConverter(project_dir, template_path)
        if args.debug_no_images:
            converter.debug_skip_images = True
        result = converter.convert(out)
        print(f"✓ PPT生成成功: {result}")
    except Exception:
        traceback.print_exc()
        print("错误: HTML转PPT失败。", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
